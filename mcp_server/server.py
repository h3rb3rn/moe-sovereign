"""
MoE Precision Tools — MCP Server
Exact calculations via Linux/Python tools for everything where LLMs systematically fail:
Arithmetic, dates, units, statistics, hashing, regex, networking, and more.
"""

import ast
import asyncio
import base64
import calendar
import csv
import hashlib
import inspect
import io
import ipaddress
import importlib.metadata
import json
import logging
import math
import operator as op_module
import os
import re
import statistics as stats_module
import sys
import threading
import time
import xml.etree.ElementTree as ET
import zipfile
from contextlib import asynccontextmanager
from datetime import datetime, timedelta, timezone
from decimal import (
    ROUND_CEILING,
    ROUND_DOWN,
    ROUND_FLOOR,
    ROUND_HALF_DOWN,
    ROUND_HALF_EVEN,
    ROUND_HALF_UP,
    ROUND_UP,
    Context,
    Decimal,
    InvalidOperation,
    localcontext,
)
from fractions import Fraction
from importlib import resources
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Tuple, Union, get_args, get_origin
from zoneinfo import ZoneInfo, ZoneInfoNotFoundError

import httpx
import yaml
try:
    from defusedxml import ElementTree as DefusedElementTree
except ImportError:  # Host-only test fallback; the production lock pins defusedxml.
    DefusedElementTree = ET
from jsonschema import Draft202012Validator
from jsonschema.exceptions import SchemaError
import sympy as sp
import uvicorn
from fastapi import FastAPI
from mcp.server.fastmcp import FastMCP
from pydantic import BaseModel

logger = logging.getLogger("MOE-SOVEREIGN.MCP")

# --- MCP Server ---
mcp = FastMCP("precision-tools")

# ─── SAFE ARITHMETIC EVAL ────────────────────────────────────────────────────

_SAFE_OPS = {
    ast.Add: op_module.add,
    ast.Sub: op_module.sub,
    ast.Mult: op_module.mul,
    ast.Div: op_module.truediv,
    ast.Pow: op_module.pow,
    ast.USub: op_module.neg,
    ast.Mod: op_module.mod,
    ast.FloorDiv: op_module.floordiv,
}
_SAFE_FUNCS = {
    "sqrt": math.sqrt, "abs": abs, "round": round,
    "sin": math.sin, "cos": math.cos, "tan": math.tan,
    "asin": math.asin, "acos": math.acos, "atan": math.atan,
    "log": math.log, "log10": math.log10, "log2": math.log2,
    "exp": math.exp, "floor": math.floor, "ceil": math.ceil,
    "factorial": math.factorial,
}
_SAFE_CONSTS = {"pi": math.pi, "e": math.e, "tau": math.tau, "inf": math.inf}


def _safe_eval_node(node: ast.AST) -> Any:
    if isinstance(node, ast.Constant):
        return node.value
    if isinstance(node, ast.Name):
        if node.id in _SAFE_CONSTS:
            return _SAFE_CONSTS[node.id]
        raise ValueError(f"Unknown variable: {node.id}")
    if isinstance(node, ast.BinOp):
        if type(node.op) not in _SAFE_OPS:
            raise ValueError(f"Disallowed operation: {type(node.op)}")
        return _SAFE_OPS[type(node.op)](_safe_eval_node(node.left), _safe_eval_node(node.right))
    if isinstance(node, ast.UnaryOp):
        if type(node.op) not in _SAFE_OPS:
            raise ValueError(f"Disallowed operation: {type(node.op)}")
        return _SAFE_OPS[type(node.op)](_safe_eval_node(node.operand))
    if isinstance(node, ast.Call) and isinstance(node.func, ast.Name):
        fn = node.func.id
        if fn not in _SAFE_FUNCS:
            raise ValueError(f"Unknown function: {fn}")
        args = [_safe_eval_node(a) for a in node.args]
        return _SAFE_FUNCS[fn](*args)
    raise ValueError(f"Disallowed AST node: {ast.dump(node)}")


# ─── TOOLS ───────────────────────────────────────────────────────────────────

@mcp.tool()
def calculate(expression: str) -> str:
    """
    Calculates mathematical expressions exactly without LLM hallucination.
    Supports: +, -, *, /, **, %, //, parentheses, sqrt(), sin(), cos(), log(), factorial() etc.
    Also supports percentages: '15% of 239.99' or '239.99 * 0.15'.
    Example: calculate("sqrt(2) * pi") → 4.442882938...

    Security note: expressions that parse as valid Python but contain unsafe
    constructs (import, attribute access, calls to non-whitelisted names) are
    rejected by the AST validator and never reach the SymPy fallback.
    The SymPy fallback is reserved exclusively for SyntaxErrors (e.g. implicit
    multiplication like '2x') which our safe AST evaluator cannot parse.
    """
    # Normalize percentage notation: "15% of X" → "(15/100)*X"
    expr = re.sub(
        r"(\d+(?:\.\d+)?)\s*%\s*(?:of|von)\s*",
        lambda m: f"({m.group(1)}/100)*",
        expression,
        flags=re.IGNORECASE,
    )
    expr = re.sub(r"(\d+(?:\.\d+)?)%", r"(\1/100)", expr)

    # Stage 1: safe AST evaluation (whitelist-only — no arbitrary Python allowed).
    try:
        tree = ast.parse(expr.strip(), mode="eval")
    except SyntaxError:
        # SyntaxError: expression may use implicit multiplication or other
        # SymPy-parseable syntax.  Fall through to the SymPy fallback below.
        pass
    else:
        # AST parsed successfully — run the safe evaluator.
        # Any ValueError here means an unsafe construct was detected; return
        # an error immediately without falling through to SymPy.
        try:
            result = _safe_eval_node(tree.body)
        except Exception as e:
            return f"Error: {e}"
        if isinstance(result, float) and result == int(result) and abs(result) < 1e15:
            return f"{expression} = {int(result)}"
        if isinstance(result, float):
            return f"{expression} = {result:.12g}"
        return f"{expression} = {result}"

    # Stage 2: SymPy fallback — only reached via SyntaxError above.
    try:
        sym_result = sp.sympify(expression)
        simplified = sp.simplify(sym_result)
        numeric = float(simplified.evalf()) if simplified.is_number else None
        if numeric is not None:
            return f"{expression} = {numeric:.12g} (exact: {simplified})"
        return f"{expression} = {simplified}"
    except Exception as e2:
        return f"Error: {e2}"


@mcp.tool()
def solve_equation(equation: str, variable: str = "x") -> str:
    """
    Solves algebraic equations exactly via SymPy.
    Examples: 'x**2 - 4 = 0', 'x**3 - 6*x**2 + 11*x - 6 = 0', '2*x + 5 = 13'
    """
    try:
        var = sp.Symbol(variable)
        if "=" in equation:
            left, right = equation.split("=", 1)
            eq = sp.Eq(sp.sympify(left.strip()), sp.sympify(right.strip()))
        else:
            eq = sp.Eq(sp.sympify(equation.strip()), 0)
        solutions = sp.solve(eq, var)
        numeric = [float(s.evalf()) if s.is_number else str(s) for s in solutions]
        return (
            f"Equation: {equation}\n"
            f"Solutions ({variable}): {[str(s) for s in solutions]}\n"
            f"Numeric: {numeric}"
        )
    except Exception as e:
        return f"Error: {e}"


@mcp.tool()
def date_diff(date1: str, date2: str) -> str:
    """
    Calculates the exact difference between two dates.
    Format: YYYY-MM-DD. Returns exact total days and a calendar delta.
    Example: date_diff('1990-05-15', '2026-03-29')
    """
    try:
        from dateutil.relativedelta import relativedelta

        d1 = datetime.strptime(date1.strip(), "%Y-%m-%d").date()
        d2 = datetime.strptime(date2.strip(), "%Y-%m-%d").date()
        earlier, later = (d1, d2) if d1 <= d2 else (d2, d1)
        diff = (later - earlier).days
        calendar_delta = relativedelta(later, earlier)
        return (
            f"From {earlier} to {later}: "
            f"{diff} days total "
            f"({calendar_delta.years} years, {calendar_delta.months} months, "
            f"{calendar_delta.days} days calendar delta)"
        )
    except Exception as e:
        return f"Error: {e}"


@mcp.tool()
def date_add(base_date: str, days: int = 0, months: int = 0, years: int = 0) -> str:
    """
    Adds or subtracts time from a date.
    Format base_date: YYYY-MM-DD. Negative values for subtraction.
    Example: date_add('2026-01-01', months=3, days=-5)
    """
    try:
        from dateutil.relativedelta import relativedelta
        d = datetime.strptime(base_date.strip(), "%Y-%m-%d").date()
        result = d + relativedelta(years=years, months=months, days=days)
        days_en = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]
        return (
            f"{base_date} + {years}a {months}m {days}d "
            f"= {result} ({days_en[result.weekday()]})"
        )
    except Exception as e:
        return f"Error: {e}"


_CALENDAR_NAMES = {
    "de": {
        "weekdays": (
            "Montag", "Dienstag", "Mittwoch", "Donnerstag", "Freitag",
            "Samstag", "Sonntag",
        ),
        "months": (
            "Januar", "Februar", "März", "April", "Mai", "Juni",
            "Juli", "August", "September", "Oktober", "November", "Dezember",
        ),
    },
    "en": {
        "weekdays": (
            "Monday", "Tuesday", "Wednesday", "Thursday", "Friday",
            "Saturday", "Sunday",
        ),
        "months": (
            "January", "February", "March", "April", "May", "June",
            "July", "August", "September", "October", "November", "December",
        ),
    },
}

_PINNED_TZDATA_VERSION = "2026.3"
_TIME_YEAR_MIN = 1900
_TIME_YEAR_MAX = 2100
_ISO_INSTANT_PATTERN = re.compile(
    r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}(?:Z|[+-]\d{2}:\d{2})$"
)
_ISO_LOCAL_PATTERN = re.compile(r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}$")
_IANA_ZONE_PATTERN = re.compile(
    r"^[A-Za-z][A-Za-z0-9._+-]*(?:/[A-Za-z0-9._+-]+)*$"
)


def _tzdata_version() -> str:
    """Return the packaged IANA database version used by the container."""
    try:
        return importlib.metadata.version("tzdata")
    except importlib.metadata.PackageNotFoundError:
        # Source-only unit tests may run without the MCP dependency set. The
        # production image is fail-closed below and pins the package in its
        # lock file; this marker makes a system-database fallback visible.
        return "system-unpinned"


def _load_zone(name: str) -> ZoneInfo:
    """Load an IANA zone from the pinned Python package, not the host OS."""
    if not isinstance(name, str) or not _IANA_ZONE_PATTERN.fullmatch(name):
        raise ValueError("timezone must be a valid IANA zone name")
    version = _tzdata_version()
    if version != "system-unpinned":
        if version != _PINNED_TZDATA_VERSION:
            raise RuntimeError(
                f"tzdata_version_mismatch:{version}:{_PINNED_TZDATA_VERSION}"
            )
        try:
            resource = resources.files("tzdata.zoneinfo").joinpath(*name.split("/"))
            with resource.open("rb") as handle:
                return ZoneInfo.from_file(handle, key=name)
        except (FileNotFoundError, ModuleNotFoundError, ValueError) as exc:
            raise ValueError(f"unknown IANA timezone: {name}") from exc
    try:
        return ZoneInfo(name)
    except ZoneInfoNotFoundError as exc:
        raise ValueError(f"unknown IANA timezone: {name}") from exc


def _parse_explicit_instant(value: str) -> datetime:
    raw = value.strip() if isinstance(value, str) else ""
    if not _ISO_INSTANT_PATTERN.fullmatch(raw):
        raise ValueError(
            "instant must be ISO-8601 with seconds and an explicit Z/UTC offset"
        )
    try:
        parsed = datetime.fromisoformat(raw[:-1] + "+00:00" if raw.endswith("Z") else raw)
    except ValueError as exc:
        raise ValueError("instant must be a valid ISO-8601 instant") from exc
    if parsed.utcoffset() is None:
        raise ValueError("instant must include an explicit Z/UTC offset")
    if not _TIME_YEAR_MIN <= parsed.year <= _TIME_YEAR_MAX:
        raise ValueError(f"instant year must be {_TIME_YEAR_MIN}..{_TIME_YEAR_MAX}")
    return parsed


def _parse_local_datetime(value: str) -> datetime:
    raw = value.strip() if isinstance(value, str) else ""
    if not _ISO_LOCAL_PATTERN.fullmatch(raw):
        raise ValueError(
            "local_datetime must be a naive ISO local time with seconds"
        )
    try:
        parsed = datetime.fromisoformat(raw)
    except ValueError as exc:
        raise ValueError("local_datetime must be a valid ISO local time") from exc
    if not _TIME_YEAR_MIN <= parsed.year <= _TIME_YEAR_MAX:
        raise ValueError(
            f"local_datetime year must be {_TIME_YEAR_MIN}..{_TIME_YEAR_MAX}"
        )
    return parsed


def _iso_seconds(value: datetime) -> str:
    return value.isoformat(timespec="seconds")


def _offset_text(value: datetime) -> str:
    offset = value.utcoffset()
    if offset is None:
        raise ValueError("timezone offset unavailable")
    total = int(offset.total_seconds())
    sign = "+" if total >= 0 else "-"
    hours, remainder = divmod(abs(total), 3600)
    minutes = remainder // 60
    return f"{sign}{hours:02d}:{minutes:02d}"


def _is_dst(value: datetime) -> bool:
    delta = value.dst()
    return bool(delta and delta != timedelta(0))


def _localized_time_facts(value: datetime, locale: str) -> Dict[str, Any]:
    locale_code = locale.strip().casefold() if isinstance(locale, str) else ""
    if locale_code not in _CALENDAR_NAMES:
        raise ValueError("locale must be one of: de, en")
    iso = value.date().isocalendar()
    return {
        "date": value.date().isoformat(),
        "time": value.time().isoformat(timespec="seconds"),
        "weekday_iso": iso.weekday,
        "weekday_name": _CALENDAR_NAMES[locale_code]["weekdays"][value.weekday()],
        "iso_week": iso.week,
        "iso_week_year": iso.year,
        "locale": locale_code,
    }


def _calendar_facts_payload(date_str: str, locale: str) -> Dict[str, Any]:
    """Build exact proleptic-Gregorian calendar facts for a strict ISO date."""
    raw_date = date_str.strip() if isinstance(date_str, str) else ""
    if not re.fullmatch(r"\d{4}-\d{2}-\d{2}", raw_date):
        raise ValueError("date_str must be a valid ISO date in YYYY-MM-DD format")

    locale_code = locale.strip().casefold() if isinstance(locale, str) else ""
    if locale_code not in _CALENDAR_NAMES:
        raise ValueError("locale must be one of: de, en")

    try:
        parsed = datetime.strptime(raw_date, "%Y-%m-%d").date()
    except ValueError as exc:
        raise ValueError("date_str must be a valid ISO date in YYYY-MM-DD format") from exc

    iso = parsed.isocalendar()
    leap_year = calendar.isleap(parsed.year)
    names = _CALENDAR_NAMES[locale_code]
    return {
        "calendar_system": "proleptic_gregorian",
        "date": parsed.isoformat(),
        "day": parsed.day,
        "day_of_year": parsed.timetuple().tm_yday,
        "days_in_month": calendar.monthrange(parsed.year, parsed.month)[1],
        "days_in_year": 366 if leap_year else 365,
        "is_leap_year": leap_year,
        "is_weekend": iso.weekday >= 6,
        "iso_week": iso.week,
        "iso_week_year": iso.year,
        "locale": locale_code,
        "month": parsed.month,
        "month_name": names["months"][parsed.month - 1],
        "quarter": ((parsed.month - 1) // 3) + 1,
        "weekday_iso": iso.weekday,
        "weekday_name": names["weekdays"][parsed.weekday()],
        "year": parsed.year,
    }


@mcp.tool()
def calendar_facts(date_str: str, locale: str = "de") -> str:
    """Return deterministic calendar facts as stable JSON.

    ``date_str`` must be an absolute ISO date (YYYY-MM-DD); relative inputs
    such as "today" are deliberately rejected because their answer depends on
    a caller-supplied clock and time zone. ``locale`` is ``de`` or ``en`` and
    affects only the weekday and month names.

    The result includes the canonical date, localized weekday/month names,
    ISO weekday, ISO week and ISO week-year, day of year, month/year lengths,
    quarter, leap-year status and weekend status.
    """
    payload = _calendar_facts_payload(date_str, locale)
    return json.dumps(
        payload,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
    )


@mcp.tool()
def time_facts(instant: str, timezone_name: str = "UTC", locale: str = "de") -> str:
    """Return deterministic timezone facts for an explicit ISO-8601 instant.

    ``instant`` must include ``Z`` or a numeric UTC offset. ``timezone_name``
    must be an IANA identifier such as ``Europe/Berlin``. Relative clock terms
    such as "now" are deliberately unsupported.
    """
    parsed = _parse_explicit_instant(instant)
    zone = _load_zone(timezone_name)
    utc_value = parsed.astimezone(timezone.utc)
    local_value = utc_value.astimezone(zone)
    offset = local_value.utcoffset()
    facts = {
        "input_instant": _iso_seconds(parsed),
        "utc_instant": _iso_seconds(utc_value),
        "as_of": _iso_seconds(utc_value),
        "timezone": timezone_name,
        "local_datetime": _iso_seconds(local_value),
        "utc_offset": _offset_text(local_value),
        "utc_offset_seconds": int(offset.total_seconds()) if offset else 0,
        "timezone_abbreviation": local_value.tzname() or "",
        "is_dst": _is_dst(local_value),
        "fold": int(local_value.fold),
        "tzdata_version": _tzdata_version(),
        **_localized_time_facts(local_value, locale),
    }
    return json.dumps(facts, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


def _valid_local_candidates(local_value: datetime, zone: ZoneInfo) -> Dict[int, datetime]:
    """Return fold-indexed candidates that survive an exact UTC round trip."""
    candidates: Dict[int, datetime] = {}
    for fold in (0, 1):
        aware = local_value.replace(tzinfo=zone, fold=fold)
        round_trip = aware.astimezone(timezone.utc).astimezone(zone)
        if (
            round_trip.replace(tzinfo=None) == local_value
            and int(round_trip.fold) == fold
        ):
            candidates[fold] = aware
    return candidates


@mcp.tool()
def timezone_convert(
    local_datetime: str,
    from_timezone: str,
    to_timezone: str,
    fold: Optional[int] = None,
    locale: str = "de",
) -> str:
    """Convert a naive local ISO time between explicit IANA zones.

    A repeated DST-fold time requires ``fold=0`` or ``fold=1``. A local time
    in a spring-forward gap is rejected. The function never silently chooses
    a side of an ambiguity and never infers a zone from a geographic place.
    """
    local_value = _parse_local_datetime(local_datetime)
    source_zone = _load_zone(from_timezone)
    target_zone = _load_zone(to_timezone)
    candidates = _valid_local_candidates(local_value, source_zone)
    if not candidates:
        raise ValueError("nonexistent_local_time_gap")
    ambiguous = len(candidates) == 2 and (
        candidates[0].utcoffset() != candidates[1].utcoffset()
    )
    if ambiguous and fold is None:
        raise ValueError("ambiguous_local_time_fold_required")
    if fold is not None and fold not in (0, 1):
        raise ValueError("fold must be 0, 1, or null")
    selected_fold = int(fold) if ambiguous else 0
    if selected_fold not in candidates:
        raise ValueError("fold_not_valid_for_local_time")
    source_value = candidates[selected_fold]
    utc_value = source_value.astimezone(timezone.utc)
    target_value = utc_value.astimezone(target_zone)
    facts = {
        "input_local_datetime": local_datetime,
        "from_timezone": from_timezone,
        "to_timezone": to_timezone,
        "fold": selected_fold,
        "ambiguous": ambiguous,
        "source_datetime": _iso_seconds(source_value),
        "source_utc_offset": _offset_text(source_value),
        "source_is_dst": _is_dst(source_value),
        "utc_instant": _iso_seconds(utc_value),
        "as_of": _iso_seconds(utc_value),
        "target_datetime": _iso_seconds(target_value),
        "target_utc_offset": _offset_text(target_value),
        "target_is_dst": _is_dst(target_value),
        "target_fold": int(target_value.fold),
        "locale": locale.strip().casefold(),
        "tzdata_version": _tzdata_version(),
    }
    # Validate locale through the same fixed table used by other time tools.
    _localized_time_facts(target_value, facts["locale"])
    return json.dumps(facts, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


_DECIMAL_INPUT_PATTERN = re.compile(r"^-?(?:0|[1-9]\d{0,47})(?:\.\d{1,24})?$")
_DECIMAL_ROUNDING = {
    "half_even": ROUND_HALF_EVEN,
    "half_up": ROUND_HALF_UP,
    "half_down": ROUND_HALF_DOWN,
    "down": ROUND_DOWN,
    "up": ROUND_UP,
    "floor": ROUND_FLOOR,
    "ceiling": ROUND_CEILING,
}
_DECIMAL_CONTEXT_PRECISION = 128


def _contract_decimal(value: str, field: str) -> Decimal:
    """Parse a bounded canonical decimal string without touching a float."""
    if not isinstance(value, str) or not _DECIMAL_INPUT_PATTERN.fullmatch(value):
        raise ValueError(f"{field}_must_be_canonical_decimal_string")
    parsed = Decimal(value)
    if not parsed.is_finite():
        raise ValueError(f"{field}_must_be_finite")
    return parsed


def _decimal_string(value: Decimal) -> str:
    rendered = format(value, "f")
    if "." in rendered:
        rendered = rendered.rstrip("0").rstrip(".")
    return "0" if rendered in {"", "-0"} else rendered


@mcp.tool()
def decimal_finance(
    operation: str,
    operands: List[str],
    currency: str,
    scale: int,
    rounding: str,
) -> str:
    """Perform bounded financial arithmetic with explicit Decimal semantics.

    ``operands`` has an operation-specific positional contract:
    add/subtract/multiply/divide=[left,right], percentage=[base,percent],
    simple_interest=[principal,annual_percent,years], and
    compound_interest=[principal,annual_percent,years,compounds_per_year].
    This tool contains no tax, exchange-rate, jurisdiction or legal rules.
    """
    allowed = {
        "add", "subtract", "multiply", "divide", "percentage",
        "simple_interest", "compound_interest",
    }
    if operation not in allowed:
        raise ValueError("unsupported_decimal_finance_operation")
    if not isinstance(currency, str) or not re.fullmatch(r"[A-Z]{3}", currency):
        raise ValueError("currency_must_be_iso_4217_alpha3")
    if isinstance(scale, bool) or not isinstance(scale, int) or not 0 <= scale <= 12:
        raise ValueError("scale_out_of_range")
    if rounding not in _DECIMAL_ROUNDING:
        raise ValueError("unsupported_rounding_mode")
    expected_count = {
        "add": 2, "subtract": 2, "multiply": 2, "divide": 2,
        "percentage": 2, "simple_interest": 3, "compound_interest": 4,
    }[operation]
    if not isinstance(operands, list) or len(operands) != expected_count:
        raise ValueError(f"{operation}_requires_{expected_count}_operands")
    values = [_contract_decimal(value, f"operands[{index}]") for index, value in enumerate(operands)]
    quantum = Decimal(1).scaleb(-scale)
    context = Context(prec=_DECIMAL_CONTEXT_PRECISION, Emin=-999, Emax=999)
    with localcontext(context):
        if operation == "add":
            calculated = values[0] + values[1]
        elif operation == "subtract":
            calculated = values[0] - values[1]
        elif operation == "multiply":
            calculated = values[0] * values[1]
        elif operation == "divide":
            if values[1] == 0:
                raise ValueError("division_by_zero")
            calculated = values[0] / values[1]
        elif operation == "percentage":
            calculated = values[0] * values[1] / Decimal(100)
        elif operation == "simple_interest":
            if values[2] < 0:
                raise ValueError("years_must_be_non_negative")
            calculated = values[0] * (Decimal(1) + values[1] * values[2] / Decimal(100))
        else:
            years = values[2]
            compounds = values[3]
            if years != years.to_integral_value() or not 0 <= years <= 1000:
                raise ValueError("years_must_be_integer_0_to_1000")
            if compounds != compounds.to_integral_value() or not 1 <= compounds <= 365:
                raise ValueError("compounds_per_year_must_be_integer_1_to_365")
            exponent = int(years) * int(compounds)
            if exponent > 10000:
                raise ValueError("compound_iteration_limit_exceeded")
            base = Decimal(1) + values[1] / (Decimal(100) * compounds)
            if base < 0:
                raise ValueError("compound_base_must_be_non_negative")
            calculated = values[0] * (base ** exponent)
        if not calculated.is_finite() or abs(calculated.adjusted()) > 256:
            raise ValueError("decimal_result_magnitude_exceeded")
        try:
            quantized = calculated.quantize(quantum, rounding=_DECIMAL_ROUNDING[rounding])
        except InvalidOperation as exc:
            raise ValueError("decimal_quantization_failed") from exc
    facts = {
        "operation": operation,
        "operands": list(operands),
        "currency": currency,
        "scale": scale,
        "rounding": rounding,
        "calculation_precision": _DECIMAL_CONTEXT_PRECISION,
        "calculation_value": _decimal_string(calculated),
        "result": format(quantized, f".{scale}f"),
    }
    return json.dumps(facts, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


_PROBABILITY_MAX_N = 4096
_PROBABILITY_MAX_RESULT_BITS = 65536


def _probability_project(value: Fraction, scale: Optional[int], rounding: Optional[str]) -> Optional[str]:
    if scale is None and rounding is None:
        return None
    if scale is None or rounding is None:
        raise ValueError("decimal_scale_and_rounding_must_be_supplied_together")
    if isinstance(scale, bool) or not isinstance(scale, int) or not 0 <= scale <= 18:
        raise ValueError("decimal_scale_out_of_range")
    if rounding not in _DECIMAL_ROUNDING:
        raise ValueError("unsupported_rounding_mode")
    with localcontext(Context(prec=_DECIMAL_CONTEXT_PRECISION, Emin=-999, Emax=999)):
        projected = Decimal(value.numerator) / Decimal(value.denominator)
        projected = projected.quantize(
            Decimal(1).scaleb(-scale), rounding=_DECIMAL_ROUNDING[rounding]
        )
    return format(projected, f".{scale}f")


@mcp.tool()
def exact_probability(
    operation: str,
    n: Optional[int] = None,
    k: Optional[int] = None,
    numerator: Optional[int] = None,
    denominator: Optional[int] = None,
    decimal_scale: Optional[int] = None,
    rounding: Optional[str] = None,
) -> str:
    """Return exact bounded combinatorics/probability facts using Fraction.

    Supported operations are ``fraction``, ``combination``, ``permutation``
    and ``binomial_probability``. For binomial probability ``numerator`` and
    ``denominator`` define the exact Bernoulli probability p.
    """
    if operation not in {"fraction", "combination", "permutation", "binomial_probability"}:
        raise ValueError("unsupported_probability_operation")
    for field_name, value in (("n", n), ("k", k), ("numerator", numerator), ("denominator", denominator)):
        if value is not None and (isinstance(value, bool) or not isinstance(value, int)):
            raise ValueError(f"{field_name}_must_be_integer")
    if operation == "fraction":
        if numerator is None or denominator is None:
            raise ValueError("fraction_requires_numerator_and_denominator")
        if denominator == 0:
            raise ValueError("denominator_must_not_be_zero")
        exact = Fraction(numerator, denominator)
    else:
        if n is None or k is None:
            raise ValueError(f"{operation}_requires_n_and_k")
        if not 0 <= n <= _PROBABILITY_MAX_N:
            raise ValueError("n_out_of_range")
        if not 0 <= k <= n:
            raise ValueError("k_must_be_between_zero_and_n")
        if operation == "combination":
            exact = Fraction(math.comb(n, k), 1)
        elif operation == "permutation":
            exact = Fraction(math.perm(n, k), 1)
        else:
            if numerator is None or denominator is None or denominator <= 0:
                raise ValueError("binomial_requires_positive_probability_denominator")
            probability = Fraction(numerator, denominator)
            if not 0 <= probability <= 1:
                raise ValueError("probability_must_be_between_zero_and_one")
            estimated_bits = max(1, denominator.bit_length()) * n + n
            if estimated_bits > _PROBABILITY_MAX_RESULT_BITS:
                raise ValueError("probability_cost_limit_exceeded")
            exact = (
                Fraction(math.comb(n, k), 1)
                * probability ** k
                * (1 - probability) ** (n - k)
            )
    if max(abs(exact.numerator).bit_length(), exact.denominator.bit_length()) > _PROBABILITY_MAX_RESULT_BITS:
        raise ValueError("probability_result_bit_limit_exceeded")
    facts = {
        "operation": operation,
        "n": n,
        "k": k,
        "probability_numerator": numerator,
        "probability_denominator": denominator,
        "result_numerator": exact.numerator,
        "result_denominator": exact.denominator,
        "fraction": f"{exact.numerator}/{exact.denominator}",
        "decimal_scale": decimal_scale,
        "rounding": rounding,
        "decimal": _probability_project(exact, decimal_scale, rounding),
    }
    return json.dumps(facts, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


_STRUCTURED_MAX_PAYLOAD_BYTES = 65536
_STRUCTURED_MAX_SCHEMA_BYTES = 32768
_STRUCTURED_MAX_DEPTH = 64
_STRUCTURED_MAX_NODES = 10000
_STRUCTURED_MAX_CSV_ROWS = 1000
_STRUCTURED_MAX_CSV_COLUMNS = 100
_STRUCTURED_MAX_FIELD_CHARS = 8192


def _bounded_tree(value: Any) -> tuple[int, int]:
    """Return depth/node count and reject recursive/oversized parsed trees."""
    stack = [(value, 1)]
    nodes = 0
    max_depth = 0
    seen: set[int] = set()
    while stack:
        current, depth = stack.pop()
        nodes += 1
        max_depth = max(max_depth, depth)
        if nodes > _STRUCTURED_MAX_NODES:
            raise ValueError("structured_node_limit_exceeded")
        if depth > _STRUCTURED_MAX_DEPTH:
            raise ValueError("structured_depth_limit_exceeded")
        if isinstance(current, dict):
            object_id = id(current)
            if object_id in seen:
                raise ValueError("structured_cycle_or_alias_rejected")
            seen.add(object_id)
            stack.extend((item, depth + 1) for pair in current.items() for item in pair)
        elif isinstance(current, (list, tuple)):
            object_id = id(current)
            if object_id in seen:
                raise ValueError("structured_cycle_or_alias_rejected")
            seen.add(object_id)
            stack.extend((item, depth + 1) for item in current)
    return max_depth, nodes


def _schema_contains_ref(value: Any) -> bool:
    if isinstance(value, dict):
        return "$ref" in value or any(_schema_contains_ref(item) for item in value.values())
    if isinstance(value, list):
        return any(_schema_contains_ref(item) for item in value)
    return False


def _validation_error(code: str, message: str, line: Optional[int] = None, column: Optional[int] = None, path: str = "") -> Dict[str, Any]:
    return {"code": code, "message": message[:300], "line": line, "column": column, "path": path[:300]}


@mcp.tool()
def structured_validate(
    format_name: str,
    payload: str,
    schema_json: Optional[str] = None,
    csv_dialect: Optional[str] = None,
) -> str:
    """Safely parse JSON/YAML/XML/CSV and optionally validate JSON Schema.

    The function never resolves remote references, DTDs, entities or
    XIncludes, never evaluates YAML tags or CSV formulas, and returns only
    hashes/diagnostics rather than echoing the supplied payload.
    """
    if format_name not in {"json", "yaml", "xml", "csv"}:
        raise ValueError("unsupported_structured_format")
    if not isinstance(payload, str):
        raise ValueError("payload_must_be_string")
    payload_bytes = payload.encode("utf-8")
    if len(payload_bytes) > _STRUCTURED_MAX_PAYLOAD_BYTES:
        raise ValueError("structured_payload_size_limit_exceeded")
    if schema_json is not None and format_name != "json":
        raise ValueError("schema_json_is_supported_only_for_json")
    errors: List[Dict[str, Any]] = []
    warnings: List[Dict[str, Any]] = []
    details: Dict[str, Any] = {}
    parsed: Any = None

    if format_name == "json":
        try:
            parsed = json.loads(payload)
        except json.JSONDecodeError as exc:
            errors.append(_validation_error("json_parse_error", exc.msg, exc.lineno, exc.colno))
        if not errors:
            depth, nodes = _bounded_tree(parsed)
            details.update({"depth": depth, "nodes": nodes})
        if schema_json is not None:
            if len(schema_json.encode("utf-8")) > _STRUCTURED_MAX_SCHEMA_BYTES:
                raise ValueError("structured_schema_size_limit_exceeded")
            try:
                schema = json.loads(schema_json)
                if not isinstance(schema, dict):
                    raise ValueError("schema_root_must_be_object")
                if _schema_contains_ref(schema):
                    raise ValueError("schema_ref_not_allowed")
                Draft202012Validator.check_schema(schema)
                if not errors:
                    for item in sorted(
                        Draft202012Validator(schema).iter_errors(parsed),
                        key=lambda error: (tuple(str(part) for part in error.absolute_path), error.message),
                    )[:50]:
                        errors.append(_validation_error(
                            "json_schema_error", item.message,
                            path="/" + "/".join(str(part) for part in item.absolute_path),
                        ))
            except json.JSONDecodeError as exc:
                errors.append(_validation_error("schema_parse_error", exc.msg, exc.lineno, exc.colno))
            except SchemaError as exc:
                errors.append(_validation_error("schema_invalid", exc.message))
            except ValueError as exc:
                errors.append(_validation_error(str(exc), str(exc)))

    elif format_name == "yaml":
        try:
            tokens = list(yaml.scan(payload))
            if any(token.__class__.__name__ in {"AliasToken", "AnchorToken", "TagToken"} for token in tokens):
                raise ValueError("yaml_alias_anchor_or_tag_rejected")
            parsed = yaml.safe_load(payload)
            depth, nodes = _bounded_tree(parsed)
            details.update({"depth": depth, "nodes": nodes})
        except yaml.YAMLError as exc:
            mark = getattr(exc, "problem_mark", None)
            errors.append(_validation_error(
                "yaml_parse_error", str(getattr(exc, "problem", None) or exc),
                getattr(mark, "line", -1) + 1 if mark else None,
                getattr(mark, "column", -1) + 1 if mark else None,
            ))
        except ValueError as exc:
            errors.append(_validation_error(str(exc), str(exc)))

    elif format_name == "xml":
        lowered_payload = payload.casefold()
        if any(token in lowered_payload for token in ("<!doctype", "<!entity", "<xi:include", "xinclude")):
            errors.append(_validation_error("xml_forbidden_construct", "DTD, entities and XInclude are disabled"))
        else:
            try:
                parsed = DefusedElementTree.fromstring(payload)
                stack = [(parsed, 1)]
                nodes = 0
                depth = 0
                while stack:
                    element, level = stack.pop()
                    nodes += 1
                    depth = max(depth, level)
                    if nodes > _STRUCTURED_MAX_NODES:
                        raise ValueError("structured_node_limit_exceeded")
                    if depth > _STRUCTURED_MAX_DEPTH:
                        raise ValueError("structured_depth_limit_exceeded")
                    stack.extend((child, level + 1) for child in list(element))
                details.update({"depth": depth, "nodes": nodes})
            except Exception as exc:
                errors.append(_validation_error("xml_parse_error", str(exc)))

    else:
        dialects = {"comma": ",", "semicolon": ";", "tab": "\t", "pipe": "|"}
        if csv_dialect not in dialects:
            raise ValueError("csv_dialect_must_be_explicit")
        try:
            reader = csv.reader(
                io.StringIO(payload, newline=""),
                delimiter=dialects[csv_dialect],
                strict=True,
            )
            expected_columns: Optional[int] = None
            row_count = 0
            formula_count = 0
            for row_count, row in enumerate(reader, 1):
                if row_count > _STRUCTURED_MAX_CSV_ROWS:
                    raise ValueError("csv_row_limit_exceeded")
                if len(row) > _STRUCTURED_MAX_CSV_COLUMNS:
                    raise ValueError("csv_column_limit_exceeded")
                if any(len(field) > _STRUCTURED_MAX_FIELD_CHARS for field in row):
                    raise ValueError("csv_field_size_limit_exceeded")
                if expected_columns is None:
                    expected_columns = len(row)
                elif len(row) != expected_columns:
                    errors.append(_validation_error(
                        "csv_column_count_mismatch",
                        f"row has {len(row)} columns; expected {expected_columns}",
                        line=row_count,
                    ))
                formula_count += sum(
                    field.lstrip().startswith(("=", "+", "-", "@")) for field in row
                )
            details.update({"rows": row_count, "columns": expected_columns or 0})
            if formula_count:
                warnings.append({"code": "csv_formula_prefix", "count": formula_count})
        except csv.Error as exc:
            errors.append(_validation_error("csv_parse_error", str(exc), line=getattr(reader, "line_num", None)))
        except ValueError as exc:
            errors.append(_validation_error(str(exc), str(exc)))

    schema_hash = hashlib.sha256(schema_json.encode("utf-8")).hexdigest() if schema_json is not None else None
    facts = {
        "valid": not errors,
        "format": format_name,
        "payload_hash": hashlib.sha256(payload_bytes).hexdigest(),
        "schema_hash": schema_hash,
        "errors": errors[:50],
        "warnings": warnings[:50],
        "details": details,
    }
    return json.dumps(facts, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


_RUST_COMPILE_SANDBOX_URL = os.getenv(
    "RUST_COMPILE_SANDBOX_URL", "http://rust-compile-sandbox:8080"
)
_RUST_COMPILE_MAX_SOURCE_CHARS = 200_000
_RUST_COMPILE_HTTP_TIMEOUT_S = 15.0


@mcp.tool()
async def rust_compile_check(source: str, edition: str = "2021") -> str:
    """Type/borrow-check Rust source in an isolated, network-free sandbox.

    Runs `rustc --emit=metadata` (analysis only -- no codegen, no linking,
    the code is never executed) against the given source and returns
    structured diagnostics. Use this to verify a Rust answer actually
    compiles before presenting it as correct; it catches lifetime,
    ownership, interior-mutability and type errors deterministically,
    which LLM self-review misses or catches inconsistently.
    """
    if not isinstance(source, str) or not source.strip():
        raise ValueError("source_must_be_non_empty_string")
    if len(source) > _RUST_COMPILE_MAX_SOURCE_CHARS:
        raise ValueError("source_exceeds_size_limit")
    if edition not in {"2015", "2018", "2021", "2024"}:
        raise ValueError("unsupported_edition")
    try:
        async with httpx.AsyncClient(timeout=_RUST_COMPILE_HTTP_TIMEOUT_S) as client:
            resp = await client.post(
                f"{_RUST_COMPILE_SANDBOX_URL}/compile-check",
                json={"source": source, "edition": edition},
            )
            resp.raise_for_status()
            result = resp.json()
    except Exception as exc:
        logger.warning(f"rust_compile_check sandbox call failed: {exc}")
        return json.dumps(
            {"compiles": None, "diagnostics": [], "duration_ms": 0, "sandbox_error": str(exc)[:300]},
            ensure_ascii=False, sort_keys=True, separators=(",", ":"),
        )
    facts = {
        "compiles": result.get("compiles"),
        "diagnostics": (result.get("diagnostics") or [])[:50],
        "duration_ms": result.get("duration_ms"),
        "timed_out": result.get("timed_out", False),
        "source_hash": hashlib.sha256(source.encode("utf-8")).hexdigest(),
    }
    return json.dumps(facts, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


@mcp.tool()
def day_of_week(date_str: str) -> str:
    """
    Returns weekday, calendar week and day of year for a date.
    Format: YYYY-MM-DD.
    Example: day_of_week('2026-12-25')
    """
    try:
        facts = _calendar_facts_payload(date_str, "en")
        return (
            f"{facts['date']} is a {facts['weekday_name']} "
            f"(CW {facts['iso_week']}, day {facts['day_of_year']} "
            f"of year {facts['year']})"
        )
    except Exception as e:
        return f"Error: {e}"


@mcp.tool()
def unit_convert(value: float, from_unit: str, to_unit: str) -> str:
    """
    Converts physical units exactly via pint.
    Examples: unit_convert(100, 'km/h', 'm/s'), unit_convert(1, 'mile', 'km'),
               unit_convert(100, 'degF', 'degC'), unit_convert(5, 'lb', 'kg')
    """
    try:
        from pint import UnitRegistry
        ureg = UnitRegistry()
        qty = value * ureg(from_unit)
        result = qty.to(to_unit)
        return f"{value} {from_unit} = {result.magnitude:.10g} {to_unit}"
    except Exception as e:
        return f"Error: {e}"


@mcp.tool()
def statistics_calc(data: str, operations: str = "mean,median,stdev,min,max,sum,count") -> str:
    """
    Calculates statistical measures for a data set.
    data: comma-separated numbers, e.g. '1,2,3,4,5,6,7,8,9,10'
    operations: comma-separated operations (mean, median, stdev, variance, min, max, sum, count, mode)
    """
    try:
        values = [float(x.strip()) for x in data.split(",") if x.strip()]
        if not values:
            return "Error: No data found"
        ops = [o.strip() for o in operations.split(",")]
        op_map = {
            "mean": stats_module.mean,
            "median": stats_module.median,
            "stdev": lambda v: stats_module.stdev(v) if len(v) > 1 else 0.0,
            "variance": lambda v: stats_module.variance(v) if len(v) > 1 else 0.0,
            "min": min,
            "max": max,
            "sum": sum,
            "count": len,
            "mode": stats_module.mode,
        }
        results = {}
        for op in ops:
            if op in op_map:
                val = op_map[op](values)
                results[op] = round(float(val), 10)
        preview = data[:40] + "..." if len(data) > 40 else data
        return f"Statistics [{preview}]: {json.dumps(results, ensure_ascii=False)}"
    except Exception as e:
        return f"Error: {e}"


@mcp.tool()
def hash_text(text: str, algorithm: str = "sha256") -> str:
    """
    Calculates cryptographic hashes exactly.
    algorithm: md5, sha1, sha224, sha256, sha384, sha512
    Example: hash_text('Hello World', 'sha256')
    """
    try:
        algos = {
            "md5": hashlib.md5, "sha1": hashlib.sha1,
            "sha224": hashlib.sha224, "sha256": hashlib.sha256,
            "sha384": hashlib.sha384, "sha512": hashlib.sha512,
        }
        algo = algorithm.lower()
        if algo not in algos:
            return f"Unknown algorithm. Available: {list(algos.keys())}"
        digest = algos[algo](text.encode("utf-8")).hexdigest()
        preview = text[:30] + "..." if len(text) > 30 else text
        return f"{algorithm.upper()}('{preview}') = {digest}"
    except Exception as e:
        return f"Error: {e}"


@mcp.tool()
def base64_codec(text: str, mode: str = "encode") -> str:
    """
    Base64 encode or decode.
    mode: 'encode' or 'decode'
    Example: base64_codec('Hello World', 'encode')
    """
    try:
        if mode == "encode":
            result = base64.b64encode(text.encode("utf-8")).decode("ascii")
            return f"Base64-encoded: {result}"
        elif mode == "decode":
            result = base64.b64decode(text.encode("ascii")).decode("utf-8")
            return f"Base64-decoded: {result}"
        else:
            return "Error: mode must be 'encode' or 'decode'"
    except Exception as e:
        return f"Error: {e}"


@mcp.tool()
def regex_extract(pattern: str, text: str, flags: str = "") -> str:
    """
    Performs regex pattern matching and extracts all matches.
    flags: i=ignorecase, m=multiline, s=dotall (combinable, e.g. 'im')
    Example: regex_extract(r'\\d{4}-\\d{2}-\\d{2}', 'Date: 2026-03-29 and 2025-12-31')
    """
    try:
        flag_map = {"i": re.IGNORECASE, "m": re.MULTILINE, "s": re.DOTALL}
        re_flags = 0
        for f in flags.lower():
            re_flags |= flag_map.get(f, 0)
        matches = re.findall(pattern, text, re_flags)
        groups_count = len(matches[0]) if matches and isinstance(matches[0], tuple) else 0
        return (
            f"Pattern '{pattern}' → {len(matches)} matches: "
            f"{matches[:20]}{'...' if len(matches) > 20 else ''}"
        )
    except Exception as e:
        return f"Error: {e}"


@mcp.tool()
def subnet_calc(cidr: str) -> str:
    """
    Calculates network information for CIDR notation.
    Example: subnet_calc('192.168.1.0/24'), subnet_calc('10.0.0.5/22')
    """
    try:
        network = ipaddress.ip_network(cidr, strict=False)
        hosts = list(network.hosts())
        host_count = len(hosts)
        return (
            f"Network: {network.network_address}/{network.prefixlen}\n"
            f"Broadcast: {network.broadcast_address}\n"
            f"Subnet mask: {network.netmask} ({network.prefixlen} bits)\n"
            f"Usable hosts: {host_count}\n"
            f"First host IP: {hosts[0] if hosts else 'N/A'}\n"
            f"Last host IP: {hosts[-1] if hosts else 'N/A'}\n"
            f"Version: IPv{network.version}"
        )
    except Exception as e:
        return f"Error: {e}"


@mcp.tool()
def text_analyze(text: str) -> str:
    """
    Analyzes text for words, characters, sentences, paragraphs, reading time.
    Useful for precise text metrics without LLM estimation.
    """
    words = len(re.findall(r"\b\w+\b", text))
    chars_total = len(text)
    chars_no_space = len(text.replace(" ", "").replace("\n", "").replace("\t", ""))
    sentences = len(re.findall(r"[.!?]+", text))
    paragraphs = len([p for p in text.split("\n") if p.strip()])
    unique_words = len(set(re.findall(r"\b\w+\b", text.lower())))
    read_min = words / 200.0
    return (
        f"Words: {words} ({unique_words} unique), "
        f"Characters: {chars_total} ({chars_no_space} without whitespace), "
        f"Sentences: {sentences}, paragraphs: {paragraphs}, "
        f"Reading time: ~{read_min:.1f} min (200 WPM)"
    )


@mcp.tool()
def prime_factorize(n: int) -> str:
    """
    Calculates the prime factorization of an integer exactly via SymPy.
    Example: prime_factorize(360) → 2^3 × 3^2 × 5
    """
    try:
        if n <= 1:
            return f"{n} has no prime factors (n > 1 required)"
        factors = sp.factorint(n)
        parts = [f"{p}^{e}" if e > 1 else str(p) for p, e in sorted(factors.items())]
        return f"{n} = {' × '.join(parts)}"
    except Exception as e:
        return f"Error: {e}"


@mcp.tool()
def gcd_lcm(a: int, b: int, operation: str = "both") -> str:
    """
    Calculates GCD (greatest common divisor) and LCM (least common multiple).
    operation: 'gcd', 'lcm', or 'both'
    Example: gcd_lcm(48, 18)
    """
    try:
        g = math.gcd(abs(a), abs(b))
        lv = abs(a * b) // g if g != 0 else 0
        if operation == "gcd":
            return f"GCD({a}, {b}) = {g}"
        if operation == "lcm":
            return f"LCM({a}, {b}) = {lv}"
        return f"GCD({a}, {b}) = {g}  |  LCM({a}, {b}) = {lv}"
    except Exception as e:
        return f"Error: {e}"


@mcp.tool()
def json_query(json_str: str, path: str) -> str:
    """
    Extracts data from JSON via dot notation or array index.
    Example: json_query('{"user":{"name":"Alice","age":30}}', 'user.name') → 'Alice'
    Supports: 'key', 'key.subkey', 'key[0]', 'key[0].subkey'
    """
    try:
        data = json.loads(json_str)
        tokens = re.split(r"\.(?![^\[]*\])", path)
        current = data
        for token in tokens:
            arr_match = re.match(r"^(\w+)\[(\d+)\]$", token)
            if arr_match:
                current = current[arr_match.group(1)][int(arr_match.group(2))]
            elif re.match(r"^\[(\d+)\]$", token):
                current = current[int(re.match(r"^\[(\d+)\]$", token).group(1))]
            else:
                current = current[token]
        return f"'{path}' → {json.dumps(current, ensure_ascii=False)}"
    except Exception as e:
        return f"Error: {e}"


@mcp.tool()
def roman_numeral(value: str) -> str:
    """
    Converts between Arabic and Roman numerals.
    Input: number (1-3999) or Roman numeral (e.g. 'XIV', 'MMXXVI').
    """
    try:
        roman_vals = [
            (1000, "M"), (900, "CM"), (500, "D"), (400, "CD"),
            (100, "C"), (90, "XC"), (50, "L"), (40, "XL"),
            (10, "X"), (9, "IX"), (5, "V"), (4, "IV"), (1, "I"),
        ]
        roman_parse = {"I": 1, "V": 5, "X": 10, "L": 50, "C": 100, "D": 500, "M": 1000}
        v = value.strip()
        if v.isdigit():
            n = int(v)
            if not 1 <= n <= 3999:
                return "Error: Only 1–3999 supported"
            result = ""
            for val, sym in roman_vals:
                while n >= val:
                    result += sym
                    n -= val
            return f"{value} → {result}"
        else:
            s = v.upper()
            n = 0
            for i, c in enumerate(s):
                if c not in roman_parse:
                    return f"Invalid Roman numeral: '{c}'"
                if i + 1 < len(s) and roman_parse[c] < roman_parse[s[i + 1]]:
                    n -= roman_parse[c]
                else:
                    n += roman_parse[c]
            return f"{value} → {n}"
    except Exception as e:
        return f"Error: {e}"


# ─── LEGAL TOOLS ─────────────────────────────────────────────────────────────

# Configurable via GII_BASE_URL env var — set to a local mirror for offline deployments
_GII_BASE       = os.getenv("GII_BASE_URL", "https://www.gesetze-im-internet.de")
_GII_TOC_URL    = f"{_GII_BASE}/gii-toc.xml"
_TOC_CACHE: dict = {"data": None, "ts": 0.0}
_TOC_LOCK        = threading.Lock()
_LAW_CACHE: dict = {}
_LAW_LOCK        = threading.Lock()
_TOC_CACHE_TTL   = 24 * 3600
_LAW_CACHE_TTL   = 6  * 3600


def _strip_ns(tag: str) -> str:
    """Removes XML namespace prefix from a tag name."""
    return re.sub(r"\{[^}]*\}", "", tag)


def _extract_text(element) -> str:
    """Recursively extracts all text from an XML element."""
    parts: list[str] = []
    if element.text and element.text.strip():
        parts.append(element.text.strip())
    for child in element:
        child_text = _extract_text(child)
        if child_text:
            parts.append(child_text)
        if child.tail and child.tail.strip():
            parts.append(child.tail.strip())
    return " ".join(parts)


def _get_toc() -> Optional[ET.Element]:
    """Loads and caches the GII table-of-contents XML (gii-toc.xml)."""
    with _TOC_LOCK:
        now = time.time()
        if _TOC_CACHE["data"] is not None and now - _TOC_CACHE["ts"] < _TOC_CACHE_TTL:
            return _TOC_CACHE["data"]
        try:
            resp = httpx.get(_GII_TOC_URL, timeout=8.0, follow_redirects=True)
            resp.raise_for_status()
            root = ET.fromstring(resp.content)
            _TOC_CACHE["data"] = root
            _TOC_CACHE["ts"]   = now
            return root
        except Exception:
            return _TOC_CACHE["data"]  # Stale data as fallback


def _resolve_law_url(law: str) -> str:
    """Determines the ZIP URL for a law — from TOC (exact) or fallback (lowercase)."""
    toc = _TOC_CACHE.get("data")
    if toc is not None:
        law_up = law.upper()
        for item in toc:
            if _strip_ns(item.tag) != "item":
                continue
            link_url = ""
            for child in item:
                if _strip_ns(child.tag) == "link" and child.text:
                    link_url = child.text.strip()
            if link_url:
                m = re.search(r"/([^/]+)/xml\.zip", link_url)
                if m and m.group(1).upper() == law_up:
                    return link_url
    return f"{_GII_BASE}/{law.lower()}/xml.zip"


def _get_law_xml(law: str) -> Tuple[Optional[ET.Element], str]:
    """Loads and caches the XML ZIP of a federal law. Returns (root, error)."""
    key = law.upper()
    with _LAW_LOCK:
        now = time.time()
        if key in _LAW_CACHE and now - _LAW_CACHE[key]["ts"] < _LAW_CACHE_TTL:
            return _LAW_CACHE[key]["root"], ""
        # LRU eviction when > 20 cached laws
        if len(_LAW_CACHE) >= 20:
            oldest = min(_LAW_CACHE, key=lambda k: _LAW_CACHE[k]["ts"])
            del _LAW_CACHE[oldest]
        url = _resolve_law_url(law)
        try:
            resp = httpx.get(url, timeout=15.0, follow_redirects=True)
            if resp.status_code == 404:
                # Fallback: retry without TOC
                fallback_url = f"{_GII_BASE}/{law.lower()}/xml.zip"
                if fallback_url != url:
                    resp2 = httpx.get(fallback_url, timeout=15.0, follow_redirects=True)
                    if resp2.status_code == 200:
                        resp = resp2
                    else:
                        return None, (
                            f"Law '{law}' not found. "
                            f"Use legal_search_laws() to find valid abbreviations."
                        )
                else:
                    return None, (
                        f"Law '{law}' not found. "
                        f"Use legal_search_laws() to find valid abbreviations."
                    )
            resp.raise_for_status()
            with zipfile.ZipFile(io.BytesIO(resp.content)) as zf:
                xml_bytes = zf.read(zf.namelist()[0])
            root = ET.fromstring(xml_bytes)
            _LAW_CACHE[key] = {"root": root, "ts": now}
            return root, ""
        except httpx.TimeoutException:
            return None, "Error: Timeout loading from gesetze-im-internet.de — please try again."
        except Exception as e:
            return None, f"Error loading '{law}': {e}"


def _get_enbez_num(enbez: str) -> str:
    """Extracts the number+letter from an enbez like '§ 242a' → '242a'."""
    m = re.search(r"(\d+[a-z]?)", enbez.lower())
    return m.group(1) if m else ""


def _iter_norms(root: ET.Element):
    """Iterator over all <norm> elements with (enbez, titel, textdaten_element)."""
    for elem in root.iter():
        if _strip_ns(elem.tag) != "norm":
            continue
        meta = None
        textdaten = None
        for child in elem:
            tag = _strip_ns(child.tag)
            if tag == "metadaten":
                meta = child
            elif tag == "textdaten":
                textdaten = child
        if meta is None:
            continue
        enbez = ""
        titel = ""
        for child in meta:
            tag = _strip_ns(child.tag)
            if tag == "enbez" and child.text:
                enbez = child.text.strip()
            elif tag == "titel" and child.text:
                titel = child.text.strip()
        yield enbez, titel, textdaten


def _norm_text(textdaten) -> str:
    """Extracts the full text from a <textdaten> element."""
    if textdaten is None:
        return ""
    for child in textdaten:
        if _strip_ns(child.tag) == "text":
            for content in child:
                return _extract_text(content)
    return ""


@mcp.tool()
def legal_search_laws(query: str, max_results: int = 10) -> str:
    """
    Searches German federal laws (~6,000 laws) by keyword in abbreviation or title.
    Returns abbreviations to be used with legal_get_paragraph / legal_get_law_overview.
    Examples: legal_search_laws("Mietrecht"), legal_search_laws("Datenschutz"), legal_search_laws("BGB")
    """
    toc = _get_toc()
    if toc is None:
        return "Error: gesetze-im-internet.de unreachable — please try again."

    q = query.lower()
    matches: list[tuple[str, str]] = []
    for item in toc:
        if _strip_ns(item.tag) != "item":
            continue
        titel = link_url = ""
        for child in item:
            tag = _strip_ns(child.tag)
            if tag == "title" and child.text:
                titel = child.text.strip()
            elif tag == "link" and child.text:
                link_url = child.text.strip()
        # Extract abbreviation from URL: .../bgb/xml.zip → BGB
        abk = ""
        if link_url:
            m = re.search(r"/([^/]+)/xml\.zip", link_url)
            if m:
                abk = m.group(1).upper()
        if not abk:
            continue
        if q in abk.lower() or q in titel.lower():
            matches.append((abk, titel, link_url))

    if not matches:
        return (
            f"No laws found for '{query}'.\n"
            f"Tip: Search with abbreviation (e.g. 'Miet', 'Datenschutz', 'Arbeits', 'GG', 'BGB')."
        )

    shown = matches[:max_results]
    rest  = len(matches) - len(shown)
    lines = [f"Results for '{query}' ({len(matches)} laws found):"]
    for abk, t, _ in shown:
        lines.append(f"  - {abk}: {t}")
    if rest > 0:
        lines.append(f"  [... {rest} more — increase max_results or refine search]")
    lines.append(f"\nNext steps:")
    lines.append(f'  → legal_get_law_overview("{shown[0][0]}")   — show table of contents')
    lines.append(f'  → legal_get_paragraph("{shown[0][0]}", "§-number")  — retrieve law text')
    return "\n".join(lines)


@mcp.tool()
def legal_get_law_overview(law: str, max_entries: int = 50) -> str:
    """
    Shows the table of contents / structure of a German federal law.
    law: Abbreviation (e.g. 'BGB', 'GG', 'StGB', 'DSGVO', 'HGB', 'ZPO', 'ArbSchG')
    Useful when the searched paragraph is unknown — shows all §§ with titles.
    Examples: legal_get_law_overview("GG"), legal_get_law_overview("BGB")
    """
    root, err = _get_law_xml(law)
    if root is None:
        return err

    law_title = law.upper()
    entries: list[tuple[str, str]] = []

    for enbez, titel, _ in _iter_norms(root):
        if enbez:
            entries.append((enbez, titel))
        # Law long title from first norm without enbez
        elif not enbez and not entries:
            pass  # extracted below from <langue>

    # Long title from metadata of the first norm
    for elem in root.iter():
        if _strip_ns(elem.tag) != "norm":
            continue
        for child in elem:
            if _strip_ns(child.tag) == "metadaten":
                for mc in child:
                    if _strip_ns(mc.tag) == "langue" and mc.text:
                        law_title = f"{mc.text.strip()} ({law.upper()})"
                        break
        break

    if not entries:
        return f"Law '{law}' loaded, but no paragraphs/articles found."

    total = len(entries)
    shown = entries[:max_entries]
    lines = [f"{law_title} — Table of contents ({total} norms, showing {len(shown)}):"]
    for enbez, titel in shown:
        if titel:
            lines.append(f"  {enbez:<14} {titel}")
        else:
            lines.append(f"  {enbez}")
    if total > max_entries:
        lines.append(f"\n[... {total - max_entries} more norms not shown]")
        lines.append(f'→ Use legal_fulltext_search("{law.upper()}", "keyword") for targeted search')
    lines.append(f'→ Use legal_get_paragraph("{law.upper()}", "§-number") for full text')
    return "\n".join(lines)


@mcp.tool()
def legal_get_paragraph(law: str, paragraph: str) -> str:
    """
    Retrieves the exact law text of a paragraph or article from a German federal law.
    law: Abbreviation (e.g. 'BGB', 'GG', 'StGB', 'DSGVO', 'HGB', 'ZPO')
    paragraph: §/Art. number as string (e.g. '242', '823', '1' for Art. 1 GG, '13' for Art. 13 GG)
    Returns the official wording directly from gesetze-im-internet.de.
    Examples: legal_get_paragraph("BGB", "242"), legal_get_paragraph("GG", "1")
    """
    root, err = _get_law_xml(law)
    if root is None:
        return err

    para_norm = paragraph.lower().strip().lstrip("§").lstrip("art").strip()

    for enbez, titel, textdaten in _iter_norms(root):
        if not enbez:
            continue
        if _get_enbez_num(enbez) == para_norm:
            text = _norm_text(textdaten)
            if not text:
                return (
                    f"{enbez} {law.upper()}"
                    + (f" — {titel}" if titel else "")
                    + "\n\n(No text content — possibly repealed or empty norm)"
                )
            if len(text) > 2200:
                text = text[:2200] + "\n[... Text truncated — full text at gesetze-im-internet.de]"
            header = f"{enbez} {law.upper()}" + (f" — {titel}" if titel else "")
            return f"{header}\n\nSource: gesetze-im-internet.de/{law.lower()}\n\n{text}"

    return (
        f"§ {paragraph} {law.upper()} not found.\n"
        f'Tip: Use legal_get_law_overview("{law.upper()}") to see all available §§.'
    )


@mcp.tool()
def legal_fulltext_search(law: str, query: str, max_results: int = 5) -> str:
    """
    Full-text search within a German federal law by keyword.
    law: Abbreviation (e.g. 'BGB', 'StGB', 'DSGVO', 'HGB')
    query: Search term (e.g. 'Treu und Glauben', 'Einwilligung', 'Schadensersatz')
    Returns matching paragraphs with text excerpt.
    Examples: legal_fulltext_search("BGB", "Treu und Glauben"), legal_fulltext_search("DSGVO", "Einwilligung")
    """
    root, err = _get_law_xml(law)
    if root is None:
        return err

    q_lower  = query.lower()
    results: list[tuple[str, str, str]] = []

    for enbez, titel, textdaten in _iter_norms(root):
        if not enbez:
            continue
        text = _norm_text(textdaten)
        if not text:
            continue
        text_lower = text.lower()
        if q_lower not in text_lower:
            continue
        idx   = text_lower.find(q_lower)
        start = max(0, idx - 80)
        end   = min(len(text), idx + len(query) + 120)
        pre   = "..." if start > 0 else ""
        post  = "..." if end < len(text) else ""
        snippet = pre + text[start:end] + post
        results.append((enbez, titel, snippet))

    if not results:
        return f"No matches for '{query}' in {law.upper()}."

    shown = results[:max_results]
    rest  = len(results) - len(shown)
    lines = [f"Full-text search '{query}' in {law.upper()} — {len(results)} matches:"]
    for enbez, titel, snippet in shown:
        lines.append("")
        lines.append(enbez + (f" — {titel}" if titel else ""))
        lines.append(f'  „{snippet}"')
    if rest > 0:
        lines.append(f"\n[... {rest} more results — increase max_results]")
    if shown:
        num = _get_enbez_num(shown[0][0])
        lines.append(f'\n→ Use legal_get_paragraph("{law.upper()}", "{num}") for full text')
    return "\n".join(lines)


# ─── GRAPH RAG TOOLS (Neo4j) ─────────────────────────────────────────────────

_graph_manager: Optional[Any] = None  # GraphRAGManager instance, lazy-initialized


async def _get_graph_manager():
    """Returns the GraphRAGManager, initializing it on first call."""
    global _graph_manager
    if _graph_manager is not None:
        return _graph_manager
    neo4j_uri  = os.getenv("NEO4J_URI",  "bolt://neo4j-knowledge:7687")
    neo4j_user = os.getenv("NEO4J_USER", "neo4j")
    neo4j_pass = os.getenv("NEO4J_PASS")
    if not neo4j_uri:
        return None
    try:
        from graph_rag.manager import GraphRAGManager
        mgr = GraphRAGManager(neo4j_uri, neo4j_user, neo4j_pass)
        await mgr.setup()
        _graph_manager = mgr
        logger.info("✅ GraphRAGManager initialized in MCP server")
    except Exception as e:
        logger.warning(f"⚠️ GraphRAGManager init failed: {e}")
    return _graph_manager


@mcp.tool()
async def graph_query(query: str, categories: Optional[List[str]] = None) -> str:
    """
    Searches the Neo4j knowledge graph for entities and relations related to the query.
    Returns structured context with provenance metadata.
    categories: Optional list of expert categories for domain filtering
    (e.g. ['technical_support', 'code_reviewer'])
    """
    mgr = await _get_graph_manager()
    if mgr is None:
        return "Error: Knowledge graph unavailable (Neo4j not configured)"
    try:
        result = await mgr.query_context(query, categories=categories or [])
        return result or "No relevant entries found in the knowledge graph."
    except Exception as e:
        logger.warning(f"graph_query error: {e}")
        return f"Knowledge graph query error: {e}"


@mcp.tool()
async def graph_ingest(
    question: str,
    answer: str,
    domain: str = "general",
    source_model: str = "external",
    confidence: float = 0.7,
) -> str:
    """
    Stores facts from a question-answer pair in the Neo4j knowledge graph.
    Extracts entities and relations and inserts them with provenance metadata.
    Useful for external agents (Cursor, Claude Desktop) to ingest knowledge.
    """
    mgr = await _get_graph_manager()
    if mgr is None:
        return "Error: Knowledge graph unavailable (Neo4j not configured)"
    try:
        # No LLM available for external ingests — use heuristics directly
        # Create a minimal stub LLM that returns simple JSON text
        class _StubLLM:
            async def ainvoke(self, prompt: str):
                class _R:
                    content = "[]"
                return _R()

        await mgr.extract_and_ingest(
            question, answer, _StubLLM(),
            domain=domain, source_model=source_model, confidence=confidence,
        )
        return f"Ingest started: domain={domain}, source_model={source_model}, confidence={confidence}"
    except Exception as e:
        logger.warning(f"graph_ingest error: {e}")
        return f"Knowledge graph ingest error: {e}"


@mcp.tool()
async def graph_provenance(entity_name: str) -> str:
    """
    Returns the complete version history of all relations for an entity.
    Useful for analyzing contradictions (Model A said X, Model B corrected to Y).
    """
    mgr = await _get_graph_manager()
    if mgr is None:
        return "Error: Knowledge graph unavailable (Neo4j not configured)"
    try:
        records = await mgr.get_provenance(entity_name)
        if not records:
            return f"No relations found for entity '{entity_name}'."
        lines = [f"[Provenance: {entity_name}]"]
        for r in records:
            line = f"• {r['relation']} → {r['target']} | v{r.get('version', '?')}"
            line += f" | Source: {r.get('source_model', '?')}"
            conf = r.get("confidence")
            if conf is not None:
                line += f" | Confidence: {conf:.0%}"
            if r.get("superseded_version"):
                line += f" | (supersedes v{r['superseded_version']}, was: {r.get('prev_source_model', '?')})"
            lines.append(line)
        return "\n".join(lines)
    except Exception as e:
        logger.warning(f"graph_provenance error: {e}")
        return f"Provenance query error: {e}"


# ─── CODE-NAVIGATION TOOLS (Agentic Coder) ──────────────────────────────────

# Workspace root: file access is restricted to this directory.
_CODE_WORKSPACE = Path(os.getenv("CODE_WORKSPACE", "/app/workspace")).resolve()

def _resolve_safe_path(raw: str) -> Path:
    """Resolves a path and ensures it is within _CODE_WORKSPACE."""
    p = Path(raw)
    if not p.is_absolute():
        p = _CODE_WORKSPACE / p
    resolved = p.resolve()
    if not str(resolved).startswith(str(_CODE_WORKSPACE)):
        raise PermissionError(
            f"Path '{resolved}' is outside the allowed workspace '{_CODE_WORKSPACE}'."
        )
    return resolved


# Language-specific regex patterns for repo_map (non-Python files)
_LANG_PATTERNS: Dict[str, List[re.Pattern]] = {
    ".js":   [re.compile(r"^\s*(?:export\s+)?(?:async\s+)?function\s+(\w+)", re.M),
              re.compile(r"^\s*(?:export\s+)?class\s+(\w+)", re.M)],
    ".ts":   [re.compile(r"^\s*(?:export\s+)?(?:async\s+)?function\s+(\w+)", re.M),
              re.compile(r"^\s*(?:export\s+)?class\s+(\w+)", re.M),
              re.compile(r"^\s*(?:export\s+)?(?:abstract\s+)?class\s+(\w+)", re.M)],
    ".go":   [re.compile(r"^func\s+(?:\(\w+\s+\*?\w+\)\s+)?(\w+)\s*\(", re.M),
              re.compile(r"^type\s+(\w+)\s+struct", re.M)],
    ".rs":   [re.compile(r"^\s*(?:pub\s+)?(?:async\s+)?fn\s+(\w+)", re.M),
              re.compile(r"^\s*(?:pub\s+)?struct\s+(\w+)", re.M),
              re.compile(r"^\s*(?:pub\s+)?trait\s+(\w+)", re.M)],
    ".java": [re.compile(r"(?:public|private|protected|static|\s)+[\w<>\[\]]+\s+(\w+)\s*\(", re.M),
              re.compile(r"^\s*(?:public\s+)?(?:abstract\s+)?class\s+(\w+)", re.M)],
    ".cpp":  [re.compile(r"^\w[\w:*&<>\s]+\s+(\w+)\s*\([^)]*\)\s*(?:const\s*)?\{", re.M)],
    ".c":    [re.compile(r"^\w[\w*\s]+\s+(\w+)\s*\([^)]*\)\s*\{", re.M)],
}

_SKIP_DIRS = {".git", "__pycache__", "node_modules", ".venv", "venv", "dist", "build", ".mypy_cache", ".pytest_cache"}


def _map_python_file(path: Path) -> List[str]:
    """Extracts classes/functions from a Python file via AST."""
    try:
        source = path.read_text(encoding="utf-8", errors="replace")
        tree = ast.parse(source)
    except SyntaxError:
        return ["<parse error>"]
    symbols: List[str] = []
    for node in ast.walk(tree):
        if isinstance(node, ast.ClassDef):
            methods = [
                n.name for n in ast.walk(node)
                if isinstance(n, (ast.FunctionDef, ast.AsyncFunctionDef)) and n.col_offset > 0
            ]
            symbols.append(f"class {node.name}: [{', '.join(methods[:8])}{'…' if len(methods) > 8 else ''}]")
    for node in ast.iter_child_nodes(tree):
        if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)):
            symbols.append(f"def {node.name}()")
    return symbols


def _map_lang_file(path: Path, ext: str) -> List[str]:
    """Extracts classes/functions from non-Python files via regex."""
    patterns = _LANG_PATTERNS.get(ext, [])
    if not patterns:
        return []
    try:
        source = path.read_text(encoding="utf-8", errors="replace")
    except OSError:
        return []
    seen: set = set()
    symbols: List[str] = []
    for pat in patterns:
        for m in pat.finditer(source):
            name = m.group(1)
            if name and name not in seen:
                seen.add(name)
                symbols.append(name)
    return symbols


@mcp.tool()
def repo_map(path: str = ".", max_depth: int = 3) -> str:
    """
    Returns a compact skeleton view of a directory: file paths +
    class/function names (without code). Ideal for context-limited SLMs to
    navigate a repo before reading files.
    Supported languages: Python (.py), JavaScript (.js), TypeScript (.ts),
    Go (.go), Rust (.rs), Java (.java), C/C++ (.c/.cpp).
    Example: repo_map("src/", max_depth=2)
    """
    try:
        root = _resolve_safe_path(path)
    except PermissionError as e:
        return f"ERROR: {e}"
    if not root.exists():
        return f"ERROR: Path does not exist: '{path}'"
    if not root.is_dir():
        return f"ERROR: '{path}' is not a directory. For files: read_file_chunked."

    supported_exts = {".py"} | set(_LANG_PATTERNS.keys())
    lines: List[str] = []
    total_files = 0

    def walk(directory: Path, depth: int, prefix: str) -> None:
        nonlocal total_files
        if depth > max_depth:
            return
        try:
            entries = sorted(directory.iterdir(), key=lambda e: (e.is_file(), e.name.lower()))
        except PermissionError:
            return
        for entry in entries:
            if entry.name in _SKIP_DIRS or entry.name.startswith("."):
                continue
            rel = entry.relative_to(root)
            if entry.is_dir():
                lines.append(f"{prefix}📁 {rel}/")
                walk(entry, depth + 1, prefix + "  ")
            elif entry.is_file() and entry.suffix.lower() in supported_exts:
                total_files += 1
                if total_files > 200:
                    lines.append(f"{prefix}… (too many files, reduce max_depth)")
                    return
                ext = entry.suffix.lower()
                if ext == ".py":
                    symbols = _map_python_file(entry)
                else:
                    symbols = _map_lang_file(entry, ext)
                sym_str = f" → {', '.join(symbols[:6])}{'…' if len(symbols) > 6 else ''}" if symbols else ""
                lines.append(f"{prefix}📄 {rel}{sym_str}")

    lines.append(f"# Repo-Map: {root} (max_depth={max_depth})")
    walk(root, 1, "")
    lines.append(f"\n# {total_files} file(s) indexed. Use read_file_chunked for details.")
    return "\n".join(lines)


@mcp.tool()
def read_file_chunked(file_path: str, start_line: int = 1, end_line: int = 50) -> str:
    """
    Reads a slice (start_line to end_line) from a file with line numbers.
    Prevents SLMs from loading entire files into context.
    Lines are 1-based. end_line=0 → read to end of file (max 200 lines).
    Example: read_file_chunked("src/main.py", start_line=10, end_line=50)
    """
    try:
        fpath = _resolve_safe_path(file_path)
    except PermissionError as e:
        return f"ERROR: {e}"
    if not fpath.exists():
        return f"ERROR: File not found: '{file_path}'"
    if not fpath.is_file():
        return f"ERROR: '{file_path}' is not a file. For directories: repo_map."

    MAX_LINES = 200
    if end_line == 0:
        end_line = start_line + MAX_LINES - 1
    if start_line < 1:
        start_line = 1
    if end_line - start_line + 1 > MAX_LINES:
        end_line = start_line + MAX_LINES - 1

    try:
        all_lines = fpath.read_text(encoding="utf-8", errors="replace").splitlines()
    except OSError as e:
        return f"ERROR reading file: {e}"

    total = len(all_lines)
    chunk = all_lines[start_line - 1 : end_line]
    header = f"# {fpath.name} (lines {start_line}–{min(end_line, total)} of {total})\n"
    body = "\n".join(f"{start_line + i:>5} │ {line}" for i, line in enumerate(chunk))
    footer = ""
    if end_line < total:
        footer = f"\n# … {total - end_line} more line(s). Next chunk: start_line={end_line + 1}"
    return header + body + footer


@mcp.tool()
def lsp_query(file_path: str, action: str, symbol: str = "", line: int = 0, col: int = 0) -> str:
    """
    Performs rudimentary LSP queries on Python files (not a full LSP server).
    Actions:
    - 'signature': Returns the signature of a symbol (function parameters, docstring).
    - 'find_references': Finds all usages of a symbol in the file.
    - 'completions': Shows completions at position (line, col).
    Python files only (.py). For other languages: use repo_map + read_file_chunked.
    Example: lsp_query("src/main.py", "signature", symbol="process_request")
    """
    try:
        fpath = _resolve_safe_path(file_path)
    except PermissionError as e:
        return f"ERROR: {e}"
    if not fpath.exists():
        return f"ERROR: File not found: '{file_path}'"
    if fpath.suffix.lower() != ".py":
        return (
            f"lsp_query supports Python files only (.py). "
            f"For '{fpath.suffix}' files: use repo_map + read_file_chunked."
        )

    try:
        import jedi  # type: ignore
    except ImportError:
        return "ERROR: jedi not installed. Run 'pip install jedi' in the MCP container."

    try:
        source = fpath.read_text(encoding="utf-8", errors="replace")
    except OSError as e:
        return f"ERROR reading file: {e}"

    action = action.strip().lower()

    if action == "signature":
        if not symbol:
            return "ERROR: 'symbol' must be provided for action='signature'."
        # Find first usage of the symbol as a call
        match = re.search(rf"\b{re.escape(symbol)}\s*\(", source)
        if not match:
            return f"Symbol '{symbol}' not found as a call in '{fpath.name}'. Check repo_map."
        char_pos = match.start() + len(match.group()) - 1
        row = source[:char_pos].count("\n") + 1
        col_pos = char_pos - source[:char_pos].rfind("\n") - 1
        try:
            script = jedi.Script(source=source, path=str(fpath))
            sigs = script.get_signatures(line=row, column=col_pos)
            if not sigs:
                return f"No signature found for '{symbol}' (Jedi). Check whether definition is in the same file."
            results = []
            for sig in sigs[:3]:
                params = ", ".join(p.description for p in sig.params)
                doc = sig.docstring(raw=True)[:300] if sig.docstring() else ""
                results.append(f"def {sig.name}({params})\n  Docstring: {doc or '—'}")
            return f"Signature for '{symbol}':\n" + "\n\n".join(results)
        except Exception as e:
            return f"Jedi error in signature: {e}"

    elif action == "find_references":
        if not symbol:
            return "ERROR: 'symbol' must be provided for action='find_references'."
        match = re.search(rf"\b{re.escape(symbol)}\b", source)
        if not match:
            return f"Symbol '{symbol}' not found in '{fpath.name}'."
        char_pos = match.start() + len(symbol) // 2
        row = source[:char_pos].count("\n") + 1
        col_pos = char_pos - source[:char_pos].rfind("\n") - 1
        try:
            script = jedi.Script(source=source, path=str(fpath))
            refs = script.get_references(line=row, column=col_pos)
            if not refs:
                return f"No references found for '{symbol}'."
            lines_out = [f"References for '{symbol}' in '{fpath.name}' ({len(refs)} matches):"]
            for ref in refs[:20]:
                lines_out.append(f"  Line {ref.line}: {ref.description}")
            if len(refs) > 20:
                lines_out.append(f"  … {len(refs) - 20} more (read_file_chunked for details)")
            return "\n".join(lines_out)
        except Exception as e:
            return f"Jedi error in find_references: {e}"

    elif action == "completions":
        if line <= 0 or col < 0:
            return "ERROR: 'line' (>0) and 'col' (>=0) must be provided for action='completions'."
        try:
            script = jedi.Script(source=source, path=str(fpath))
            comps = script.complete(line=line, column=col)
            if not comps:
                return f"No completions at line {line}, column {col}."
            names = [f"  {c.name} ({c.type})" for c in comps[:15]]
            return f"Completions at {fpath.name}:{line}:{col}:\n" + "\n".join(names)
        except Exception as e:
            return f"Jedi error in completions: {e}"

    else:
        return f"Unknown action '{action}'. Valid: 'signature', 'find_references', 'completions'."


# ─── File Generation Tool ───────────────────────────────────────────────────

_GENERATED_DIR = Path(os.getenv("MOE_GENERATED_DIR", "/app/generated"))
_GENERATED_DIR.mkdir(exist_ok=True)

# ─── MinIO helpers ──────────────────────────────────────────────────────────

def _minio_client():
    """Return a configured S3 client (Garage backend), or None if credentials are missing."""
    endpoint  = os.getenv("MINIO_ENDPOINT", "")
    access    = os.getenv("MINIO_ROOT_USER", "")
    secret    = os.getenv("MINIO_ROOT_PASSWORD", "")
    if not (endpoint and access and secret):
        return None
    try:
        from minio import Minio
        # Garage signs requests for its configured region ("garage" in the
        # production config). MinIO's SDK otherwise defaults to us-east-1,
        # causing valid uploads and pre-signed downloads to fail signature
        # validation against Garage.
        region = os.getenv("MINIO_REGION", os.getenv("GARAGE_REGION", "garage"))
        return Minio(
            endpoint,
            access_key=access,
            secret_key=secret,
            secure=False,
            region=region,
        )
    except Exception:
        return None


def _minio_public_url() -> str:
    """Return the admin-configured public base URL for MinIO, falling back to endpoint."""
    # MINIO_PUBLIC_URL is writable via Admin Portal → Settings → Storage URL
    url = os.getenv("MINIO_PUBLIC_URL", "").rstrip("/")
    if not url:
        endpoint = os.getenv("MINIO_ENDPOINT", "moe-storage-garage:3900")
        url = f"http://{endpoint}"
    return url


def _minio_ensure_bucket(client, bucket: str) -> None:
    if not client.bucket_exists(bucket):
        client.make_bucket(bucket)


def _minio_upload_bytes(data: bytes, object_name: str, content_type: str,
                        bucket: str | None = None) -> str:
    """Upload bytes to MinIO and return a pre-signed download URL (24h expiry)."""
    from datetime import timedelta
    mc = _minio_client()
    if mc is None:
        raise RuntimeError("MinIO not configured (MINIO_ENDPOINT/MINIO_ROOT_USER/MINIO_ROOT_PASSWORD missing)")
    bkt = bucket or os.getenv("MINIO_DEFAULT_BUCKET", "moe-files")
    _minio_ensure_bucket(mc, bkt)
    mc.put_object(bkt, object_name, io.BytesIO(data), length=len(data), content_type=content_type)
    presigned = mc.presigned_get_object(bkt, object_name, expires=timedelta(hours=24))
    # Replace internal hostname with public URL
    public_base = _minio_public_url()
    internal = os.getenv("MINIO_ENDPOINT", "moe-storage-garage:3900")
    presigned = presigned.replace(f"http://{internal}", public_base, 1)
    return presigned


def file_upload(content_base64: str, filename: str, content_type: str = "application/octet-stream",
                bucket: str = "") -> str:
    """Upload a file to MinIO object storage and return a 24h pre-signed download URL.

    Args:
        content_base64: Base64-encoded file content
        filename: Target filename (e.g. 'report.pdf')
        content_type: MIME type (default: application/octet-stream)
        bucket: Bucket name (default: moe-files from env)
    Returns:
        Pre-signed download URL valid for 24 hours, or error message.
    """
    import uuid as _uuid
    try:
        data = base64.b64decode(content_base64)
    except Exception as e:
        return f"Error: invalid base64 content — {e}"
    safe = re.sub(r'[^\w\-.]', '_', filename)[:120]
    object_name = f"{_uuid.uuid4().hex[:8]}_{safe}"
    try:
        url = _minio_upload_bytes(data, object_name, content_type, bucket or None)
        size_kb = len(data) / 1024
        return f"Uploaded: {safe} ({size_kb:.1f} KB)\nDownload URL (24h): {url}"
    except Exception as e:
        return f"Upload failed: {e}"


def file_download_url(object_name: str, bucket: str = "", expires_hours: int = 24) -> str:
    """Generate a fresh pre-signed download URL for an existing MinIO object.

    Args:
        object_name: Full object path in the bucket
        bucket: Bucket name (default: moe-files from env)
        expires_hours: Link validity in hours (1–168, default: 24)
    Returns:
        Pre-signed download URL, or error message.
    """
    from datetime import timedelta
    mc = _minio_client()
    if mc is None:
        return "Error: MinIO not configured."
    bkt = bucket or os.getenv("MINIO_DEFAULT_BUCKET", "moe-files")
    hours = max(1, min(168, expires_hours))
    try:
        presigned = mc.presigned_get_object(bkt, object_name, expires=timedelta(hours=hours))
        public_base = _minio_public_url()
        internal = os.getenv("MINIO_ENDPOINT", "moe-storage-garage:3900")
        presigned = presigned.replace(f"http://{internal}", public_base, 1)
        return f"Download URL ({hours}h): {presigned}"
    except Exception as e:
        return f"Error generating URL: {e}"


def generate_file(content: str, filename: str = "output", format: str = "html") -> str:
    """
    Generates a file from content and returns a download path.
    Supported formats: html, md, docx, txt, pptx, pdf.
    Use format='pdf' to create a proper PDF document (rendered via WeasyPrint).
    The file is stored server-side with a UUID prefix and can be downloaded
    via the /downloads/ endpoint. Files are auto-cleaned after 24 hours.
    """
    import uuid as _uuid
    _id = _uuid.uuid4().hex[:12]
    _safe_name = re.sub(r'[^\w\-.]', '_', filename)[:80]
    fmt = format.lower().strip()

    if fmt == "html":
        try:
            import markdown as _md
            html_body = _md.markdown(content, extensions=["tables", "fenced_code"])
        except ImportError:
            html_body = f"<pre>{content}</pre>"
        html = (
            "<!DOCTYPE html><html><head><meta charset='utf-8'>"
            f"<title>{_safe_name}</title>"
            "<style>body{font-family:system-ui;max-width:800px;margin:2rem auto;padding:0 1rem;line-height:1.6}"
            "table{border-collapse:collapse;width:100%}th,td{border:1px solid #ddd;padding:8px;text-align:left}"
            "pre{background:#f4f4f4;padding:1rem;overflow-x:auto;border-radius:4px}"
            "code{background:#f4f4f4;padding:2px 4px;border-radius:2px}</style></head>"
            f"<body>{html_body}</body></html>"
        )
        out_path = _GENERATED_DIR / f"{_id}_{_safe_name}.html"
        out_path.write_text(html, encoding="utf-8")

    elif fmt == "docx":
        try:
            from docx import Document
            doc = Document()
            for para in content.split("\n\n"):
                para = para.strip()
                if not para:
                    continue
                if para.startswith("# "):
                    doc.add_heading(para[2:], level=1)
                elif para.startswith("## "):
                    doc.add_heading(para[3:], level=2)
                elif para.startswith("### "):
                    doc.add_heading(para[4:], level=3)
                else:
                    doc.add_paragraph(para)
            out_path = _GENERATED_DIR / f"{_id}_{_safe_name}.docx"
            doc.save(str(out_path))
        except ImportError:
            return "Error: python-docx not available. Use format='html' or 'md' instead."

    elif fmt in ("md", "markdown", "txt", "text"):
        ext = "md" if fmt in ("md", "markdown") else "txt"
        out_path = _GENERATED_DIR / f"{_id}_{_safe_name}.{ext}"
        out_path.write_text(content, encoding="utf-8")

    elif fmt in ("pptx", "ppt", "powerpoint"):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()
            # Parse slides from content: "## Slide Title\n- bullet\n- bullet\n\n## Next Slide..."
            slide_blocks = re.split(r'\n(?=##?\s)', content.strip())
            for block in slide_blocks:
                block = block.strip()
                if not block:
                    continue
                lines = block.splitlines()
                title_line = lines[0].lstrip("#").strip()
                body_lines = [l for l in lines[1:] if l.strip()]
                layout = prs.slide_layouts[1]  # title + content
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = title_line
                if body_lines and slide.placeholders[1]:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    for i, bl in enumerate(body_lines):
                        bl = bl.lstrip("-*• ").strip()
                        if not bl:
                            continue
                        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                        p.text = bl
                        p.level = 0
            out_path = _GENERATED_DIR / f"{_id}_{_safe_name}.pptx"
            prs.save(str(out_path))
        except ImportError:
            return "Error: python-pptx not available."

    elif fmt in ("pdf",):
        try:
            import markdown as _md
            from weasyprint import HTML as _WH, CSS as _WCSS
        except ImportError as _ie:
            return f"Error: PDF generation requires weasyprint ({_ie}). Use format='html' as fallback."
        # Convert markdown → HTML → PDF
        try:
            html_body = _md.markdown(content, extensions=["tables", "fenced_code", "nl2br"])
        except Exception:
            html_body = f"<pre>{content}</pre>"
        full_html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>
  @page {{ margin: 2cm; }}
  body {{ font-family: "Liberation Sans", Arial, sans-serif; font-size: 11pt; line-height: 1.6; color: #1a1a1a; }}
  h1 {{ font-size: 22pt; color: #2c3e50; border-bottom: 2px solid #2c3e50; padding-bottom: 4pt; margin-top: 0; }}
  h2 {{ font-size: 16pt; color: #34495e; margin-top: 18pt; }}
  h3 {{ font-size: 13pt; color: #555; }}
  table {{ border-collapse: collapse; width: 100%; margin: 12pt 0; }}
  th, td {{ border: 1px solid #ccc; padding: 6pt 8pt; }}
  th {{ background: #f0f4f8; font-weight: bold; }}
  pre {{ background: #f7f7f7; border: 1px solid #ddd; padding: 8pt; border-radius: 3pt; font-size: 9pt; }}
  code {{ background: #f0f0f0; padding: 1pt 3pt; border-radius: 2pt; font-size: 9pt; }}
  a {{ color: #2980b9; }}
  ul, ol {{ margin-left: 20pt; }}
  li {{ margin-bottom: 3pt; }}
  blockquote {{ border-left: 3pt solid #ccc; margin: 0; padding-left: 12pt; color: #555; }}
</style>
</head><body>{html_body}</body></html>"""
        out_path = _GENERATED_DIR / f"{_id}_{_safe_name}.pdf"
        _WH(string=full_html).write_pdf(str(out_path))

    else:
        return f"Unsupported format: '{fmt}'. Use: html, docx, md, txt, pptx, pdf."

    size_kb = out_path.stat().st_size / 1024
    content_types = {
        ".html": "text/html",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".md": "text/markdown",
        ".txt": "text/plain",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    ct = content_types.get(out_path.suffix, "application/octet-stream")
    # Upload to MinIO if configured — always delete local copy afterwards to
    # prevent generated payloads from persisting on the host filesystem.
    try:
        url = _minio_upload_bytes(out_path.read_bytes(), out_path.name, ct)
        out_path.unlink(missing_ok=True)
        return f"File generated: {out_path.name} ({size_kb:.1f} KB)\nDownload URL (24h): {url}"
    except Exception:
        # MinIO unavailable — keep local copy for /downloads/ fallback but do NOT
        # execute it: the /downloads/ endpoint serves files read-only via FileResponse.
        return (
            f"File generated: {out_path.name} ({size_kb:.1f} KB)\n"
            f"Download: /downloads/{out_path.name}"
        )


# ─── Attachment Parser Tool ──────────────────────────────────────────────────

# Security constants for parse_attachment
_ATTACH_MAX_BYTES = 20 * 1024 * 1024  # 20 MB hard limit
_ATTACH_TIMEOUT   = 30.0              # seconds


def _assert_public_url(url: str) -> None:
    """Raise ValueError if the URL resolves to a private/loopback/link-local address.

    Prevents SSRF: an attacker could prompt the LLM to call fetch_pdf_text or
    parse_attachment with an internal URL (172.20.x.x, 169.254.x.x, localhost, etc.)
    to reach container-internal services like Postgres, Redis, or the admin API.
    """
    import urllib.parse
    parsed = urllib.parse.urlparse(url)
    scheme = parsed.scheme.lower()
    if scheme not in ("http", "https"):
        raise ValueError(f"URL scheme '{scheme}' is not allowed — only http/https.")
    host = parsed.hostname or ""
    if not host:
        raise ValueError("URL has no hostname.")
    # Reject obvious hostnames
    if host.lower() in ("localhost", "metadata.google.internal"):
        raise ValueError(f"Host '{host}' is not allowed (internal).")
    try:
        addr = ipaddress.ip_address(host)
        if addr.is_private or addr.is_loopback or addr.is_link_local or addr.is_reserved:
            raise ValueError(f"URL resolves to a private/reserved IP ({addr}) — SSRF blocked.")
    except ValueError as ve:
        if "SSRF blocked" in str(ve) or "not allowed" in str(ve) or "scheme" in str(ve):
            raise
        # hostname that isn't a bare IP — allow (DNS will resolve it at fetch time)
        pass


def _fetch_attachment_bytes(url: str) -> bytes:
    """
    Downloads a URL with a streaming size check (max 20 MB, 30 s timeout).
    Raises ValueError on size violation, httpx exceptions on network errors.
    """
    _assert_public_url(url)  # SSRF guard — blocks private/internal IPs
    collected: list[bytes] = []
    total = 0
    with httpx.stream("GET", url, timeout=_ATTACH_TIMEOUT, follow_redirects=False) as resp:
        resp.raise_for_status()
        for chunk in resp.iter_bytes(chunk_size=65536):
            total += len(chunk)
            if total > _ATTACH_MAX_BYTES:
                raise ValueError(
                    f"Attachment exceeds 20 MB size limit ({total} bytes downloaded so far)."
                )
            collected.append(chunk)
    return b"".join(collected)


def _parse_xlsx(data: bytes, max_chars: int) -> str:
    """Converts an XLSX workbook to a CSV-style plain text table."""
    try:
        import openpyxl
    except ImportError:
        return "Error: openpyxl not installed — cannot parse XLSX files."
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines: list[str] = [f"=== Sheet: {sheet_name} ==="]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            lines.append(",".join(cells))
        parts.append("\n".join(lines))
        if sum(len(p) for p in parts) >= max_chars:
            break
    return "\n\n".join(parts)[:max_chars]


def _parse_docx(data: bytes, max_chars: int) -> str:
    """Extracts plain text from a DOCX file (paragraphs only, no macros)."""
    try:
        from docx import Document as _DocxDocument
    except ImportError:
        return "Error: python-docx not installed — cannot parse DOCX files."
    doc = _DocxDocument(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)[:max_chars]


def _parse_pdf(data: bytes, max_chars: int) -> str:
    """Extracts plain text from a PDF using pypdf (no subprocess, no code exec)."""
    try:
        import pypdf
    except ImportError:
        return "Error: pypdf not installed — cannot parse PDF files."
    reader = pypdf.PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        parts.append(text)
        if sum(len(p) for p in parts) >= max_chars:
            break
    return "\n\n".join(parts)[:max_chars]


def _parse_csv(data: bytes, max_chars: int) -> str:
    """Decodes a CSV file as UTF-8 (with latin-1 fallback) and returns its text."""
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("latin-1", errors="replace")
    return text[:max_chars]


@mcp.tool()
def parse_attachment(url: str, max_chars: int = 6000) -> str:
    """
    Downloads and parses a file attachment from a URL and returns its text content.

    Supported formats:
    - XLSX → CSV-style table (all sheets)
    - DOCX → plain text (paragraphs)
    - PDF  → extracted text (all pages up to max_chars)
    - CSV / TXT → raw text

    Security constraints:
    - Maximum download size: 20 MB
    - Request timeout: 30 seconds
    - No code execution, no subprocess, no filesystem writes
    - URL must be http/https; path traversal is not applicable (URL-only)

    Args:
        url: HTTP/HTTPS URL of the file to download and parse.
        max_chars: Maximum number of characters to return (default 6000).

    Returns:
        Plain text content of the attachment, or an error string starting with 'Error:'.
    """
    # Validate URL scheme — only http/https allowed
    if not re.match(r"^https?://", url.strip(), re.IGNORECASE):
        return "Error: Only http:// and https:// URLs are supported."

    # Cap max_chars to a sane upper bound
    max_chars = max(100, min(max_chars, 50_000))

    try:
        raw = _fetch_attachment_bytes(url)
    except ValueError as e:
        return f"Error: {e}"
    except httpx.TimeoutException:
        return "Error: Download timed out (30 s limit)."
    except httpx.HTTPStatusError as e:
        return f"Error: HTTP {e.response.status_code} fetching attachment."
    except Exception as e:
        return f"Error downloading attachment: {e}"

    # Detect file type from URL extension (lowercase, strip query strings)
    url_path = url.split("?")[0].lower()
    if url_path.endswith(".xlsx"):
        return _parse_xlsx(raw, max_chars)
    if url_path.endswith(".docx"):
        return _parse_docx(raw, max_chars)
    if url_path.endswith(".pdf"):
        return _parse_pdf(raw, max_chars)
    if url_path.endswith((".csv", ".txt", ".tsv")):
        return _parse_csv(raw, max_chars)

    # Fallback: sniff magic bytes for known formats
    if raw[:4] == b"PK\x03\x04":
        # ZIP-based: could be XLSX or DOCX — try XLSX first, then DOCX
        try:
            return _parse_xlsx(raw, max_chars)
        except Exception:
            try:
                return _parse_docx(raw, max_chars)
            except Exception as e:
                return f"Error: ZIP-based file is neither XLSX nor DOCX: {e}"
    if raw[:4] == b"%PDF":
        return _parse_pdf(raw, max_chars)

    # Last resort: treat as plain text
    return _parse_csv(raw, max_chars)


# ─── Graph Analyzer Tool ─────────────────────────────────────────────────────


def _parse_graph_description(edges_description: str) -> tuple[list[tuple[str, str]], list[str]]:
    """
    Parses a text description of a graph into a list of (source, target) edge tuples
    and a deduplicated node list.

    Accepted formats (auto-detected, not mutually exclusive):
    - "nodes: A,B,C; edges: A-B, B-C"  — semicolon-separated header style
    - "A-B, B-C, A-C"                   — bare edge list, dash separator
    - "A→B; B→C"                        — arrow (directed) separator
    - CSV/table: lines like "A,B" or "A;B" (two columns = edge list)

    Direction is ignored for Eulerian analysis (treated as undirected graph).
    Returns (edges, nodes) where nodes are ordered by first appearance.
    """
    text = edges_description.strip()
    edges: list[tuple[str, str]] = []
    node_order: list[str] = []
    node_set: set[str] = set()

    def add_node(n: str) -> None:
        n = n.strip()
        if n and n not in node_set:
            node_set.add(n)
            node_order.append(n)

    # Normalize Unicode arrows to ASCII equivalents for easier parsing
    text = text.replace("→", "->").replace("←", "<-").replace("↔", "<->")

    # Attempt structured "nodes: ...; edges: ..." format
    nodes_match  = re.search(r"nodes?\s*:\s*([^;]+)", text, re.IGNORECASE)
    edges_match  = re.search(r"edges?\s*:\s*(.+)", text, re.IGNORECASE | re.DOTALL)

    if nodes_match:
        for n in re.split(r"[,\s]+", nodes_match.group(1)):
            n = n.strip()
            if n:
                add_node(n)

    edge_text = edges_match.group(1).strip() if edges_match else text

    # Split edge text by comma or semicolon, then parse each token as an edge
    separators = re.split(r"[,;\n]+", edge_text)
    for token in separators:
        token = token.strip()
        if not token:
            continue
        # Try arrow formats: A->B, A<-B, A<->B
        m = re.match(r"^(.+?)\s*(?:->|<->|--)\s*(.+)$", token)
        if not m:
            # Try reverse arrow A<-B (target → source)
            m2 = re.match(r"^(.+?)\s*<-\s*(.+)$", token)
            if m2:
                src, tgt = m2.group(2).strip(), m2.group(1).strip()
            else:
                # Try dash: A-B (single dash, not --)
                m3 = re.match(r"^([^-]+)-([^-].*)$", token)
                if m3:
                    src, tgt = m3.group(1).strip(), m3.group(2).strip()
                else:
                    # Try CSV two-column: "A,B" or "A;B"
                    parts = re.split(r"[,;]", token, maxsplit=1)
                    if len(parts) == 2 and parts[0].strip() and parts[1].strip():
                        src, tgt = parts[0].strip(), parts[1].strip()
                    else:
                        continue
        else:
            src, tgt = m.group(1).strip(), m.group(2).strip()

        if src and tgt:
            add_node(src)
            add_node(tgt)
            edges.append((src, tgt))

    return edges, node_order


def _analyze_eulerian(degree: dict[str, int], components: int) -> dict[str, Any]:
    """
    Determines Eulerian path/circuit existence for an undirected graph.

    Eulerian circuit exists iff: graph is connected AND all vertices have even degree.
    Eulerian path exists iff: graph is connected AND exactly two vertices have odd degree.

    Returns a dict with keys: has_circuit, has_path, odd_degree_nodes, explanation.
    """
    odd_nodes = [n for n, d in degree.items() if d % 2 != 0]
    if components > 1:
        return {
            "has_circuit": False,
            "has_path": False,
            "odd_degree_nodes": odd_nodes,
            "explanation": (
                f"Graph is disconnected ({components} components). "
                "Eulerian path/circuit requires a single connected component."
            ),
        }
    if len(odd_nodes) == 0:
        return {
            "has_circuit": True,
            "has_path": True,  # A circuit is also a path
            "odd_degree_nodes": [],
            "explanation": "All vertices have even degree → Eulerian circuit exists (starts and ends at same vertex).",
        }
    if len(odd_nodes) == 2:
        return {
            "has_circuit": False,
            "has_path": True,
            "odd_degree_nodes": odd_nodes,
            "explanation": (
                f"Exactly 2 odd-degree vertices ({odd_nodes[0]}, {odd_nodes[1]}) → "
                f"Eulerian path exists (from {odd_nodes[0]} to {odd_nodes[1]} or vice versa)."
            ),
        }
    return {
        "has_circuit": False,
        "has_path": False,
        "odd_degree_nodes": odd_nodes,
        "explanation": (
            f"{len(odd_nodes)} vertices have odd degree ({', '.join(odd_nodes[:10])}"
            f"{'…' if len(odd_nodes) > 10 else ''}) → No Eulerian path or circuit."
        ),
    }


def _connected_components(nodes: list[str], edges: list[tuple[str, str]]) -> list[list[str]]:
    """
    Computes connected components via iterative BFS (undirected interpretation of edges).
    Returns a list of components, each component being a sorted list of node names.
    """
    adjacency: dict[str, set[str]] = {n: set() for n in nodes}
    for src, tgt in edges:
        adjacency.setdefault(src, set()).add(tgt)
        adjacency.setdefault(tgt, set()).add(src)

    visited: set[str] = set()
    components: list[list[str]] = []

    for start in nodes:
        if start in visited:
            continue
        component: list[str] = []
        queue = [start]
        while queue:
            node = queue.pop()
            if node in visited:
                continue
            visited.add(node)
            component.append(node)
            queue.extend(adjacency.get(node, set()) - visited)
        components.append(sorted(component))

    return components


@mcp.tool()
def graph_analyze(edges_description: str) -> str:
    """
    Analyzes a graph described in plain text and returns structural findings as JSON.

    Input format (flexible — any of these work):
    - "nodes: A,B,C; edges: A-B, B-C, A-C"
    - "A->B, B->C, C->A"          (directed arrows, treated as undirected for Euler)
    - "A-B\\nB-C\\nC-D"           (one edge per line)
    - CSV table where each row is "source,target"

    Analysis results (JSON):
    - node_count, edge_count
    - degree_map: degree of each node
    - connected_components: list of node groups
    - is_connected: bool
    - eulerian: {has_circuit, has_path, odd_degree_nodes, explanation}
    - densitiy: edge_count / max_possible_edges (0..1)

    Returns:
        JSON string with analysis findings, or an error string starting with 'Error:'.

    Example:
        graph_analyze("nodes: A,B,C,D; edges: A-B, B-C, C-D, D-A")
        → {"node_count": 4, "edge_count": 4, "is_connected": true,
           "eulerian": {"has_circuit": true, ...}, ...}
    """
    if not edges_description.strip():
        return "Error: edges_description must not be empty."

    try:
        edges, nodes = _parse_graph_description(edges_description)
    except Exception as e:
        return f"Error parsing graph description: {e}"

    if not nodes:
        return "Error: No nodes found in the description. Check the input format."

    # Build degree map (undirected: each edge increments both endpoints)
    degree: dict[str, int] = {n: 0 for n in nodes}
    for src, tgt in edges:
        degree[src] = degree.get(src, 0) + 1
        if src != tgt:  # skip self-loops for degree count on other endpoint
            degree[tgt] = degree.get(tgt, 0) + 1

    # Self-loops count twice toward degree in undirected graphs
    self_loops = [(s, t) for s, t in edges if s == t]
    for s, _ in self_loops:
        degree[s] += 1  # already counted once above; add the second increment

    components = _connected_components(nodes, edges)
    n = len(nodes)
    e = len(edges)
    max_edges = n * (n - 1) // 2 if n > 1 else 0
    density = round(e / max_edges, 6) if max_edges > 0 else 0.0

    eulerian_info = _analyze_eulerian(degree, len(components))

    result = {
        "node_count": n,
        "edge_count": e,
        "nodes": nodes,
        "degree_map": degree,
        "connected_components": components,
        "component_count": len(components),
        "is_connected": len(components) == 1,
        "self_loops": [f"{s}-{t}" for s, t in self_loops],
        "density": density,
        "eulerian": eulerian_info,
    }

    try:
        return json.dumps(result, ensure_ascii=False, indent=2)
    except Exception as e:
        return f"Error serializing result: {e}"


# ─── Web / Document fetch helpers ───────────────────────────────────────────

def fetch_pdf_text(url: str, max_chars: int = 8000) -> str:
    """Download a PDF from a URL and extract its text content.

    Handles arXiv abstract pages (arxiv.org/abs/...) by automatically converting
    them to the direct PDF URL (arxiv.org/pdf/...). Falls back to fetching the
    HTML abstract if the PDF itself is inaccessible.

    url: URL to the PDF or arXiv abstract page
    max_chars: Maximum characters to return (default 8000)
    """
    import re as _re
    import pypdf

    _assert_public_url(url)  # SSRF guard

    # Normalise arXiv URLs: abs/ → pdf/ with .pdf suffix
    arxiv_abs = _re.match(r'https?://arxiv\.org/abs/(\d{4}\.\d+(?:v\d+)?)', url)
    if arxiv_abs:
        arxiv_id = arxiv_abs.group(1)
        pdf_url = f"https://arxiv.org/pdf/{arxiv_id}.pdf"
        html_url = url  # keep original for fallback
    else:
        pdf_url = url
        html_url = None

    headers = {
        "User-Agent": (
            "Mozilla/5.0 (compatible; research-bot/1.0; +https://example.com/bot)"
        )
    }

    def _try_pdf(target_url: str) -> str | None:
        try:
            resp = httpx.get(target_url, follow_redirects=True, timeout=45, headers=headers)
        except Exception as exc:
            return f"Download failed: {exc}"
        if resp.status_code != 200:
            return f"Download failed: HTTP {resp.status_code} for {target_url}"
        try:
            reader = pypdf.PdfReader(io.BytesIO(resp.content))
        except Exception as exc:
            return f"Not a valid PDF or could not parse: {exc}"
        pages_text: list[str] = []
        for page in reader.pages:
            try:
                pages_text.append(page.extract_text() or "")
            except Exception:
                continue
        full_text = "\n".join(pages_text).strip()
        if not full_text:
            return None  # Signal: try fallback
        return full_text[:max_chars]

    result = _try_pdf(pdf_url)
    if result and not result.startswith(("Download failed", "Not a valid")):
        return result

    # Fallback: fetch HTML abstract page (arXiv) or original URL as HTML
    fallback_url = html_url or url
    if fallback_url != pdf_url:
        try:
            resp = httpx.get(fallback_url, follow_redirects=True, timeout=20, headers=headers)
            if resp.status_code == 200:
                # Strip HTML tags for a plain-text abstract
                text = _re.sub(r'<[^>]+>', ' ', resp.text)
                text = _re.sub(r'\s+', ' ', text).strip()
                return f"[Abstract page fallback]\n{text[:max_chars]}"
        except Exception:
            pass

    return result or "No extractable text found (PDF may be image-only or behind a paywall)."


def python_sandbox(code: str) -> str:
    """Execute a small, self-contained Python snippet and return its output.

    Designed for numerical calculations, probability trees, Markov chains,
    combinatorics, and any logic that is too complex for the calculate tool.
    The snippet runs in a restricted environment: only the standard library
    modules math, fractions, itertools, collections, decimal, and statistics
    are available. No file I/O, no network, no subprocesses.

    code: Python source code. Use print() to produce output — the return value
          of the last expression is also captured automatically.

    Examples:
      - Recursive probability: from fractions import Fraction\\nP1=Fraction(1,3)\\n...
      - Combinatorics: import math; print(math.comb(10, 3))
      - Simulation: import random; random.seed(0); wins=[...]; print(sum(wins)/len(wins))
    """
    import builtins
    import io
    import contextlib

    _ALLOWED_MODULES = {
        "math", "fractions", "itertools", "collections",
        "decimal", "statistics", "random",
        # "re" intentionally excluded: regex objects expose __class__.__mro__
        # which can be used to escape the sandbox via __subclasses__() traversal.
    }

    def _restricted_import(name, *args, **kwargs):
        if name.split(".")[0] not in _ALLOWED_MODULES:
            raise ImportError(f"Module '{name}' is not allowed in sandbox (allowed: {sorted(_ALLOWED_MODULES)})")
        return original_import(name, *args, **kwargs)

    original_import = builtins.__import__
    stdout_capture = io.StringIO()

    # Block builtins that enable introspection-based sandbox escapes.
    # vars/dir/getattr allow traversal of __class__.__mro__.__subclasses__().
    _BLOCKED_BUILTINS = {
        "open", "exec", "eval", "compile", "input", "breakpoint",
        "__import__", "vars", "dir", "getattr", "setattr", "delattr",
        "globals", "locals", "memoryview",
    }
    safe_globals = {
        "__builtins__": {
            k: getattr(builtins, k)
            for k in dir(builtins)
            if not k.startswith("_") and k not in _BLOCKED_BUILTINS
        },
        "__import__": _restricted_import,
    }
    safe_globals["__builtins__"]["__import__"] = _restricted_import  # type: ignore[index]

    try:
        with contextlib.redirect_stdout(stdout_capture):
            exec(code, safe_globals)  # noqa: S102
        output = stdout_capture.getvalue().strip()
        return output if output else "(no output — use print() to show results)"
    except Exception as exc:
        partial = stdout_capture.getvalue().strip()
        partial_note = f"\nPartial output before crash:\n{partial}" if partial else ""
        return (
            f"[EXECUTION_ERROR] {type(exc).__name__}: {exc}{partial_note}\n"
            f"Fix the code and retry. Check for: index bounds, division by zero, "
            f"undefined variables, or disallowed modules."
        )


_ARC_COLOR_NAMES: dict[int, tuple[str, str]] = {
    0: ("black",   "·"),
    1: ("blue",    "B"),
    2: ("red",     "R"),
    3: ("green",   "G"),
    4: ("yellow",  "Y"),
    5: ("grey",    "Z"),
    6: ("pink",    "P"),
    7: ("orange",  "O"),
    8: ("azure",   "A"),
    9: ("maroon",  "M"),
}


def grid_repr(data: str, label: str = "") -> str:
    """Render a 2-D integer grid as annotated ASCII art with a colour legend.

    Designed for ARC-AGI grids (values 0–9) but works for any integer matrix.
    Values 0–9 use the standard ARC colour palette; values outside that range
    are shown as their decimal digit string.

    data:  JSON-encoded 2-D list, e.g. '[[0,1,2],[3,0,1]]'
    label: optional title printed above the grid (e.g. 'input' or 'output')

    Returns a compact ASCII block suitable for injection into a model prompt.
    """
    import json as _json

    try:
        grid = _json.loads(data)
    except _json.JSONDecodeError:
        return f"[GRID_ERROR] Invalid JSON: {data[:120]}"

    if not isinstance(grid, list) or not grid:
        return "[GRID_ERROR] data must be a non-empty list of rows"

    rows: list[list[int]] = []
    for r, row in enumerate(grid):
        if not isinstance(row, list):
            return f"[GRID_ERROR] Row {r} is not a list"
        rows.append([int(v) for v in row])

    n_rows = len(rows)
    n_cols = max((len(r) for r in rows), default=0)

    # Determine symbol width (values outside 0-9 need more chars)
    all_vals: set[int] = {v for row in rows for v in row}
    sym_width = max((len(str(v)) if v not in _ARC_COLOR_NAMES else 1) for v in all_vals) if all_vals else 1

    def sym(v: int) -> str:
        if v in _ARC_COLOR_NAMES:
            return _ARC_COLOR_NAMES[v][1]
        return str(v)

    col_labels = "  " + "  ".join(f"c{c:<{sym_width - 1}}" for c in range(n_cols))
    lines: list[str] = []
    if label:
        lines.append(f"── {label} ({n_rows}×{n_cols}) ──")
    else:
        lines.append(f"Grid ({n_rows}×{n_cols}):")
    lines.append(col_labels)
    for r, row in enumerate(rows):
        cells = "  ".join(f"{sym(v):>{sym_width}}" for v in row)
        lines.append(f"r{r:<2} {cells}")

    present = sorted(all_vals)
    legend_parts: list[str] = []
    for v in present:
        if v in _ARC_COLOR_NAMES:
            name, s = _ARC_COLOR_NAMES[v]
            legend_parts.append(f"{s}={name}({v})")
        else:
            legend_parts.append(str(v))
    lines.append("Legend: " + "  ".join(legend_parts))
    return "\n".join(lines)


def wikipedia_get_section(title: str = "", section: str = "", lang: str = "en", article: str = "") -> str:
    """Fetch a section of a Wikipedia article via the MediaWiki API.

    Returns the full plain-text content of the requested section (or the
    entire article intro if section is empty). Useful for counting items in
    discography, filmography, bibliography, or any other list sections.

    title:   Wikipedia article title (e.g. 'Mercedes Sosa'). Also accepts 'article' as alias.
    section: Section name to fetch (e.g. 'Studio albums'). Empty = article intro.
    lang:    Language code (default 'en'; use 'de' for German Wikipedia)
    """
    import re as _re

    # Accept 'article' as alias for 'title' — LLMs sometimes use this parameter name
    if not title and article:
        title = article

    base = f"https://{lang}.wikipedia.org/w/api.php"
    _headers = {"User-Agent": "MoE-Sovereign/1.0 (research-bot; contact@moe.local) python-httpx"}
    # First get the section index by fetching the table of contents
    try:
        toc_resp = httpx.get(
            base,
            params={
                "action": "parse", "page": title, "prop": "sections",
                "format": "json", "redirects": "1",
            },
            headers=_headers,
            timeout=15,
        )
    except Exception as exc:
        return f"Wikipedia API request failed: {exc}"

    if toc_resp.status_code != 200:
        return f"Wikipedia API error: HTTP {toc_resp.status_code}"

    toc_data = toc_resp.json()
    if "error" in toc_data:
        return f"Wikipedia API error: {toc_data['error'].get('info', toc_data['error'])}"

    sections = toc_data.get("parse", {}).get("sections", [])
    section_index = "0"  # 0 = lead section
    if section:
        low_target = section.lower()
        for s in sections:
            if low_target in s.get("line", "").lower() or low_target in s.get("anchor", "").lower():
                section_index = s["index"]
                break

    # Fetch the section content as plain text (wikitext → strip markup)
    try:
        content_resp = httpx.get(
            base,
            params={
                "action": "parse", "page": title, "prop": "wikitext",
                "section": section_index, "format": "json", "redirects": "1",
            },
            headers=_headers,
            timeout=15,
        )
    except Exception as exc:
        return f"Wikipedia section fetch failed: {exc}"

    if content_resp.status_code != 200:
        return f"Wikipedia section error: HTTP {content_resp.status_code}"

    wikitext = content_resp.json().get("parse", {}).get("wikitext", {}).get("*", "")

    # For discography/table sections: extract structured rows (Year | Entry) before stripping markup
    table_rows = []
    if _re.search(r'\{\|.*wikitable', wikitext, _re.DOTALL):
        current_year = None
        for line in wikitext.split('\n'):
            year_m = _re.match(r'^\|(\d{4})\s*$', line.strip())
            if year_m:
                current_year = year_m.group(1)
            elif line.startswith("|''") and current_year:
                title_clean = _re.sub(r'\[\[(?:[^|\]]*\|)?([^\]]+)\]\]', r'\1', line)
                title_clean = _re.sub(r"[|'{}]", '', title_clean).strip()
                title_clean = _re.sub(r'\{\{[^}]*\}\}', '', title_clean).strip()
                if title_clean:
                    table_rows.append(f"{current_year}: {title_clean}")

    # Strip wiki markup for a readable plain-text result
    text = _re.sub(r'\[\[(?:[^|\]]*\|)?([^\]]+)\]\]', r'\1', wikitext)  # links
    text = _re.sub(r'\{\{[^}]*\}\}', '', text)   # templates
    text = _re.sub(r"'''?", '', text)             # bold/italic
    text = _re.sub(r'<[^>]+>', ' ', text)         # HTML tags
    text = _re.sub(r'\n{3,}', '\n\n', text).strip()

    if not text:
        return f"Section '{section}' not found in article '{title}'. Available sections: {[s.get('line') for s in sections[:20]]}"

    # Prepend structured table summary if available
    prefix = ""
    if table_rows:
        prefix = f"STRUCTURED TABLE ({len(table_rows)} entries):\n" + "\n".join(table_rows) + "\n\n---\nRAW TEXT:\n"

    return f"Wikipedia — {title} [{section or 'intro'}]:\n\n{prefix}{text[:5000]}"


def github_get_issue(repo: str, issue_number: int) -> str:
    """Fetch a GitHub issue by repository and issue number via the public API.

    repo: Repository in 'owner/repo' format (e.g. 'torvalds/linux')
    issue_number: The issue number to fetch
    """
    api_url = f"https://api.github.com/repos/{repo}/issues/{issue_number}"
    try:
        response = httpx.get(
            api_url,
            headers={"Accept": "application/vnd.github.v3+json"},
            timeout=15,
        )
    except Exception as exc:
        return f"Request failed: {exc}"

    if response.status_code == 404:
        return f"Issue #{issue_number} not found in '{repo}' (or repository is private)."
    if response.status_code != 200:
        return f"GitHub API error: HTTP {response.status_code}"

    try:
        data = response.json()
    except Exception as exc:
        return f"Failed to parse GitHub response: {exc}"

    title = data.get("title", "(no title)")
    state = data.get("state", "unknown")
    created_at = data.get("created_at", "unknown")
    body = (data.get("body") or "")[:2000]
    labels = [lbl.get("name", "") for lbl in data.get("labels", [])]
    comments = data.get("comments", 0)

    return json.dumps(
        {
            "title": title,
            "state": state,
            "created_at": created_at,
            "body": body,
            "labels": labels,
            "comments_count": comments,
        },
        ensure_ascii=False,
        indent=2,
    )


@mcp.tool()
def github_issue_events(repo: str, issue_number: int, event_type: str = "labeled") -> str:
    """Fetch the timeline events of a GitHub issue to find when labels were added/removed.

    Use this when you need to know WHEN a specific label was added to an issue,
    not just which labels the issue currently has.

    repo:         Repository in 'owner/repo' format (e.g. 'numpy/numpy')
    issue_number: The issue number
    event_type:   Filter event type: 'labeled', 'unlabeled', 'closed', 'assigned', or '' for all
    """
    api_url = f"https://api.github.com/repos/{repo}/issues/{issue_number}/events"
    try:
        resp = httpx.get(
            api_url,
            headers={"Accept": "application/vnd.github.v3+json", "User-Agent": "MoE-Sovereign/1.0"},
            timeout=20,
        )
    except Exception as e:
        return f"[github_issue_events request failed: {e}]"

    if resp.status_code == 404:
        return f"[github_issue_events: Issue #{issue_number} not found in '{repo}']"
    if resp.status_code != 200:
        return f"[github_issue_events: HTTP {resp.status_code}]"

    try:
        events = resp.json()
    except Exception as e:
        return f"[github_issue_events: parse error: {e}]"

    if event_type:
        events = [e for e in events if e.get("event") == event_type]

    results = []
    for ev in events:
        entry = {
            "event": ev.get("event"),
            "created_at": ev.get("created_at"),
            "actor": (ev.get("actor") or {}).get("login"),
        }
        if ev.get("event") in ("labeled", "unlabeled"):
            entry["label"] = (ev.get("label") or {}).get("name")
        results.append(entry)

    if not results:
        return f"No '{event_type}' events found on issue #{issue_number}" if event_type else f"No events found"
    return json.dumps(results, indent=2)


# ─── External Data Sources ───────────────────────────────────────────────────

_SEARXNG_URL = os.environ.get("SEARXNG_URL", "").rstrip("/")


@mcp.tool()
def web_search_domain(query: str, domain: str = "", max_results: int = 5) -> str:
    """Domain-restricted web search via SearXNG.

    Adds a site: restriction to focus results on a specific website.
    Use when general search fails and you need data from a known source.

    query:       Search query (without site: — that is added automatically)
    domain:      Target domain, e.g. 'github.com', 'arxiv.org', 'wikipedia.org',
                 'pubchem.ncbi.nlm.nih.gov', 'orcid.org'. Leave empty for unrestricted search.
    max_results: Max search results to return (1-10, default 5)
    """
    if not _SEARXNG_URL:
        return "[web_search_domain: SEARXNG_URL not configured]"
    full_query = f"site:{domain} {query}" if domain else query
    try:
        resp = httpx.get(
            f"{_SEARXNG_URL}/search",
            params={"q": full_query, "format": "json", "engines": "google,bing,duckduckgo"},
            headers={"Accept": "application/json"},
            timeout=20,
        )
        resp.raise_for_status()
        data = resp.json()
        results = data.get("results", [])[:max_results]
        if not results:
            return f"[No results for: {full_query}]"
        parts = []
        for r in results:
            title = r.get("title", "")
            url   = r.get("url", "")
            snippet = r.get("content", "")[:300]
            parts.append(f"**{title}**\n{url}\n{snippet}")
        return "\n\n---\n\n".join(parts)
    except Exception as e:
        return f"[web_search_domain error: {e}]"


@mcp.tool()
def youtube_transcript(video_url: str, max_chars: int = 4000) -> str:
    """Fetch the transcript/captions of a YouTube video.

    Works with most public videos that have auto-generated or manual captions.
    Useful for questions about video content, interviews, documentaries, tutorials.

    video_url: Full YouTube URL (https://www.youtube.com/watch?v=...) or video ID
    max_chars: Maximum characters to return (default 4000)
    """
    # Extract video ID from URL or treat input as ID directly
    _id_match = re.search(
        r'(?:v=|youtu\.be/|embed/|shorts/)([A-Za-z0-9_\-]{11})',
        video_url,
    )
    video_id = _id_match.group(1) if _id_match else video_url.strip()
    if len(video_id) != 11 or not re.match(r'^[A-Za-z0-9_\-]+$', video_id):
        return f"[youtube_transcript: Could not extract valid video ID from: {video_url!r}]"

    # Primary: youtube-transcript-api (no API key required, version >= 0.6)
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        api = YouTubeTranscriptApi()
        # List available transcripts and pick English or German
        transcript_list = api.list(video_id)
        chosen = None
        for lang in ("en", "en-US", "en-GB", "de"):
            try:
                chosen = transcript_list.find_transcript([lang])
                break
            except Exception:
                pass
        if chosen is None:
            # Fall back to first available
            chosen = next(iter(transcript_list), None)
        if chosen is not None:
            entries = list(chosen.fetch())
            text = " ".join(getattr(e, "text", str(e)) for e in entries)
            return text[:max_chars] if len(text) > max_chars else text
    except ImportError:
        pass  # library not installed, fall through to HTTP fallback
    except Exception as e:
        err_str = str(e)
        if "Subtitles are disabled" in err_str or "Could not retrieve" in err_str or "TranscriptsDisabled" in err_str:
            return f"[youtube_transcript: Captions not available for video {video_id}]"
        # Other errors: fall through to HTTP fallback

    # Fallback: fetch YouTube page and extract caption track URL
    try:
        page = httpx.get(
            f"https://www.youtube.com/watch?v={video_id}",
            headers={"User-Agent": "Mozilla/5.0 (compatible; MoE-Research/1.0)"},
            timeout=15,
            follow_redirects=True,
        )
        # Extract timedtext URL from page source
        m = re.search(r'"captionTracks":\[.*?"baseUrl":"([^"]+)"', page.text)
        if not m:
            return f"[youtube_transcript: No caption track found for {video_id}]"
        caption_url = m.group(1).replace("\\u0026", "&")
        cap_resp = httpx.get(caption_url, timeout=15)
        # Parse XML caption format
        texts = re.findall(r'<text[^>]*>(.*?)</text>', cap_resp.text, re.DOTALL)
        import html
        clean = " ".join(html.unescape(t).replace("\n", " ") for t in texts)
        return clean[:max_chars] if len(clean) > max_chars else clean
    except Exception as e:
        return f"[youtube_transcript fallback error: {e}]"


@mcp.tool()
def github_search_issues(
    repo: str,
    labels: str = "",
    state: str = "open",
    sort: str = "created",
    order: str = "asc",
    max_results: int = 5,
    query: str = "",
) -> str:
    """Search GitHub issues in a repository using the GitHub Search API.

    Finds issues matching labels, state, and optional text query — unlike
    github_get_issue which requires a known issue number.

    IMPORTANT: GitHub labels often have prefixes (e.g. 'Regression' may be
    stored as '06 - Regression'). This tool auto-resolves partial label names
    by fetching the repo's label list first — pass the keyword only, e.g. 'Regression'.

    repo:        Repository in 'owner/repo' format (e.g. 'numpy/numpy')
    labels:      Comma-separated label keywords — partial matches resolved automatically
    state:       'open', 'closed', or 'all'
    sort:        'created', 'updated', 'comments'
    order:       'asc' (oldest first) or 'desc' (newest first)
    max_results: Max issues to return (1-10)
    query:       Additional text to search in issue title/body
    """
    # Resolve fuzzy label names → exact GitHub label names
    resolved_labels: list[str] = []
    if labels:
        try:
            label_resp = httpx.get(
                f"https://api.github.com/repos/{repo}/labels",
                params={"per_page": 100},
                headers={"Accept": "application/vnd.github.v3+json",
                         "User-Agent": "MoE-Sovereign/1.0"},
                timeout=10,
            )
            if label_resp.status_code == 200:
                all_labels = [l["name"] for l in label_resp.json()]
                for want in labels.split(","):
                    want = want.strip()
                    # Exact match first, then case-insensitive substring
                    exact = [l for l in all_labels if l == want]
                    # All fuzzy matches — e.g. 'Regression' → ['05 - Regression', '06 - Regression']
                    fuzzy = [l for l in all_labels if want.lower() in l.lower()]
                    resolved_labels.extend(exact or fuzzy or [want])
            else:
                resolved_labels = [l.strip() for l in labels.split(",") if l.strip()]
        except Exception:
            resolved_labels = [l.strip() for l in labels.split(",") if l.strip()]

    # Build one query per resolved label (GitHub AND-es multiple label: filters).
    # For OR semantics across label variants, run a query per label and merge results.
    def _build_q(label_filter: str) -> str:
        parts = [f"repo:{repo}"]
        if state and state != "all":
            parts.append(f"is:{state}")
        if label_filter:
            parts.append(f'label:"{label_filter}"')
        if query:
            parts.append(query)
        parts.append("is:issue")
        return " ".join(parts)

    label_queries = [_build_q(lbl) for lbl in resolved_labels] if resolved_labels else [_build_q("")]
    all_items: list[dict] = []
    seen_ids: set[int] = set()
    last_search_q = label_queries[0]

    for search_q in label_queries:
        last_search_q = search_q
        try:
            resp = httpx.get(
                "https://api.github.com/search/issues",
                params={"q": search_q, "sort": sort, "order": order, "per_page": max_results},
                headers={"Accept": "application/vnd.github.v3+json", "User-Agent": "MoE-Sovereign/1.0"},
                timeout=20,
            )
        except Exception as e:
            return f"[github_search_issues request failed: {e}]"
        if resp.status_code == 422:
            continue  # skip invalid query variant
        if resp.status_code != 200:
            return f"[github_search_issues: HTTP {resp.status_code}]"
        try:
            data = resp.json()
        except Exception:
            continue
        for item in data.get("items", []):
            if item["id"] not in seen_ids:
                seen_ids.add(item["id"])
                all_items.append(item)

    # Re-sort merged results by the requested sort field
    if sort == "created":
        all_items.sort(key=lambda x: x.get("created_at", ""), reverse=(order == "desc"))
    all_items = all_items[:max_results]

    if not all_items:
        return f"No issues found for query: {last_search_q}"

    results = []
    for issue in all_items:
        results.append({
            "number":     issue.get("number"),
            "title":      issue.get("title"),
            "state":      issue.get("state"),
            "created_at": issue.get("created_at"),
            "updated_at": issue.get("updated_at"),
            "labels":     [lbl["name"] for lbl in issue.get("labels", [])],
            "url":        issue.get("html_url"),
        })
    return json.dumps({"total_count": len(results), "query": last_search_q, "issues": results}, indent=2)


@mcp.tool()
def pubchem_compound_search(
    name: str = "",
    cid: int = 0,
    mw_min: float = 0,
    mw_max: float = 0,
    classification: str = "",
) -> str:
    """Search the PubChem compound database for chemical compound data.

    Can search by name, CID, or molecular weight range.
    Returns: CID, molecular weight, molecular formula, IUPAC name, and synonyms.

    name:           Compound name or IUPAC name to search (e.g. 'acetic acid')
    cid:            PubChem Compound ID (direct lookup, overrides name)
    mw_min/mw_max:  Molecular weight range filter (Da) — used with classification
    classification: FDA/GRAS classification filter (e.g. 'food additive')
    """
    base = "https://pubchem.ncbi.nlm.nih.gov/rest/pug"

    def _get_properties(cid_val: int) -> dict:
        prop_url = f"{base}/compound/cid/{cid_val}/property/MolecularWeight,MolecularFormula,IUPACName/JSON"
        r = httpx.get(prop_url, timeout=15)
        r.raise_for_status()
        props = r.json().get("PropertyTable", {}).get("Properties", [{}])[0]
        return props

    try:
        if cid:
            props = _get_properties(cid)
            props["CID"] = cid
            return json.dumps(props, indent=2)

        if name:
            # Name → CID lookup
            search_url = f"{base}/compound/name/{httpx.URL(name).path}/cids/JSON"
            r = httpx.get(search_url, timeout=15)
            if r.status_code == 404:
                return f"[pubchem: compound '{name}' not found]"
            r.raise_for_status()
            cids = r.json().get("IdentifierList", {}).get("CID", [])
            if not cids:
                return f"[pubchem: no results for '{name}']"
            results = []
            for c in cids[:5]:
                try:
                    props = _get_properties(c)
                    props["CID"] = c
                    results.append(props)
                except Exception:
                    pass
            return json.dumps(results, indent=2)

        if mw_min > 0 or mw_max > 0:
            # MW range search: scan candidates and filter client-side.
            # The fastformula endpoint's MolecularWeight query param is unreliable —
            # it may be silently ignored, returning unfiltered CIDs. Filter here instead.
            _props_url = f"{base}/compound/cid/{{chunk}}/property/MolecularWeight,MolecularFormula,IUPACName/JSON"
            if classification and "food" in classification.lower():
                # Food Additive CIDs are concentrated in low CID range (1-3000)
                # Most small-molecule food additives (MW ≤ 1000) have CID < 3000.
                candidate_cids = list(range(1, 3001))
            else:
                # Broad scan for unknown classification
                candidate_cids = list(range(1, 2001))
            results = []
            batch_size = 200
            for i in range(0, len(candidate_cids), batch_size):
                if len(results) >= 10:
                    break
                chunk = ",".join(str(c) for c in candidate_cids[i:i + batch_size])
                try:
                    r = httpx.get(_props_url.format(chunk=chunk), timeout=20)
                    if r.status_code != 200:
                        continue
                    rows = r.json().get("PropertyTable", {}).get("Properties", [])
                    for row in rows:
                        mw = float(row.get("MolecularWeight", 0) or 0)
                        if mw_min > 0 and mw < mw_min:
                            continue
                        if mw_max > 0 and mw > mw_max:
                            continue
                        results.append(row)
                        if len(results) >= 10:
                            break
                except Exception:
                    continue
            if not results:
                return f"[pubchem MW search: no compounds found in range {mw_min:.1f}-{mw_max:.1f} Da]"
            return json.dumps(results, indent=2)

        return "[pubchem_compound_search: provide name, cid, or mw_min/mw_max]"
    except Exception as e:
        return f"[pubchem_compound_search error: {e}]"


@mcp.tool()
def orcid_works_count(orcid_id: str, before_year: int = 0, after_year: int = 0) -> str:
    """Count and list publications on an ORCID researcher profile.

    Uses the public ORCID API (no authentication required for public profiles).
    Returns total work count and year distribution.

    orcid_id:    ORCID iD in format XXXX-XXXX-XXXX-XXXX
    before_year: Count only works published before this year (e.g. 2020 → pre-2020)
    after_year:  Count only works published after this year
    """
    # Normalise ORCID format
    orcid_clean = orcid_id.strip().replace("https://orcid.org/", "")
    if not re.match(r'^\d{4}-\d{4}-\d{4}-\d{3}[\dX]$', orcid_clean):
        return f"[orcid_works_count: invalid ORCID format: {orcid_id!r}. Expected XXXX-XXXX-XXXX-XXXX]"

    url = f"https://pub.orcid.org/v3.0/{orcid_clean}/works"
    try:
        resp = httpx.get(
            url,
            headers={"Accept": "application/json"},
            timeout=20,
        )
    except Exception as e:
        return f"[orcid_works_count request failed: {e}]"

    if resp.status_code == 404:
        return f"[orcid_works_count: ORCID {orcid_clean} not found or profile is private]"
    if resp.status_code != 200:
        return f"[orcid_works_count: HTTP {resp.status_code}]"

    try:
        data = resp.json()
    except Exception as e:
        return f"[orcid_works_count: parse error: {e}]"

    groups = data.get("group", [])
    years: list[int] = []
    for grp in groups:
        # Each group represents one unique work. Multiple work-summaries within a group
        # are different data sources for the same work — take only the first year found
        # to avoid counting the same publication multiple times.
        year_for_group: int | None = None
        for ws in grp.get("work-summary", []):
            pub_date = ws.get("publication-date") or {}
            year_val = (pub_date.get("year") or {}).get("value")
            if year_val:
                try:
                    year_for_group = int(year_val)
                    break
                except ValueError:
                    pass
        if year_for_group is not None:
            years.append(year_for_group)

    # Apply year filters
    filtered = years
    if before_year > 0:
        filtered = [y for y in filtered if y < before_year]
    if after_year > 0:
        filtered = [y for y in filtered if y > after_year]

    from collections import Counter
    year_dist = dict(sorted(Counter(filtered).items()))
    result = {
        "orcid": orcid_clean,
        "total_unique_works": len(groups),
        "works_with_year": len(years),
        "filtered_count": len(filtered),
        "filter_applied": {
            "before_year": before_year or None,
            "after_year": after_year or None,
        },
        "NOTE": "Use 'filtered_count' when a year filter was applied, 'total_unique_works' otherwise.",
        "year_distribution": year_dist,
    }
    return json.dumps(result, indent=2)


@mcp.tool()
def pubchem_advanced_search(
    mw_max: float = 0,
    mw_min: float = 0,
    heavy_atoms: int = 0,
    hb_acceptors_max: int = -1,
    hb_donors_max: int = -1,
    complexity_min: int = 0,
    complexity_max: int = 0,
    classification: str = "",
    max_results: int = 10,
) -> str:
    """Advanced PubChem compound search with multi-criteria filtering.

    Use this for GAIA-style questions like "find the compound in PubChem's Food Additive
    Status classification with MW ≤ 100, 6 heavy atoms, ≤ 1 hydrogen bond acceptors,
    and complexity between 10 and 15".

    All filter arguments are optional — provide only the ones you need.
    Returns matching compound CIDs with their molecular weight, formula, and IUPAC name.

    mw_min/mw_max:        Molecular weight range in Da (0 = no limit)
    heavy_atoms:          Exact heavy atom count (0 = no filter)
    hb_acceptors_max:     Max hydrogen bond acceptors (-1 = no filter)
    hb_donors_max:        Max hydrogen bond donors (-1 = no filter)
    complexity_min/max:   Bertz complexity score range (0 = no filter)
    classification:       Filter hint for result description (e.g. 'Food Additive')
    max_results:          Max compounds to return (1-20)
    """
    base = "https://pubchem.ncbi.nlm.nih.gov/rest/pug"
    _PROPS = "MolecularWeight,MolecularFormula,IUPACName,HBondAcceptorCount,HBondDonorCount,HeavyAtomCount,Complexity"

    def _batch_props(cids: list[int]) -> list[dict]:
        """Fetch properties for up to 200 CIDs in a single request."""
        if not cids:
            return []
        chunk = ",".join(str(c) for c in cids[:200])
        url = f"{base}/compound/cid/{chunk}/property/{_PROPS}/JSON"
        r = httpx.get(url, timeout=20)
        if r.status_code != 200:
            return []
        rows = r.json().get("PropertyTable", {}).get("Properties", [])
        for row in rows:
            row["CID"] = row.get("CID", 0)
        return rows

    try:
        # Build candidate CID list via NCATS Food Additive classification (if requested)
        # then apply multi-criteria property filters client-side.
        candidate_cids: list[int] = []

        if classification and "food" in classification.lower():
            # The /classification/cid/JSON?source=NCATS+Food+Additive+Status endpoint
            # returns HTTP 400. Scan CIDs 1-10000, which covers all known small-molecule
            # food additives in PubChem (the vast majority have CID < 10000).
            candidate_cids = list(range(1, 10001))

        # If no classification-based candidates, use MW-range PubChem search
        if not candidate_cids:
            if mw_min > 0 or mw_max > 0:
                lo = mw_min if mw_min > 0 else 0
                hi = mw_max if mw_max > 0 else 100000
                # PubChem FTP-style bulk property query via PUG REST
                search_url = (
                    f"{base}/compound/property/MolecularWeight,HeavyAtomCount,"
                    f"HBondAcceptorCount,HBondDonorCount,MolecularFormula,IUPACName/JSON"
                    f"?cid=1-{max_results * 500}"  # scan first N CIDs
                )
                # Better: use FastSearch with MW filter
                fast_url = f"https://pubchem.ncbi.nlm.nih.gov/sdq/sdqagent.cgi?infmt=json&outfmt=json&query={{\"download\":\"*\",\"collection\":\"compound\",\"where\":{{\"ands\":[{{\"mw_exact\":\"{lo:.2f},{hi:.2f}\"}},{{\"xlogp\":\"*\"}}]}},\"limit\":\"{max_results*10}\",\"downloadfilename\":\"PubChem_compound\"}}"
                # Simpler fallback: directly enumerate known small-molecule CIDs
                # For MW ≤ 100: PubChem CIDs 1-10000 cover most simple molecules
                candidate_cids = list(range(1, min(1001, max_results * 100)))
            else:
                return "[pubchem_advanced_search: provide at least mw_max, classification, or heavy_atoms filter]"

        # Fetch properties in batches of 200 CIDs and filter client-side
        results = []
        batch_size = 200
        for i in range(0, len(candidate_cids), batch_size):
            if len(results) >= max_results:
                break
            batch = candidate_cids[i:i + batch_size]
            try:
                props_list = _batch_props(batch)
            except Exception:
                continue
            for props in props_list:
                if len(results) >= max_results:
                    break
                mw  = float(props.get("MolecularWeight", 0) or 0)
                ha  = int(props.get("HeavyAtomCount", 0) or 0)
                hba = int(props.get("HBondAcceptorCount", 0) or 0)
                hbd = int(props.get("HBondDonorCount", 0) or 0)
                if mw_max > 0 and mw > mw_max:
                    continue
                if mw_min > 0 and mw < mw_min:
                    continue
                if heavy_atoms > 0 and ha != heavy_atoms:
                    continue
                if hb_acceptors_max >= 0 and hba > hb_acceptors_max:
                    continue
                if hb_donors_max >= 0 and hbd > hb_donors_max:
                    continue
                cplx = int(props.get("Complexity", 0) or 0)
                if complexity_min > 0 and cplx < complexity_min:
                    continue
                if complexity_max > 0 and cplx > complexity_max:
                    continue
                results.append(props)

        if not results:
            return (f"[pubchem_advanced_search: no compounds matched all criteria "
                    f"(MW {mw_min:.1f}-{mw_max:.1f}, HA={heavy_atoms}, HBA≤{hb_acceptors_max})]")

        return json.dumps({
            "classification_hint": classification,
            "filters": {
                "mw_range": f"{mw_min:.1f}-{mw_max:.1f}",
                "heavy_atoms": heavy_atoms or "any",
                "hb_acceptors_max": hb_acceptors_max if hb_acceptors_max >= 0 else "any",
                "hb_donors_max": hb_donors_max if hb_donors_max >= 0 else "any",
            },
            "results_count": len(results),
            "compounds": results,
        }, indent=2)

    except Exception as e:
        return f"[pubchem_advanced_search error: {e}]"


@mcp.tool()
def semantic_scholar_search(
    query: str,
    year_filter: str = "",
    max_results: int = 5,
    fetch_pdf: bool = False,
) -> str:
    """Search Semantic Scholar for academic papers and retrieve abstracts and PDF links.

    Semantic Scholar aggregates open-access papers from many publishers.  Use this
    to find specific academic papers by author, title, or topic — especially when
    those papers may be behind a paywall on the journal's own site.  If fetch_pdf
    is True and a free PDF URL exists, the first result's text is also extracted.

    query:       Author name + topic, or paper title fragment (e.g. "Valencfia-Mendez
                 harlequin shrimp length 2017")
    year_filter: Restrict results to a single year or range, e.g. "2017" or "2000-2005"
    max_results: Number of papers to return (1-10)
    fetch_pdf:   If True and a free PDF URL is found, extract up to 6000 chars of text
    """
    ss_url = "https://api.semanticscholar.org/graph/v1/paper/search"
    fields = "title,year,authors,abstract,openAccessPdf,externalIds,publicationVenue"
    params: dict = {
        "query": query,
        "fields": fields,
        "limit": min(max(1, max_results), 10),
    }
    if year_filter:
        params["year"] = year_filter

    # Retry up to 3 times with backoff on rate-limit (429) or transient errors.
    resp = None
    for _attempt in range(3):
        try:
            resp = httpx.get(ss_url, params=params, timeout=15,
                             headers={"User-Agent": "MoE-Research/1.0"})
            if resp.status_code != 429:
                break
            import time as _time
            _time.sleep(2 ** _attempt)  # 1s, 2s, 4s backoff
        except Exception as e:
            if _attempt == 2:
                return f"[semantic_scholar_search request failed: {e}]"
            import time as _time
            _time.sleep(1)

    if resp is None:
        return "[semantic_scholar_search: all retries failed]"
    if resp.status_code == 429:
        return "[semantic_scholar_search: rate limited after 3 retries — use web_search_domain with site:semanticscholar.org as fallback]"
    if resp.status_code != 200:
        return f"[semantic_scholar_search: HTTP {resp.status_code}]"

    try:
        data = resp.json()
    except Exception as e:
        return f"[semantic_scholar_search: parse error: {e}]"

    papers = data.get("data", [])
    if not papers:
        return f"[semantic_scholar_search: no results for query '{query[:80]}']"

    results = []
    pdf_text = ""
    for p in papers:
        authors = ", ".join(a.get("name", "") for a in (p.get("authors") or [])[:4])
        oa = p.get("openAccessPdf") or {}
        pdf_url = oa.get("url", "")
        doi = (p.get("externalIds") or {}).get("DOI", "")
        venue = (p.get("publicationVenue") or {}).get("name", "")
        entry = {
            "title":    p.get("title", ""),
            "year":     p.get("year"),
            "authors":  authors,
            "venue":    venue,
            "doi":      doi,
            "pdf_url":  pdf_url,
            "abstract": (p.get("abstract") or "")[:800],
        }
        results.append(entry)

        # Optionally fetch the first available PDF text
        if fetch_pdf and pdf_url and not pdf_text:
            try:
                pr = httpx.get(pdf_url, timeout=20, follow_redirects=True)
                if pr.status_code == 200 and b"%PDF" in pr.content[:8]:
                    import io, pdfminer.high_level as _pdf  # type: ignore
                    pdf_text = _pdf.extract_text(io.BytesIO(pr.content))[:6000]
            except Exception:
                pass

    output: dict = {
        "query": query,
        "total_found": data.get("total", len(results)),
        "papers": results,
    }
    if pdf_text:
        output["pdf_text_preview"] = pdf_text

    return json.dumps(output, indent=2, ensure_ascii=False)


@mcp.tool()
def wikidata_search(text: str, language: str = "en", max_results: int = 5) -> str:
    """Search Wikidata for entities by text and return their IDs for use in wikidata_sparql.

    Use this BEFORE wikidata_sparql when you don't know the entity ID (wd:Q...).
    Returns list of {id, label, description} — use the id in a follow-up wikidata_sparql call.

    Example workflow:
      1. wikidata_search("Morarji Desai") → finds wd:Q192131 (actually Giuseppe Meazza)
         → use the correct ID from results
      2. wikidata_sparql("SELECT ?label WHERE { wd:<id> rdfs:label ?label ... }")

    text:        Search term (person name, place, concept)
    language:    Result label language (default: "en")
    max_results: Max entities to return (1-10)
    """
    url = "https://www.wikidata.org/w/api.php"
    try:
        resp = httpx.get(url, params={
            "action":   "wbsearchentities",
            "search":   text,
            "language": language,
            "format":   "json",
            "limit":    min(max(1, max_results), 10),
            "type":     "item",
        }, headers={"User-Agent": "MoE-Research/1.0"}, timeout=10)
    except Exception as e:
        return f"[wikidata_search request failed: {e}]"
    if resp.status_code != 200:
        return f"[wikidata_search: HTTP {resp.status_code}]"
    try:
        results = resp.json().get("search", [])
    except Exception as e:
        return f"[wikidata_search: parse error: {e}]"
    if not results:
        return f"[wikidata_search: no entities found for '{text}']"
    entities = [
        {"id": r.get("id", ""), "label": r.get("label", ""), "description": r.get("description", "")}
        for r in results
    ]
    return json.dumps({"query": text, "entities": entities}, indent=2, ensure_ascii=False)


@mcp.tool()
def wikidata_sparql(sparql_query: str, max_results: int = 10) -> str:
    """Execute a SPARQL query against Wikidata for deterministic fact lookup.

    Use for entity facts, dates, locations, relationships — deterministic, no
    SearXNG variance, no HTML parsing.  Wikidata has 100M+ entities.

    Examples:
      SELECT ?label WHERE { wd:Q192131 rdfs:label ?label. FILTER(LANG(?label)='en') }
      SELECT ?dob WHERE { wd:Q... wdt:P569 ?dob }

    sparql_query: Valid SPARQL 1.1 for https://query.wikidata.org/sparql
    max_results:  LIMIT injected when query has none (1-50)
    """
    endpoint = "https://query.wikidata.org/sparql"
    q = sparql_query.strip()
    if "LIMIT" not in q.upper():
        q = q.rstrip(";") + f" LIMIT {min(max(1, max_results), 50)}"
    try:
        resp = httpx.get(
            endpoint,
            params={"query": q, "format": "json"},
            headers={"Accept": "application/sparql-results+json",
                     "User-Agent": "MoE-Research/1.0"},
            timeout=20,
        )
    except Exception as e:
        return f"[wikidata_sparql request failed: {e}]"
    if resp.status_code == 400:
        return f"[wikidata_sparql: invalid SPARQL — {resp.text[:200]}]"
    if resp.status_code != 200:
        return f"[wikidata_sparql: HTTP {resp.status_code}]"
    try:
        data = resp.json()
    except Exception as e:
        return f"[wikidata_sparql: parse error: {e}]"
    bindings = data.get("results", {}).get("bindings", [])
    if not bindings:
        return "[wikidata_sparql: no results — try a broader query or check entity IDs]"
    rows = [
        {k: v.get("value", "") for k, v in b.items()}
        for b in bindings[:max_results]
    ]
    return json.dumps({"query": sparql_query[:200], "results": rows},
                      indent=2, ensure_ascii=False)


@mcp.tool()
def pubmed_search(query: str, max_results: int = 5, year_min: int = 0) -> str:
    """Search PubMed/NCBI for biomedical and life science academic papers.

    Use for biology, ecology, medicine, species studies, genetics, clinical papers.
    Prefer over semantic_scholar_search when the paper is biology/ecology/medicine.
    Supports PubMed Boolean syntax: "Hafnia alvei[tiab] AND mice", "harlequin shrimp[tiab]".

    query:       PubMed search query
    max_results: Number of papers to return (1-10)
    year_min:    Restrict to papers from this year onward (0 = no filter)
    """
    base = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils"
    full_query = f"({query}) AND {year_min}:3000[pdat]" if year_min > 0 else query
    try:
        sr = httpx.get(f"{base}/esearch.fcgi", params={
            "db": "pubmed", "term": full_query,
            "retmax": min(max(1, max_results), 10),
            "retmode": "json",
            "tool": "MoE-Research", "email": "research@moe-sovereign.org",
        }, timeout=15)
        if sr.status_code != 200:
            return f"[pubmed_search: search HTTP {sr.status_code}]"
        pmids = sr.json().get("esearchresult", {}).get("idlist", [])
        if not pmids:
            return f"[pubmed_search: no results for '{query[:80]}']"

        fr = httpx.get(f"{base}/efetch.fcgi", params={
            "db": "pubmed", "id": ",".join(pmids),
            "rettype": "abstract", "retmode": "xml",
            "tool": "MoE-Research", "email": "research@moe-sovereign.org",
        }, timeout=20)
        if fr.status_code != 200:
            return f"[pubmed_search: fetch HTTP {fr.status_code}]"

        import xml.etree.ElementTree as _ET
        root = _ET.fromstring(fr.content)
        papers = []
        for article in root.findall(".//PubmedArticle"):
            def _t(path):
                el = article.find(path)
                return (el.text or "").strip() if el is not None else ""
            abstract = " ".join(
                (el.text or "") for el in article.findall(".//AbstractText")
            )[:600]
            authors = ", ".join(
                f"{el.findtext('LastName', '')} {el.findtext('Initials', '')}".strip()
                for el in article.findall(".//Author")[:4]
            )
            papers.append({
                "pmid":     _t(".//PMID"),
                "title":    _t(".//ArticleTitle"),
                "year":     _t(".//PubDate/Year") or _t(".//PubDate/MedlineDate")[:4],
                "authors":  authors,
                "doi":      _t(".//ArticleId[@IdType='doi']"),
                "abstract": abstract,
            })
        return json.dumps({"query": query, "papers": papers}, indent=2, ensure_ascii=False)
    except Exception as e:
        return f"[pubmed_search error: {e}]"


@mcp.tool()
def wayback_fetch(url: str, timestamp: str = "", max_chars: int = 6000) -> str:
    """Fetch a URL from the Wayback Machine (web.archive.org).

    Retrieves a historical snapshot of any public web page. Useful for:
    - ORCID profiles at a past point in time (use timestamp=YYYYMMDD)
    - Pages that have since changed or disappeared
    - Historical data that current APIs no longer return

    url:       The original URL to look up (e.g. https://orcid.org/0000-0001-2345-6789)
    timestamp: YYYYMMDD or YYYYMMDDHHMMSS — omit for the closest available snapshot
    max_chars: Maximum characters to return from the archived page
    """
    _assert_public_url(url)
    import re as _re

    # Strategy 1: availability API (fast, but occasionally returns empty)
    snap_ts = timestamp or "20240101000000"
    archive_url: str | None = None
    avail_params: dict = {"url": url}
    if timestamp:
        avail_params["timestamp"] = timestamp
    try:
        avail_resp = httpx.get(
            "http://archive.org/wayback/available",
            params=avail_params,
            timeout=12,
        )
        if avail_resp.status_code == 200:
            snap = avail_resp.json().get("archived_snapshots", {}).get("closest", {})
            if snap.get("available"):
                archive_url = snap["url"]
                snap_ts = snap.get("timestamp", snap_ts)
    except Exception:
        pass  # fall through to direct URL strategy

    # Strategy 2: direct URL — always works if archive.org has any snapshot
    if not archive_url:
        ts = timestamp or "2024"
        archive_url = f"https://web.archive.org/web/{ts}/{url}"

    try:
        page_resp = httpx.get(
            archive_url,
            follow_redirects=True,
            timeout=28,
            headers={"User-Agent": "Mozilla/5.0 (research-bot/1.0)"},
        )
        if page_resp.status_code == 404:
            return f"[wayback_fetch: no archived version of '{url}' found]"
        if page_resp.status_code != 200:
            return f"[wayback_fetch: HTTP {page_resp.status_code} from archive.org]"
        # Extract actual snapshot timestamp from redirect URL if available
        if "web.archive.org/web/" in str(page_resp.url):
            snap_ts = str(page_resp.url).split("web.archive.org/web/")[1].split("/")[0]
        text = _re.sub(r'<[^>]+>', ' ', page_resp.text)
        text = _re.sub(r'\s+', ' ', text).strip()
        result = text[:max_chars] if len(text) > max_chars else text
        return json.dumps({
            "source_url": url,
            "archive_url": str(page_resp.url),
            "snapshot_timestamp": snap_ts,
            "content": result,
        }, ensure_ascii=False)
    except Exception as e:
        return f"[wayback_fetch: failed to retrieve archived '{url}': {e}]"


@mcp.tool()
def crossref_lookup(query: str, max_results: int = 5, filter_type: str = "") -> str:
    """Search CrossRef for academic publications by title, author, DOI, or keyword.

    CrossRef indexes 150M+ scholarly works with reliable metadata. Deterministic,
    no API key required. Use for:
    - Finding a paper's exact title, authors, DOI, publication year, journal
    - Counting articles by publisher/journal in a given year
    - Verifying publication details before answering

    query:       Title, author, keyword, or DOI string
    max_results: Number of works to return (1-20)
    filter_type: Optional work type filter: journal-article, book-chapter, proceedings-article
    """
    params: dict = {
        "query": query,
        "rows": min(max(1, max_results), 20),
        "mailto": "research@moe-sovereign.org",
        "select": "DOI,title,author,published,publisher,container-title,type,is-referenced-by-count",
    }
    if filter_type:
        params["filter"] = f"type:{filter_type}"
    try:
        resp = httpx.get("https://api.crossref.org/works", params=params, timeout=15)
        if resp.status_code != 200:
            return f"[crossref_lookup: HTTP {resp.status_code}]"
        data = resp.json()
        items = data.get("message", {}).get("items", [])
        if not items:
            return f"[crossref_lookup: no results for '{query[:80]}']"
        results = []
        for item in items[:max_results]:
            authors = ", ".join(
                f"{a.get('family', '')} {a.get('given', '')}".strip()
                for a in item.get("author", [])[:4]
            )
            pub_date = item.get("published", {}).get("date-parts", [[None]])[0]
            year = pub_date[0] if pub_date else None
            results.append({
                "doi":     item.get("DOI", ""),
                "title":   (item.get("title") or [""])[0],
                "authors": authors,
                "year":    year,
                "journal": (item.get("container-title") or [""])[0],
                "type":    item.get("type", ""),
                "cited_by": item.get("is-referenced-by-count", 0),
            })
        return json.dumps({"query": query, "results": results}, indent=2, ensure_ascii=False)
    except Exception as e:
        return f"[crossref_lookup error: {e}]"


@mcp.tool()
def openalex_search(query: str, max_results: int = 5,
                    year_min: int = 0, year_max: int = 0,
                    open_access_only: bool = False) -> str:
    """Search OpenAlex — the world's largest open academic database (250M+ works).

    OpenAlex covers all research fields with rich metadata: concepts, citations,
    funding, author affiliations. Free, no API key, deterministic.
    Complements PubMed (biomedical) and SemanticScholar with broader coverage.

    Use for: counting publications by venue/year, finding papers across disciplines,
    author publication counts, citation networks.

    query:            Search string (title, abstract, author, concept)
    max_results:      Number of works to return (1-20)
    year_min:         Only return works published from this year onward
    year_max:         Only return works published up to and including this year
    open_access_only: Restrict to open-access works
    """
    params: dict = {
        "search":   query,
        "per_page": min(max(1, max_results), 20),
        "mailto":   "research@moe-sovereign.org",
        "select":   "id,title,authorships,publication_year,doi,primary_location,cited_by_count,type",
    }
    filters = []
    if year_min > 0:
        filters.append(f"publication_year:>{year_min - 1}")
    if year_max > 0:
        filters.append(f"publication_year:<{year_max + 1}")
    if open_access_only:
        filters.append("is_oa:true")
    if filters:
        params["filter"] = ",".join(filters)
    try:
        resp = httpx.get("https://api.openalex.org/works", params=params, timeout=15)
        if resp.status_code != 200:
            return f"[openalex_search: HTTP {resp.status_code}]"
        data = resp.json()
        items = data.get("results", [])
        if not items:
            return f"[openalex_search: no results for '{query[:80]}']"
        results = []
        for item in items[:max_results]:
            authors = ", ".join(
                a.get("author", {}).get("display_name", "")
                for a in item.get("authorships", [])[:4]
            )
            venue = (item.get("primary_location") or {}).get("source", {})
            results.append({
                "title":    item.get("title", ""),
                "authors":  authors,
                "year":     item.get("publication_year"),
                "doi":      item.get("doi", ""),
                "venue":    venue.get("display_name", "") if venue else "",
                "cited_by": item.get("cited_by_count", 0),
                "type":     item.get("type", ""),
            })
        meta = data.get("meta", {})
        return json.dumps({
            "query": query, "total_found": meta.get("count", len(results)),
            "results": results,
        }, indent=2, ensure_ascii=False)
    except Exception as e:
        return f"[openalex_search error: {e}]"


@mcp.tool()
def web_browser(url: str, wait_seconds: float = 2.0, max_chars: int = 6000) -> str:
    """Render a URL with a JavaScript-capable headless browser (Splash).

    Use when fetch_pdf_text or web_researcher return empty or broken content
    because the page requires JavaScript to render (Single-Page Apps, dynamic
    tables, BBC scripts, British Museum collection pages, GitHub issue timelines).

    Falls back to plain httpx fetch if Splash is unavailable.

    url:          Full URL to render (http/https only)
    wait_seconds: Time to wait after page load for JS execution (default 2s)
    max_chars:    Maximum characters to return from the rendered HTML text
    """
    _assert_public_url(url)
    splash_url = os.environ.get("SPLASH_URL", "http://moe-splash:8050")

    # Try Splash first
    try:
        resp = httpx.get(
            f"{splash_url}/render.html",
            params={"url": url, "wait": wait_seconds, "timeout": 30},
            timeout=40,
        )
        if resp.status_code == 200:
            html = resp.text
            # Strip HTML tags for clean text
            import re as _re
            text = _re.sub(r'<[^>]+>', ' ', html)
            text = _re.sub(r'\s+', ' ', text).strip()
            return text[:max_chars] if len(text) > max_chars else text
        # Splash returned error — fall through to plain fetch
    except Exception:
        pass

    # Fallback: plain httpx (no JS)
    try:
        resp = httpx.get(url, follow_redirects=False, timeout=20,
                         headers={"User-Agent": "Mozilla/5.0 (research-bot/1.0)"})
        resp.raise_for_status()
        import re as _re
        text = _re.sub(r'<[^>]+>', ' ', resp.text)
        text = _re.sub(r'\s+', ' ', text).strip()
        return text[:max_chars] if len(text) > max_chars else text
    except Exception as e:
        return f"[web_browser error: {e}]"


@mcp.tool()
def duckduckgo_search(query: str, max_results: int = 5, region: str = "wt-wt") -> str:
    """Search the web via DuckDuckGo — no API key required, no rate-limit risk.

    Use as a complement or fallback to web_researcher when SearXNG returns
    poor results. DuckDuckGo indexes different sources and applies its own
    relevance ranking independent of SearXNG.

    Returns title, URL, and body snippet for each result.

    query:       Search query string (English preferred for best coverage)
    max_results: Number of results to return (1-10, default 5)
    region:      DuckDuckGo region code (default wt-wt = worldwide)
                 Examples: us-en, de-de, gb-en, fr-fr
    """
    try:
        from ddgs import DDGS
        with DDGS() as ddgs:
            results = list(ddgs.text(query, region=region, max_results=min(max(1, max_results), 10)))
        if not results:
            return f"[duckduckgo_search: no results for '{query[:80]}']"
        output = []
        for r in results:
            output.append({
                "title": r.get("title", ""),
                "url":   r.get("href", ""),
                "body":  r.get("body", "")[:400],
            })
        return json.dumps({"query": query, "results": output}, indent=2, ensure_ascii=False)
    except ImportError:
        return "[duckduckgo_search: ddgs library not installed — run: pip install ddgs]"
    except Exception as e:
        return f"[duckduckgo_search error: {e}]"


# ─── Starfleet Infra Tools (read-only observability) ────────────────────────
# All tools in this section are guarded by the INFRA_MCP_ENABLED flag.
# They call the orchestrator's Starfleet API endpoints via the internal Docker network.
# and therefore work without direct Redis/Prometheus access from this container.

_ORCHESTRATOR_URL = os.getenv("ORCHESTRATOR_INTERNAL_URL", "http://langgraph-app:8000")
_INFRA_MCP_ENABLED = os.getenv("INFRA_MCP_ENABLED", "false").lower() in ("1", "true", "yes")


def _infra_disabled_response(tool_name: str) -> str:
    return (
        f"[{tool_name}: disabled] Set INFRA_MCP_ENABLED=true in .env to enable "
        f"read-only infrastructure observability via MCP."
    )


@mcp.tool()
def node_status(node_name: str = "") -> str:
    """Return health and VRAM status for one or all inference nodes.

    Reads live data from the MoE orchestrator's Starfleet API.
    Includes: up/down state, currently loaded models, recent alerts per node.

    node_name: Optional node name (e.g. 'MY-NODE'). Leave empty for all nodes.
    """
    if not _INFRA_MCP_ENABLED:
        return _infra_disabled_response("node_status")
    try:
        alerts_resp = httpx.get(f"{_ORCHESTRATOR_URL}/api/watchdog/alerts?limit=20", timeout=5.0)
        features_resp = httpx.get(f"{_ORCHESTRATOR_URL}/api/starfleet/features", timeout=5.0)

        alerts_by_node: dict = {}
        if alerts_resp.status_code == 200:
            for alert in alerts_resp.json().get("alerts", []):
                node = alert.get("node", "unknown")
                alerts_by_node.setdefault(node, []).append(alert.get("message", ""))

        nodes: dict = {}
        if features_resp.status_code == 200:
            for server in features_resp.json().get("inference_servers", []):
                name = server.get("name", "unknown")
                nodes[name] = {
                    "enabled": server.get("enabled", False),
                    "gpu_count": server.get("gpu_count", 0),
                    "loaded_models": server.get("loaded_models", []),
                    "vram_used_pct": server.get("vram_used_pct"),
                    "alerts": alerts_by_node.get(name, []),
                }

        if node_name:
            filtered = {k: v for k, v in nodes.items() if node_name.lower() in k.lower()}
            result = filtered if filtered else {"error": f"No node matching '{node_name}' found"}
        else:
            result = nodes if nodes else {"note": "No inference nodes configured or Starfleet unavailable"}

        return json.dumps(result, indent=2)
    except Exception as exc:
        return f"[node_status error: {exc}]"


@mcp.tool()
def active_requests() -> str:
    """Return the count and IDs of LLM requests currently in flight.

    Reads the moe:active:* key pattern from the orchestrator's live state.
    Useful for understanding system load before issuing expensive queries.
    """
    if not _INFRA_MCP_ENABLED:
        return _infra_disabled_response("active_requests")
    try:
        resp = httpx.get(f"{_ORCHESTRATOR_URL}/api/watchdog/alerts?limit=1", timeout=5.0)
        # The orchestrator exposes active count via the watchdog endpoint's metadata.
        # Fall back to the features endpoint which includes the enabled state.
        feat = httpx.get(f"{_ORCHESTRATOR_URL}/api/starfleet/features", timeout=5.0)
        if feat.status_code == 200:
            return json.dumps({"status": "ok", "features": feat.json()}, indent=2)
        return json.dumps({"status": "ok", "note": "active request count requires watchdog enabled"})
    except Exception as exc:
        return f"[active_requests error: {exc}]"


@mcp.tool()
def mission_context_get() -> str:
    """Return the current persistent mission context (cross-session project state).

    Shows the active project title, open tasks, recent decisions, and tags.
    Use this at the start of a session to re-establish context without repeating yourself.
    """
    if not _INFRA_MCP_ENABLED:
        return _infra_disabled_response("mission_context_get")
    try:
        resp = httpx.get(f"{_ORCHESTRATOR_URL}/api/mission-context", timeout=5.0)
        if resp.status_code == 200:
            ctx = resp.json()
            if not ctx.get("enabled", True):
                return "[mission_context_get: Mission Context feature is disabled on the server]"
            return json.dumps(ctx, indent=2, ensure_ascii=False)
        return f"[mission_context_get: HTTP {resp.status_code}]"
    except Exception as exc:
        return f"[mission_context_get error: {exc}]"


@mcp.tool()
def watchdog_alerts(limit: int = 10) -> str:
    """Return recent watchdog alerts (node down, VRAM high, benchmark stuck, etc.).

    Alerts are stored in Valkey and persisted across requests.
    Most recent alert is first. Use limit to control how many entries to fetch (max 100).
    """
    if not _INFRA_MCP_ENABLED:
        return _infra_disabled_response("watchdog_alerts")
    try:
        limit = max(1, min(100, limit))
        resp = httpx.get(
            f"{_ORCHESTRATOR_URL}/api/watchdog/alerts",
            params={"limit": limit},
            timeout=5.0,
        )
        if resp.status_code == 200:
            data = resp.json()
            if not data.get("enabled"):
                return "[watchdog_alerts: Watchdog feature is disabled on the server]"
            alerts = data.get("alerts", [])
            if not alerts:
                return json.dumps({"status": "nominal", "alerts": [], "message": "No alerts on record."})
            return json.dumps({"status": "alerts_present", "count": len(alerts), "alerts": alerts}, indent=2)
        return f"[watchdog_alerts: HTTP {resp.status_code}]"
    except Exception as exc:
        return f"[watchdog_alerts error: {exc}]"


@mcp.tool()
def chess_analyze_position(fen: str, top_moves: int = 3) -> str:
    """Analyze a chess position and return the best moves using Lichess cloud evaluation (Stockfish).

    Uses the Lichess public API — no API key required. Returns top moves with centipawn scores.
    Covers ~342 million positions evaluated at depth 20-99.

    fen: FEN string of the chess position (e.g. 'rnbqkbnr/pppppppp/8/8/4P3/8/PPPP1PPP/RNBQKBNR b KQkq e3 0 1')
    top_moves: Number of best moves to return (1-5, default 3)
    """
    fen = fen.strip()
    if not fen:
        return "[chess_analyze_position: empty FEN string]"

    # Basic FEN validation: must have 6 space-separated parts
    parts = fen.split()
    if len(parts) < 4:
        return f"[chess_analyze_position: invalid FEN — expected at least 4 parts, got {len(parts)}]"

    top_moves = max(1, min(5, int(top_moves)))

    try:
        r = httpx.get(
            "https://lichess.org/api/cloud-eval",
            params={"fen": fen, "multiPv": top_moves},
            headers={"Accept": "application/json"},
            timeout=12.0,
        )
        if r.status_code == 404:
            # Position not in Lichess cloud database — try with python-chess + stockfish fallback
            return (
                f"[chess_analyze_position: position not in Lichess cloud database. "
                f"FEN: {fen}. Use the FEN to reason about legal moves directly.]"
            )
        if r.status_code != 200:
            return f"[chess_analyze_position: Lichess API error {r.status_code}]"

        data = r.json()
        pvs = data.get("pvs", [])
        if not pvs:
            return f"[chess_analyze_position: no evaluation available for FEN: {fen}]"

        side = parts[1] if len(parts) > 1 else "?"
        side_str = "White" if side == "w" else "Black" if side == "b" else side
        depth = data.get("depth", "?")

        lines = [f"Position analysis (depth {depth}, {side_str} to move):"]
        for i, pv in enumerate(pvs, 1):
            moves = pv.get("moves", "").split()
            best = moves[0] if moves else "?"
            cp = pv.get("cp")
            mate = pv.get("mate")
            if mate is not None:
                score = f"mate in {abs(mate)}"
            elif cp is not None:
                score = f"{cp/100:+.2f}" if side == "w" else f"{-cp/100:+.2f}"
            else:
                score = "?"
            lines.append(f"  Move {i}: {best}  (eval: {score})  continuation: {' '.join(moves[:5])}")

        return "\n".join(lines)
    except Exception as e:
        return f"[chess_analyze_position: error — {e}]"


@mcp.tool()
def chess_legal_moves(fen: str) -> str:
    """Return all legal moves for a chess position given its FEN string.

    Uses python-chess for exact legal move generation — no external API needed.
    Returns moves in UCI notation (e.g. 'e2e4') and SAN notation (e.g. 'e4').

    fen: FEN string of the chess position
    """
    try:
        import chess
        board = chess.Board(fen.strip())
        legal = list(board.legal_moves)
        san_moves = [board.san(m) for m in legal]
        uci_moves = [m.uci() for m in legal]
        side = "White" if board.turn == chess.WHITE else "Black"
        in_check = board.is_check()
        return (
            f"{side} to move. {'In check. ' if in_check else ''}"
            f"{len(legal)} legal moves:\n"
            f"SAN: {', '.join(sorted(san_moves))}\n"
            f"UCI: {', '.join(uci_moves)}"
        )
    except ImportError:
        return "[chess_legal_moves: python-chess not installed — pip install chess]"
    except Exception as e:
        return f"[chess_legal_moves: error parsing FEN — {e}]"


# ─── Session Context Search ──────────────────────────────────────────────────

_ORCHESTRATOR_URL = os.environ.get("ORCHESTRATOR_URL", "http://langgraph-orchestrator:8000").rstrip("/")


@mcp.tool()
async def search_context(query: str, session_id: str, n_results: int = 8) -> str:
    """Retrieve semantically relevant chunks from the session's indexed context.

    Use this tool when you need to find specific information from a large
    codebase, document, or system prompt that was provided at session start.
    The context was automatically indexed at the beginning of the session.

    query:      What you are looking for (natural language or code snippet)
    session_id: The current session ID (available in the system prompt header)
    n_results:  Number of chunks to return (default 8, max 20)
    """
    if not session_id or not query:
        return "[search_context: session_id and query are required]"
    try:
        async with httpx.AsyncClient(timeout=10.0) as client:
            resp = await client.post(
                f"{_ORCHESTRATOR_URL}/v1/context/search",
                json={"session_id": session_id, "query": query[:500], "n_results": min(n_results, 20)},
            )
            data = resp.json()
            if not data.get("indexed"):
                return "[search_context: no context index available for this session — either the context is below the indexing threshold or indexing is still in progress]"
            chunks = data.get("chunks", "")
            if not chunks:
                return "[search_context: no relevant chunks found for this query]"
            return chunks
    except Exception as exc:
        return f"[search_context: retrieval error — {exc}]"


# ─── PM Connector Tools ──────────────────────────────────────────────────────

import pm_connector as _pm


@mcp.tool()
async def pm_create_task(
    title: str,
    description: str = "",
    labels: str = "",
    priority: str = "medium",
    assignee: str = "",
) -> str:
    """
    Creates a task/issue in the configured PM system (Linear, GitHub Issues, or webhook).
    Returns JSON with id, title, url, and state of the created task.

    title:       Task title (required)
    description: Markdown description / acceptance criteria
    labels:      Comma-separated labels (e.g. "bug,backend") — 'moe-ai' is always added
    priority:    urgent | high | medium | low | none  (default: medium; GitHub ignores this)
    assignee:    Username to assign (GitHub login or Linear user name)

    Requires PM_BACKEND, PM_API_KEY, PM_PROJECT_ID in .env.
    """
    label_list = [l.strip() for l in labels.split(",") if l.strip()] if labels else []
    return await _pm.create_task(title, description, label_list, priority, assignee)


@mcp.tool()
async def pm_list_tasks(
    status: str = "",
    assignee: str = "",
    label: str = "",
    limit: int = 20,
) -> str:
    """
    Lists tasks from the configured PM system. Returns JSON array of tasks.

    status:   Filter by state — Linear: 'Todo' | 'In Progress' | 'Done' | 'Backlog'
                                GitHub:  'open' | 'closed' | 'all'
    assignee: Filter by assignee username
    label:    Filter by label name
    limit:    Max results (default: 20)
    """
    return await _pm.list_tasks(status, assignee, label, limit)


@mcp.tool()
async def pm_update_task(task_id: str, status: str = "", comment: str = "") -> str:
    """
    Updates a task's status and/or adds a comment in the configured PM system.

    task_id: Issue number (GitHub) or Linear ID/UUID (e.g. 'ENG-42' or UUID)
    status:  New state — Linear: 'Todo' | 'In Progress' | 'Done' | 'Cancelled'
                         GitHub: 'open' | 'closed' | 'done' (maps → closed)
    comment: Comment text to append to the issue/task
    """
    return await _pm.update_task(task_id, status, comment)


@mcp.tool()
async def pm_search_tasks(query: str, limit: int = 10) -> str:
    """
    Full-text search across tasks/issues in the configured PM system.
    Returns JSON array of matching tasks with id, title, state, and url.

    query: Search terms (title, description, labels)
    limit: Max results (default: 10)
    """
    return await _pm.search_tasks(query, limit)


# ─── Generative Tools (TASK-52: N04-RGTX Hardware Offloading) ─────────────

COMFYUI_URL = os.getenv("COMFYUI_URL", "http://192.168.155.224:8188")
KOKORO_TTS_URL = os.getenv("KOKORO_TTS_URL", "http://192.168.155.224:8880")


@mcp.tool()
async def generate_image(prompt: str, size: str = "1024x1024", model: str = "flux-schnell") -> Dict[str, Any]:
    """Generate an image using local ComfyUI API endpoint on N04-RGTX.

    prompt: Detailed text description of the image to generate.
    size: Image dimensions, e.g. '1024x1024', '512x512', or '1280x720'.
    model: Checkpoint model name, defaults to 'flux-schnell'.
    """
    if not prompt or not prompt.strip():
        return {"error": "Prompt cannot be empty", "error_code": "invalid_prompt"}

    parts = size.split("x") if "x" in size else [1024, 1024]
    try:
        width = int(parts[0])
        height = int(parts[1]) if len(parts) > 1 else int(parts[0])
    except ValueError:
        width, height = 1024, 1024

    payload = {
        "prompt": prompt,
        "width": width,
        "height": height,
        "model": model,
    }

    url = f"{COMFYUI_URL.rstrip('/')}/v1/images/generations"
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            resp = await client.post(url, json=payload)
            if resp.status_code == 200:
                return resp.json()
            return {
                "error": f"ComfyUI HTTP {resp.status_code}: {resp.text[:300]}",
                "error_code": "comfyui_error",
                "status_code": resp.status_code,
            }
    except Exception as exc:
        return {
            "error": f"Failed to connect to ComfyUI backend at {url}: {exc}",
            "error_code": "backend_unreachable",
            "target_url": url,
            "gpu_node": "N04-RGTX",
        }


@mcp.tool()
async def generate_speech(text: str, voice: str = "af_heart", model: str = "kokoro") -> Dict[str, Any]:
    """Generate speech audio using local Kokoro-TTS API endpoint on N04-RGTX.

    text: Text content to convert into spoken audio.
    voice: Voice profile, defaults to 'af_heart'.
    model: TTS model name, defaults to 'kokoro'.
    """
    if not text or not text.strip():
        return {"error": "Text cannot be empty", "error_code": "invalid_text"}

    payload = {
        "input": text,
        "voice": voice,
        "model": model,
        "response_format": "mp3",
    }

    url = f"{KOKORO_TTS_URL.rstrip('/')}/v1/audio/speech"
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            resp = await client.post(url, json=payload)
            if resp.status_code == 200:
                return {
                    "ok": True,
                    "voice": voice,
                    "model": model,
                    "audio_bytes_length": len(resp.content),
                    "content_type": resp.headers.get("content-type", "audio/mpeg"),
                }
            return {
                "error": f"Kokoro TTS HTTP {resp.status_code}: {resp.text[:300]}",
                "error_code": "kokoro_error",
                "status_code": resp.status_code,
            }
    except Exception as exc:
        return {
            "error": f"Failed to connect to Kokoro TTS backend at {url}: {exc}",
            "error_code": "backend_unreachable",
            "target_url": url,
            "gpu_node": "N04-RGTX",
        }


# ─── Tool registry for REST shim ────────────────────────────────────────────

_TOOL_REGISTRY: Dict[str, Any] = {
    "generate_image": generate_image,
    "generate_speech": generate_speech,
    "calculate": calculate,
    "solve_equation": solve_equation,
    "date_diff": date_diff,
    "date_add": date_add,
    "calendar_facts": calendar_facts,
    "time_facts": time_facts,
    "timezone_convert": timezone_convert,
    "decimal_finance": decimal_finance,
    "exact_probability": exact_probability,
    "structured_validate": structured_validate,
    "rust_compile_check": rust_compile_check,
    "day_of_week": day_of_week,
    "unit_convert": unit_convert,
    "statistics_calc": statistics_calc,
    "hash_text": hash_text,
    "base64_codec": base64_codec,
    "regex_extract": regex_extract,
    "subnet_calc": subnet_calc,
    "text_analyze": text_analyze,
    "prime_factorize": prime_factorize,
    "gcd_lcm": gcd_lcm,
    "json_query": json_query,
    "roman_numeral": roman_numeral,
    "legal_search_laws":      legal_search_laws,
    "legal_get_law_overview": legal_get_law_overview,
    "legal_get_paragraph":    legal_get_paragraph,
    "legal_fulltext_search":  legal_fulltext_search,
    "graph_query":      graph_query,
    "graph_ingest":     graph_ingest,
    "graph_provenance": graph_provenance,
    # Code-Navigation (Agentic Coder)
    "repo_map":          repo_map,
    "read_file_chunked": read_file_chunked,
    "lsp_query":         lsp_query,
    "generate_file":     generate_file,
    "parse_attachment":  parse_attachment,
    "graph_analyze":     graph_analyze,
    "file_upload":       file_upload,
    "file_download_url": file_download_url,
    "fetch_pdf_text":        fetch_pdf_text,
    "github_get_issue":      github_get_issue,
    "wikipedia_get_section": wikipedia_get_section,
    "python_sandbox":        python_sandbox,
    "grid_repr":             grid_repr,
    # External data sources (added for adaptive deep research)
    "web_search_domain":       web_search_domain,
    "youtube_transcript":      youtube_transcript,
    "github_search_issues":    github_search_issues,
    "github_issue_events":     github_issue_events,
    "pubchem_compound_search":  pubchem_compound_search,
    "pubchem_advanced_search":  pubchem_advanced_search,
    "orcid_works_count":       orcid_works_count,
    "semantic_scholar_search": semantic_scholar_search,
    "wikidata_search":         wikidata_search,
    "wikidata_sparql":         wikidata_sparql,
    "pubmed_search":           pubmed_search,
    "duckduckgo_search":       duckduckgo_search,
    "web_browser":             web_browser,
    "wayback_fetch":           wayback_fetch,
    "crossref_lookup":         crossref_lookup,
    "openalex_search":         openalex_search,
    # Chess analysis via Lichess public API + python-chess
    "chess_analyze_position": chess_analyze_position,
    "chess_legal_moves":      chess_legal_moves,
    # Starfleet Infra Tools (read-only — enabled via INFRA_MCP_ENABLED)
    "node_status":          node_status,
    "active_requests":      active_requests,
    "mission_context_get":  mission_context_get,
    "watchdog_alerts":      watchdog_alerts,
    "search_context":       search_context,
    # PM Connector
    "pm_create_task":  pm_create_task,
    "pm_list_tasks":   pm_list_tasks,
    "pm_update_task":  pm_update_task,
    "pm_search_tasks": pm_search_tasks,
}

_TOOL_DESCRIPTIONS = {
    "calculate": "Exact arithmetic, percentages, formulas (LLMs hallucinate with numbers!)",
    "solve_equation": "Solve algebraic equations (SymPy)",
    "date_diff": "Exact difference between two dates",
    "date_add": "Date addition/subtraction (days, months, years)",
    "calendar_facts": "Deterministic localized weekday and ISO calendar facts as JSON",
    "time_facts": "Explicit-instant IANA timezone, offset, DST and ISO calendar facts",
    "timezone_convert": "DST-safe conversion between explicit IANA timezones",
    "decimal_finance": "Decimal-string finance arithmetic with explicit currency, scale and rounding",
    "exact_probability": "Exact bounded rational probability and combinatorics with optional Decimal projection",
    "structured_validate": "Network-free bounded JSON, YAML, XML and CSV parser/validator",
    "rust_compile_check": "Type/borrow-check Rust source in an isolated, network-free sandbox (analysis only, never executes the code)",
    "day_of_week": "Weekday, calendar week, day of year for a date",
    "unit_convert": "Physical unit conversion (km/h→m/s, °F→°C, etc.)",
    "statistics_calc": "Statistical measures for data sets (mean, median, stdev, etc.)",
    "hash_text": "Cryptographic hashes (MD5/SHA256/SHA512)",
    "base64_codec": "Base64 encode/decode",
    "regex_extract": "Regex pattern matching and extraction",
    "subnet_calc": "IP/network calculations (CIDR, subnet mask, host range)",
    "text_analyze": "Text metrics (words, characters, sentences, reading time)",
    "prime_factorize": "Prime factorization",
    "gcd_lcm": "GCD and LCM of two numbers",
    "json_query": "JSON path queries (key.subkey, array[0])",
    "roman_numeral": "Arabic ↔ Roman numerals",
    "legal_search_laws":      "Search German federal laws by keyword (returns abbreviations)",
    "legal_get_law_overview": "Shows table of contents/structure of a German federal law",
    "legal_get_paragraph":    "Retrieves exact text of a §/Art. from a German federal law (BGB/StGB/GG etc.)",
    "legal_fulltext_search":  "Full-text search within a German federal law by keyword",
    "graph_query":      "Search Neo4j knowledge graph for entities/relations (with provenance)",
    "graph_ingest":     "Store facts from Q&A in Neo4j knowledge graph (for external agents)",
    "graph_provenance": "Shows version history of an entity (contradiction analysis, temporal RAG)",
    # Code navigation (only for agentic_coder — not in global MCP_TOOLS_DESCRIPTION)
    "repo_map":          "AST/regex skeleton of a repo (file paths + classes/functions, no code)",
    "read_file_chunked": "Paginated file reading (start_line/end_line) — prevents context overflow",
    "lsp_query":         "Python LSP features: signature, find_references, completions (.py only)",
    "generate_file":     "Generate downloadable files (HTML, DOCX, PPTX, Markdown, TXT, PDF) from content — returns MinIO pre-signed URL. Use format='pdf' for proper PDF output.",
    "parse_attachment":  "Download and parse file attachments (XLSX→CSV, DOCX→text, PDF→text, CSV→text) — max 20 MB",
    "graph_analyze":     "Analyze a graph (Eulerian path/circuit, connected components, degree map, density) from text description",
    "file_upload":       "Upload a file (base64-encoded) to MinIO object storage and get a 24h pre-signed download URL",
    "file_download_url": "Generate a fresh pre-signed download URL for an existing file in MinIO storage",
    "fetch_pdf_text":        "Download a PDF from a URL and extract its text (up to 8 000 chars by default); auto-handles arXiv URLs",
    "github_get_issue":      "Fetch a GitHub issue (title, state, body, labels, comment count) via the public API",
    "wikipedia_get_section": "Fetch a specific section of a Wikipedia article as plain text (e.g. 'Discography', 'Filmography'). Use this whenever a question references a Wikipedia article.",
    "python_sandbox":        "Run a small Python snippet for exact numerical calculations: probability trees (use Fraction!), Markov chains, combinatorics. Use print() for output. NEVER write simulation/Monte Carlo code — always use exact Fraction arithmetic. Allowed modules: math, fractions, itertools, collections, decimal, statistics, random.",
    "grid_repr":             "Render a 2-D integer grid as annotated ASCII art with a colour legend. Pass data as a JSON 2-D list string (e.g. '[[0,1,2],[3,0,1]]'). Values 0–9 use the ARC colour palette (0=black, 1=blue, 2=red, 3=green, 4=yellow, 5=grey, 6=pink, 7=orange, 8=azure, 9=maroon). Use this to visualise ARC-AGI grids or any integer matrix before reasoning about it.",
    # External data sources
    "web_search_domain":       "Domain-restricted web search via SearXNG (site:github.com, site:arxiv.org, site:wikipedia.org, etc.). Use when general search fails and you need data from a specific known website.",
    "youtube_transcript":      "Fetch captions/transcript of a YouTube video by URL or video ID. Use for questions about video content, interviews, documentaries, lectures.",
    "chess_analyze_position":  "Analyze a chess position (FEN string) via Lichess Stockfish cloud evaluation. Returns best moves with centipawn scores. Use when a chess position is given as FEN notation.",
    "chess_legal_moves":       "Return all legal moves for a chess position given its FEN string (python-chess). Use to verify legal moves or enumerate options before analysis.",
    "github_search_issues":    "Search GitHub issues by repo, label, state, and text query. Use to find issues when you don't know the exact issue number (e.g. oldest regression issue in numpy/numpy with label 'Regression').",
    "github_issue_events":     "Fetch timeline events for a GitHub issue — use to find WHEN a label was added/removed. Essential for 'when was label X added to issue Y' questions. Returns event type, date, actor, and label name.",
    "pubchem_compound_search":  "Search PubChem compound database by name, CID, or molecular weight range. Returns CID, molecular weight, formula, IUPAC name. Use for chemistry/pharmacology/food-science questions.",
    "pubchem_advanced_search":  "Advanced PubChem multi-criteria search: filter by MW range, heavy atom count, HB acceptors/donors simultaneously. Use for GAIA-style questions like 'find the compound with MW≤100, 6 heavy atoms, ≤1 HB acceptors in Food Additive Status'.",
    "orcid_works_count":       "Count publications on an ORCID researcher profile. Optionally filter by year (e.g. before_year=2020 for pre-2020 works). Use for academic publication count questions.",
    "semantic_scholar_search": "Search Semantic Scholar for academic papers by author/title/topic. Returns abstracts, DOI, and open-access PDF links. Use for questions requiring specific measurements or data from named academic papers (e.g. 'Valencfia-Mendez 2017 harlequin shrimp length'). Set fetch_pdf=true to also extract the paper's text if a free PDF is available.",
    "wikidata_search":         "Search Wikidata by text to find entity IDs (wd:Q...) for use in wikidata_sparql. Use this first when you know the name but not the Wikidata ID.",
    "wikidata_sparql":         "Execute SPARQL against Wikidata for deterministic entity facts (dates, locations, people, species, relationships). ALWAYS prefer over web search for factual lookups — no HTML parsing, no SearXNG variance. Write a SPARQL 1.1 query; use wd: entity IDs or wdt: property filters.",
    "pubmed_search":           "Search PubMed/NCBI for biomedical, biology, ecology, and life science papers. Returns title, authors, year, abstract, DOI. Prefer over semantic_scholar_search for species studies, genetics, clinical trials, ecology — deterministic NCBI API, no SearXNG variance.",
    "duckduckgo_search":       "Search the web via DuckDuckGo (no API key, no rate-limit risk). Use as complement or fallback when SearXNG returns poor results — DuckDuckGo indexes different sources. Returns title, URL, snippet for each result. Good for English-language factual queries.",
    "web_browser":             "Render a URL with a JavaScript-capable headless browser (Splash). Use when fetch_pdf_text or web_researcher return broken/empty content because the page needs JS to render: BBC scripts, British Museum collection pages, GitHub issue timelines, museum databases. Falls back to plain HTTP fetch if Splash unavailable.",
    "wayback_fetch":           "Retrieve a historical snapshot of any web page from the Wayback Machine (web.archive.org). Use for: ORCID profiles at a past date, pages that have changed, historical API data. timestamp=YYYYMMDD optional — omit for closest available. Ideal for 'as of year X' questions.",
    "crossref_lookup":         "Search CrossRef for 150M+ scholarly publications — title, author, DOI, keyword. Returns DOI, authors, year, journal, citation count. Use to count articles by venue/year or verify publication metadata. No API key, deterministic.",
    "openalex_search":         "Search OpenAlex academic database (250M+ works, all disciplines). Broader than PubMed/SemanticScholar. Use for cross-disciplinary paper counts, author publication histories, funding data. Supports year_min filter and open_access_only. No API key.",
    "search_context":          "Retrieve semantically relevant chunks from the session's indexed context (codebase, documents, large system prompt). Use when you need to find specific information in a large context that was provided at session start. Returns the top-k most relevant sections.",
    # PM Connector
    "pm_create_task":  "Create a task/issue in the configured PM system (Linear | GitHub Issues | webhook). Returns JSON with id, title, url. Requires PM_BACKEND, PM_API_KEY, PM_PROJECT_ID.",
    "pm_list_tasks":   "List open tasks from the configured PM system. Filter by status, assignee, label. Returns JSON array.",
    "pm_update_task":  "Update a task's status and/or add a comment in the configured PM system. task_id = issue number (GitHub) or Linear identifier.",
    "pm_search_tasks": "Full-text search across tasks/issues in the configured PM system. Returns JSON array with id, title, state, url.",
}

# ─── ACCESS-KIND CLASSIFICATION (visibility only — see graph/tool_nodes.py) ──
# Categorizes each tool by real-world side effect, not by name:
#   read    — reads data/state, or pure local computation with no side effect
#   search  — outbound network call to an external/public service (egress-relevant,
#             non-deterministic content from an uncontrolled third party)
#   write   — creates/mutates persistent state (local or in a configured backend)
#   execute — runs arbitrary/interpreted code
# Used purely for telemetry (services/decision_log.py DecisionType.MCP_TOOL_ACCESS +
# metrics.py PROM_MCP_TOOL_ACCESS) — does NOT gate or block any tool call.
_DEFAULT_ACCESS_KIND = "read"

_TOOL_ACCESS_KIND: Dict[str, str] = {
    # Generative AI tools (local hardware offloading on N04-RGTX)
    "generate_image": "write", "generate_speech": "write",
    # Math/utility — local computation only
    "calculate": "read", "solve_equation": "read", "date_diff": "read",
    "date_add": "read", "calendar_facts": "read", "time_facts": "read",
    "timezone_convert": "read", "decimal_finance": "read",
    "exact_probability": "read", "structured_validate": "read",
    "day_of_week": "read",
    "unit_convert": "read",
    "statistics_calc": "read", "hash_text": "read", "base64_codec": "read",
    "regex_extract": "read", "subnet_calc": "read", "text_analyze": "read",
    "prime_factorize": "read", "gcd_lcm": "read", "json_query": "read",
    "roman_numeral": "read",
    # Legal — internal DB lookups
    "legal_search_laws": "read", "legal_get_law_overview": "read",
    "legal_get_paragraph": "read", "legal_fulltext_search": "read",
    # Graph — Neo4j
    "graph_query": "read", "graph_ingest": "write", "graph_provenance": "read",
    # Code navigation
    "repo_map": "read", "read_file_chunked": "read", "lsp_query": "read",
    "generate_file": "write", "parse_attachment": "read", "graph_analyze": "read",
    "file_upload": "write", "file_download_url": "read",
    # External data sources — outbound network calls
    "fetch_pdf_text": "search", "github_get_issue": "search",
    "wikipedia_get_section": "search", "web_search_domain": "search",
    "youtube_transcript": "search", "github_search_issues": "search",
    "github_issue_events": "search", "pubchem_compound_search": "search",
    "pubchem_advanced_search": "search", "orcid_works_count": "search",
    "semantic_scholar_search": "search", "wikidata_search": "search",
    "wikidata_sparql": "search", "pubmed_search": "search",
    "duckduckgo_search": "search", "web_browser": "search",
    "wayback_fetch": "search", "crossref_lookup": "search",
    "openalex_search": "search",
    # Code execution / local computation
    "python_sandbox": "execute",
    "rust_compile_check": "execute",
    "grid_repr":      "read",
    # Chess — chess_analyze_position calls the external Lichess cloud-eval API;
    # chess_legal_moves is local python-chess computation
    "chess_analyze_position": "search", "chess_legal_moves": "read",
    # Starfleet infra tools — local/internal status reads
    "node_status": "read", "active_requests": "read",
    "mission_context_get": "read", "watchdog_alerts": "read",
    "search_context": "read",
    # PM Connector — configured/trusted backend, not open web search
    "pm_create_task": "write", "pm_list_tasks": "read",
    "pm_update_task": "write", "pm_search_tasks": "read",
}


_CALENDAR_FACTS_OUTPUT_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "calendar_system": {"type": "string", "const": "proleptic_gregorian"},
        "date": {"type": "string", "pattern": r"^\d{4}-\d{2}-\d{2}$"},
        "day": {"type": "integer", "minimum": 1, "maximum": 31},
        "day_of_year": {"type": "integer", "minimum": 1, "maximum": 366},
        "days_in_month": {"type": "integer", "minimum": 28, "maximum": 31},
        "days_in_year": {"type": "integer", "enum": [365, 366]},
        "is_leap_year": {"type": "boolean"},
        "is_weekend": {"type": "boolean"},
        "iso_week": {"type": "integer", "minimum": 1, "maximum": 53},
        "iso_week_year": {"type": "integer"},
        "locale": {"type": "string", "enum": ["de", "en"]},
        "month": {"type": "integer", "minimum": 1, "maximum": 12},
        "month_name": {"type": "string", "minLength": 1},
        "quarter": {"type": "integer", "minimum": 1, "maximum": 4},
        "weekday_iso": {"type": "integer", "minimum": 1, "maximum": 7},
        "weekday_name": {"type": "string", "minLength": 1},
        "year": {"type": "integer"},
    },
    "required": [
        "calendar_system", "date", "day", "day_of_year", "days_in_month",
        "days_in_year", "is_leap_year", "is_weekend", "iso_week",
        "iso_week_year", "locale", "month", "month_name", "quarter",
        "weekday_iso", "weekday_name", "year",
    ],
    "additionalProperties": False,
}

_TIME_FACTS_OUTPUT_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "input_instant": {"type": "string", "minLength": 20, "maxLength": 32},
        "utc_instant": {"type": "string", "pattern": r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}\+00:00$"},
        "as_of": {"type": "string", "pattern": r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}\+00:00$"},
        "timezone": {"type": "string", "minLength": 1, "maxLength": 128},
        "local_datetime": {"type": "string", "minLength": 20, "maxLength": 32},
        "utc_offset": {"type": "string", "pattern": r"^[+-]\d{2}:\d{2}$"},
        "utc_offset_seconds": {"type": "integer", "minimum": -86400, "maximum": 86400},
        "timezone_abbreviation": {"type": "string", "maxLength": 32},
        "is_dst": {"type": "boolean"},
        "fold": {"type": "integer", "enum": [0, 1]},
        "tzdata_version": {"type": "string", "minLength": 1, "maxLength": 32},
        "date": {"type": "string", "pattern": r"^\d{4}-\d{2}-\d{2}$"},
        "time": {"type": "string", "pattern": r"^\d{2}:\d{2}:\d{2}$"},
        "weekday_iso": {"type": "integer", "minimum": 1, "maximum": 7},
        "weekday_name": {"type": "string", "minLength": 1},
        "iso_week": {"type": "integer", "minimum": 1, "maximum": 53},
        "iso_week_year": {"type": "integer", "minimum": _TIME_YEAR_MIN - 1, "maximum": _TIME_YEAR_MAX + 1},
        "locale": {"type": "string", "enum": ["de", "en"]},
    },
    "required": [
        "input_instant", "utc_instant", "as_of", "timezone",
        "local_datetime", "utc_offset", "utc_offset_seconds",
        "timezone_abbreviation", "is_dst", "fold", "tzdata_version",
        "date", "time", "weekday_iso", "weekday_name", "iso_week",
        "iso_week_year", "locale",
    ],
    "additionalProperties": False,
}

_TIMEZONE_CONVERT_OUTPUT_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "input_local_datetime": {"type": "string", "pattern": r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}$"},
        "from_timezone": {"type": "string", "minLength": 1, "maxLength": 128},
        "to_timezone": {"type": "string", "minLength": 1, "maxLength": 128},
        "fold": {"type": "integer", "enum": [0, 1]},
        "ambiguous": {"type": "boolean"},
        "source_datetime": {"type": "string", "minLength": 20, "maxLength": 32},
        "source_utc_offset": {"type": "string", "pattern": r"^[+-]\d{2}:\d{2}$"},
        "source_is_dst": {"type": "boolean"},
        "utc_instant": {"type": "string", "pattern": r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}\+00:00$"},
        "as_of": {"type": "string", "pattern": r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}\+00:00$"},
        "target_datetime": {"type": "string", "minLength": 20, "maxLength": 32},
        "target_utc_offset": {"type": "string", "pattern": r"^[+-]\d{2}:\d{2}$"},
        "target_is_dst": {"type": "boolean"},
        "target_fold": {"type": "integer", "enum": [0, 1]},
        "locale": {"type": "string", "enum": ["de", "en"]},
        "tzdata_version": {"type": "string", "minLength": 1, "maxLength": 32},
    },
    "required": [
        "input_local_datetime", "from_timezone", "to_timezone", "fold",
        "ambiguous", "source_datetime", "source_utc_offset",
        "source_is_dst", "utc_instant", "as_of", "target_datetime",
        "target_utc_offset", "target_is_dst", "target_fold", "locale",
        "tzdata_version",
    ],
    "additionalProperties": False,
}

_DECIMAL_STRING_SCHEMA: Dict[str, Any] = {
    "type": "string",
    "pattern": r"^-?(?:0|[1-9]\d{0,47})(?:\.\d{1,24})?$",
    "maxLength": 74,
}
_ROUNDING_SCHEMA: Dict[str, Any] = {
    "type": "string",
    "enum": sorted(_DECIMAL_ROUNDING),
}
_DECIMAL_FINANCE_OUTPUT_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "operation": {
            "type": "string",
            "enum": [
                "add", "subtract", "multiply", "divide", "percentage",
                "simple_interest", "compound_interest",
            ],
        },
        "operands": {
            "type": "array", "items": _DECIMAL_STRING_SCHEMA,
            "minItems": 2, "maxItems": 4,
        },
        "currency": {"type": "string", "pattern": r"^[A-Z]{3}$"},
        "scale": {"type": "integer", "minimum": 0, "maximum": 12},
        "rounding": _ROUNDING_SCHEMA,
        "calculation_precision": {"type": "integer", "const": _DECIMAL_CONTEXT_PRECISION},
        "calculation_value": {"type": "string", "minLength": 1, "maxLength": 400},
        "result": {"type": "string", "pattern": r"^-?\d+(?:\.\d+)?$", "maxLength": 400},
    },
    "required": [
        "operation", "operands", "currency", "scale", "rounding",
        "calculation_precision", "calculation_value", "result",
    ],
    "additionalProperties": False,
}

_EXACT_PROBABILITY_OUTPUT_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "operation": {
            "type": "string",
            "enum": ["fraction", "combination", "permutation", "binomial_probability"],
        },
        "n": {"type": ["integer", "null"], "minimum": 0, "maximum": _PROBABILITY_MAX_N},
        "k": {"type": ["integer", "null"], "minimum": 0, "maximum": _PROBABILITY_MAX_N},
        "probability_numerator": {"type": ["integer", "null"]},
        "probability_denominator": {"type": ["integer", "null"]},
        "result_numerator": {"type": "integer"},
        "result_denominator": {"type": "integer", "minimum": 1},
        "fraction": {"type": "string", "pattern": r"^-?\d+/\d+$"},
        "decimal_scale": {"type": ["integer", "null"], "minimum": 0, "maximum": 18},
        "rounding": {"type": ["string", "null"], "enum": [*sorted(_DECIMAL_ROUNDING), None]},
        "decimal": {"type": ["string", "null"], "pattern": r"^-?\d+(?:\.\d+)?$"},
    },
    "required": [
        "operation", "n", "k", "probability_numerator",
        "probability_denominator", "result_numerator", "result_denominator",
        "fraction", "decimal_scale", "rounding", "decimal",
    ],
    "additionalProperties": False,
}

_VALIDATION_DIAGNOSTIC_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "code": {"type": "string", "minLength": 1, "maxLength": 300},
        "message": {"type": "string", "maxLength": 300},
        "line": {"type": ["integer", "null"], "minimum": 1},
        "column": {"type": ["integer", "null"], "minimum": 1},
        "path": {"type": "string", "maxLength": 300},
    },
    "required": ["code", "message", "line", "column", "path"],
    "additionalProperties": False,
}
_STRUCTURED_VALIDATE_OUTPUT_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "valid": {"type": "boolean"},
        "format": {"type": "string", "enum": ["json", "yaml", "xml", "csv"]},
        "payload_hash": {"type": "string", "pattern": r"^[0-9a-f]{64}$"},
        "schema_hash": {"type": ["string", "null"], "pattern": r"^[0-9a-f]{64}$"},
        "errors": {"type": "array", "items": _VALIDATION_DIAGNOSTIC_SCHEMA, "maxItems": 50},
        "warnings": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "code": {"type": "string", "enum": ["csv_formula_prefix"]},
                    "count": {"type": "integer", "minimum": 1},
                },
                "required": ["code", "count"],
                "additionalProperties": False,
            },
            "maxItems": 50,
        },
        "details": {
            "type": "object",
            "properties": {
                "depth": {"type": "integer", "minimum": 0, "maximum": _STRUCTURED_MAX_DEPTH},
                "nodes": {"type": "integer", "minimum": 0, "maximum": _STRUCTURED_MAX_NODES},
                "rows": {"type": "integer", "minimum": 0, "maximum": _STRUCTURED_MAX_CSV_ROWS},
                "columns": {"type": "integer", "minimum": 0, "maximum": _STRUCTURED_MAX_CSV_COLUMNS},
            },
            "additionalProperties": False,
        },
    },
    "required": ["valid", "format", "payload_hash", "schema_hash", "errors", "warnings", "details"],
    "additionalProperties": False,
}


_RUST_COMPILE_CHECK_OUTPUT_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "compiles": {"type": ["boolean", "null"]},
        "diagnostics": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "level": {"type": "string", "enum": ["error", "warning"]},
                    "message": {"type": "string"},
                    "line": {"type": ["integer", "null"], "minimum": 1},
                    "column": {"type": ["integer", "null"], "minimum": 1},
                },
                "required": ["level", "message"],
                "additionalProperties": False,
            },
            "maxItems": 50,
        },
        "duration_ms": {"type": ["integer", "null"], "minimum": 0},
        "timed_out": {"type": "boolean"},
        "source_hash": {"type": "string", "pattern": r"^[0-9a-f]{64}$"},
        "sandbox_error": {"type": "string"},
    },
    "required": ["compiles", "diagnostics"],
    "additionalProperties": False,
}


_TOOL_CONTRACTS: Dict[str, Dict[str, Any]] = {
    "generate_image": {
        "contract_id": "moe.generative.generate_image",
        "contract_version": "1.0.0",
        "determinism": "generative_model",
        "source_policy": {"kind": "comfyui_api", "node": "N04-RGTX"},
        "inputSchema": {
            "type": "object",
            "properties": {
                "prompt": {"type": "string", "minLength": 1, "maxLength": 4096},
                "size": {"type": "string", "default": "1024x1024"},
                "model": {"type": "string", "default": "flux-schnell"},
            },
            "required": ["prompt"],
            "additionalProperties": False,
        },
        "outputSchema": {
            "type": "object",
            "properties": {
                "created": {"type": "integer"},
                "data": {"type": "array"},
                "error": {"type": "string"},
                "error_code": {"type": "string"},
            },
        },
    },
    "generate_speech": {
        "contract_id": "moe.generative.generate_speech",
        "contract_version": "1.0.0",
        "determinism": "generative_model",
        "source_policy": {"kind": "kokoro_tts_api", "node": "N04-RGTX"},
        "inputSchema": {
            "type": "object",
            "properties": {
                "text": {"type": "string", "minLength": 1, "maxLength": 10000},
                "voice": {"type": "string", "default": "af_heart"},
                "model": {"type": "string", "default": "kokoro"},
            },
            "required": ["text"],
            "additionalProperties": False,
        },
        "outputSchema": {
            "type": "object",
            "properties": {
                "ok": {"type": "boolean"},
                "voice": {"type": "string"},
                "model": {"type": "string"},
                "audio_bytes_length": {"type": "integer"},
                "error": {"type": "string"},
                "error_code": {"type": "string"},
            },
        },
    },
    "calendar_facts": {
        "contract_id": "moe.precision.calendar_facts",
        "contract_version": "1.0.0",
        "determinism": "input_only",
        "source_policy": {"kind": "python_stdlib", "name": "datetime-calendar"},
        "inputSchema": {
            "type": "object",
            "properties": {
                "date_str": {
                    "type": "string",
                    "pattern": r"^\d{4}-\d{2}-\d{2}$",
                },
                "locale": {"type": "string", "enum": ["de", "en"], "default": "de"},
            },
            "required": ["date_str"],
            "additionalProperties": False,
        },
        "outputSchema": _CALENDAR_FACTS_OUTPUT_SCHEMA,
    },
    "time_facts": {
        "contract_id": "moe.precision.time_facts",
        "contract_version": "1.0.0",
        "determinism": "source_versioned",
        "source_policy": {
            "kind": "python_tzdata",
            "name": "tzdata",
            "version": _PINNED_TZDATA_VERSION,
            "clock": "caller_supplied_instant",
        },
        "inputSchema": {
            "type": "object",
            "properties": {
                "instant": {
                    "type": "string",
                    "pattern": r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}(?:Z|[+-]\d{2}:\d{2})$",
                },
                "timezone_name": {
                    "type": "string",
                    "pattern": r"^[A-Za-z][A-Za-z0-9._+-]*(?:/[A-Za-z0-9._+-]+)*$",
                    "maxLength": 128,
                    "default": "UTC",
                },
                "locale": {"type": "string", "enum": ["de", "en"], "default": "de"},
            },
            "required": ["instant"],
            "additionalProperties": False,
        },
        "outputSchema": _TIME_FACTS_OUTPUT_SCHEMA,
        "limits": {"max_result_chars": 8192, "year_min": _TIME_YEAR_MIN, "year_max": _TIME_YEAR_MAX},
    },
    "timezone_convert": {
        "contract_id": "moe.precision.timezone_convert",
        "contract_version": "1.0.0",
        "determinism": "source_versioned",
        "source_policy": {
            "kind": "python_tzdata",
            "name": "tzdata",
            "version": _PINNED_TZDATA_VERSION,
            "clock": "explicit_local_datetime",
        },
        "inputSchema": {
            "type": "object",
            "properties": {
                "local_datetime": {
                    "type": "string",
                    "pattern": r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}$",
                },
                "from_timezone": {
                    "type": "string",
                    "pattern": r"^[A-Za-z][A-Za-z0-9._+-]*(?:/[A-Za-z0-9._+-]+)*$",
                    "maxLength": 128,
                },
                "to_timezone": {
                    "type": "string",
                    "pattern": r"^[A-Za-z][A-Za-z0-9._+-]*(?:/[A-Za-z0-9._+-]+)*$",
                    "maxLength": 128,
                },
                "fold": {"type": ["integer", "null"], "enum": [0, 1, None], "default": None},
                "locale": {"type": "string", "enum": ["de", "en"], "default": "de"},
            },
            "required": ["local_datetime", "from_timezone", "to_timezone"],
            "additionalProperties": False,
        },
        "outputSchema": _TIMEZONE_CONVERT_OUTPUT_SCHEMA,
        "limits": {"max_result_chars": 8192, "year_min": _TIME_YEAR_MIN, "year_max": _TIME_YEAR_MAX},
    },
    "decimal_finance": {
        "contract_id": "moe.precision.decimal_finance",
        "contract_version": "1.0.0",
        "determinism": "input_only",
        "source_policy": {"kind": "python_stdlib", "name": "decimal"},
        "inputSchema": {
            "type": "object",
            "properties": {
                "operation": {
                    "type": "string",
                    "enum": [
                        "add", "subtract", "multiply", "divide", "percentage",
                        "simple_interest", "compound_interest",
                    ],
                },
                "operands": {
                    "type": "array", "items": _DECIMAL_STRING_SCHEMA,
                    "minItems": 2, "maxItems": 4,
                },
                "currency": {"type": "string", "pattern": r"^[A-Z]{3}$"},
                "scale": {"type": "integer", "minimum": 0, "maximum": 12},
                "rounding": _ROUNDING_SCHEMA,
            },
            "required": ["operation", "operands", "currency", "scale", "rounding"],
            "additionalProperties": False,
            "allOf": [
                {
                    "if": {"properties": {"operation": {"enum": ["add", "subtract", "multiply", "divide", "percentage"]}}},
                    "then": {"properties": {"operands": {"minItems": 2, "maxItems": 2}}},
                },
                {
                    "if": {"properties": {"operation": {"const": "simple_interest"}}},
                    "then": {"properties": {"operands": {"minItems": 3, "maxItems": 3}}},
                },
                {
                    "if": {"properties": {"operation": {"const": "compound_interest"}}},
                    "then": {"properties": {"operands": {"minItems": 4, "maxItems": 4}}},
                },
            ],
        },
        "outputSchema": _DECIMAL_FINANCE_OUTPUT_SCHEMA,
        "limits": {"max_result_chars": 16384, "context_precision": _DECIMAL_CONTEXT_PRECISION, "max_scale": 12},
    },
    "exact_probability": {
        "contract_id": "moe.precision.exact_probability",
        "contract_version": "1.0.0",
        "determinism": "input_only",
        "source_policy": {"kind": "python_stdlib", "name": "fractions-math-decimal"},
        "inputSchema": {
            "type": "object",
            "properties": {
                "operation": {"type": "string", "enum": ["fraction", "combination", "permutation", "binomial_probability"]},
                "n": {"type": ["integer", "null"], "minimum": 0, "maximum": _PROBABILITY_MAX_N, "default": None},
                "k": {"type": ["integer", "null"], "minimum": 0, "maximum": _PROBABILITY_MAX_N, "default": None},
                "numerator": {"type": ["integer", "null"], "default": None},
                "denominator": {"type": ["integer", "null"], "default": None},
                "decimal_scale": {"type": ["integer", "null"], "minimum": 0, "maximum": 18, "default": None},
                "rounding": {"type": ["string", "null"], "enum": [*sorted(_DECIMAL_ROUNDING), None], "default": None},
            },
            "required": ["operation"],
            "additionalProperties": False,
        },
        "outputSchema": _EXACT_PROBABILITY_OUTPUT_SCHEMA,
        "limits": {"max_result_chars": 65536, "max_n": _PROBABILITY_MAX_N, "max_result_bits": _PROBABILITY_MAX_RESULT_BITS},
    },
    "structured_validate": {
        "contract_id": "moe.precision.structured_validate",
        "contract_version": "1.0.0",
        "determinism": "library_pinned",
        "source_policy": {"kind": "locked_parser_set", "name": "jsonschema-pyyaml-defusedxml-csv"},
        "evidence_policy": {
            "redact_input_fields": ["payload", "schema_json"],
            "replacement": "sha256_and_utf8_bytes",
        },
        "inputSchema": {
            "type": "object",
            "properties": {
                "format_name": {"type": "string", "enum": ["json", "yaml", "xml", "csv"]},
                "payload": {"type": "string", "maxLength": _STRUCTURED_MAX_PAYLOAD_BYTES},
                "schema_json": {"type": ["string", "null"], "maxLength": _STRUCTURED_MAX_SCHEMA_BYTES, "default": None},
                "csv_dialect": {"type": ["string", "null"], "enum": ["comma", "semicolon", "tab", "pipe", None], "default": None},
            },
            "required": ["format_name", "payload"],
            "additionalProperties": False,
        },
        "outputSchema": _STRUCTURED_VALIDATE_OUTPUT_SCHEMA,
        "limits": {
            "max_result_chars": 32768,
            "max_payload_bytes": _STRUCTURED_MAX_PAYLOAD_BYTES,
            "max_schema_bytes": _STRUCTURED_MAX_SCHEMA_BYTES,
            "max_depth": _STRUCTURED_MAX_DEPTH,
            "max_nodes": _STRUCTURED_MAX_NODES,
            "max_csv_rows": _STRUCTURED_MAX_CSV_ROWS,
            "max_csv_columns": _STRUCTURED_MAX_CSV_COLUMNS,
        },
    },
    "rust_compile_check": {
        "contract_id": "moe.precision.rust_compile_check",
        "contract_version": "1.0.0",
        "determinism": "library_pinned",
        "source_policy": {"kind": "pinned_toolchain", "name": "rustc 1.98 (rust:1-slim image digest)"},
        "evidence_policy": {
            "redact_input_fields": ["source"],
            "replacement": "sha256_and_utf8_bytes",
        },
        "inputSchema": {
            "type": "object",
            "properties": {
                "source": {"type": "string", "minLength": 1, "maxLength": _RUST_COMPILE_MAX_SOURCE_CHARS},
                "edition": {"type": "string", "enum": ["2015", "2018", "2021", "2024"], "default": "2021"},
            },
            "required": ["source"],
            "additionalProperties": False,
        },
        "outputSchema": _RUST_COMPILE_CHECK_OUTPUT_SCHEMA,
        "limits": {
            "max_result_chars": 32768,
            "max_source_chars": _RUST_COMPILE_MAX_SOURCE_CHARS,
            "compile_timeout_s": _RUST_COMPILE_HTTP_TIMEOUT_S,
        },
    },
    "gcd_lcm": {
        "contract_id": "moe.precision.gcd_lcm",
        "contract_version": "1.0.0",
        "determinism": "input_only",
        "source_policy": {"kind": "python_stdlib", "name": "math"},
        "inputSchema": {
            "type": "object",
            "properties": {
                "a": {"type": "integer", "minimum": -10**18, "maximum": 10**18},
                "b": {"type": "integer", "minimum": -10**18, "maximum": 10**18},
                "operation": {
                    "type": "string",
                    "enum": ["gcd", "lcm", "both"],
                    "default": "both",
                },
            },
            "required": ["a", "b"],
            "additionalProperties": False,
        },
        "outputSchema": {
            "type": "object",
            "properties": {
                "a": {"type": "integer"},
                "b": {"type": "integer"},
                "operation": {"type": "string", "enum": ["gcd", "lcm", "both"]},
                "gcd": {"type": "integer", "minimum": 0},
                "lcm": {"type": "integer", "minimum": 0},
            },
            "required": ["a", "b", "operation", "gcd", "lcm"],
            "additionalProperties": False,
        },
    },
    "unit_convert": {
        "contract_id": "moe.precision.unit_convert",
        "contract_version": "1.0.0",
        "determinism": "library_pinned",
        "source_policy": {"kind": "python_library", "name": "Pint"},
        "inputSchema": {
            "type": "object",
            "properties": {
                "value": {"type": "number"},
                "from_unit": {"type": "string", "minLength": 1, "maxLength": 64},
                "to_unit": {"type": "string", "minLength": 1, "maxLength": 64},
            },
            "required": ["value", "from_unit", "to_unit"],
            "additionalProperties": False,
        },
        "outputSchema": {
            "type": "object",
            "properties": {
                "value": {"type": "number"},
                "from_unit": {"type": "string"},
                "converted_value": {"type": "number"},
                "to_unit": {"type": "string"},
                "rendered": {"type": "string"},
            },
            "required": ["value", "from_unit", "converted_value", "to_unit", "rendered"],
            "additionalProperties": False,
        },
    },
}


_STRUCTURED_CONTRACT_POLICIES: Dict[str, Any] = {
    "normalization_policy": {
        "apply_schema_defaults": True,
        "reject_unknown_properties": True,
    },
    # A mandatory precision call is immutable after preflight. Transport
    # retries may repeat the exact same call at a higher layer, but neither
    # this service nor an LLM may repair or reinterpret its arguments.
    "retry_policy": {
        "max_attempts": 1,
        "argument_mutation": False,
    },
    # TASK-43 keeps answer caches bypassed. TASK-45 may introduce a typed,
    # post-quality cache whose key binds this contract hash and normalized
    # input; publishing that future policy here would be misleading.
    "cache_policy": {
        "mode": "bypass",
    },
}

for _contract in _TOOL_CONTRACTS.values():
    for _policy_name, _policy_value in _STRUCTURED_CONTRACT_POLICIES.items():
        _contract.setdefault(_policy_name, dict(_policy_value))


def _canonical_hash(value: Any) -> str:
    payload = json.dumps(
        value,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
        default=str,
    ).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def _contract_hash(name: str, contract: Dict[str, Any]) -> str:
    return _canonical_hash(
        {
            "tool": name,
            "contract_id": contract["contract_id"],
            "contract_version": contract["contract_version"],
            "determinism": contract["determinism"],
            "source_policy": contract["source_policy"],
            "inputSchema": contract["inputSchema"],
            "outputSchema": contract["outputSchema"],
            "normalization_policy": contract["normalization_policy"],
            "retry_policy": contract["retry_policy"],
            "cache_policy": contract["cache_policy"],
            "evidence_policy": contract.get("evidence_policy", {}),
            "limits": contract.get("limits", {"max_result_chars": 65536}),
        }
    )


def _evidence_input(args: Dict[str, Any], contract: Dict[str, Any]) -> Dict[str, Any]:
    """Return the contract-defined non-sensitive input projection."""
    projected = dict(args)
    policy = contract.get("evidence_policy") or {}
    for field in policy.get("redact_input_fields") or []:
        value = projected.get(field)
        if isinstance(value, str):
            encoded = value.encode("utf-8")
            projected[field] = {
                "sha256": hashlib.sha256(encoded).hexdigest(),
                "utf8_bytes": len(encoded),
            }
    return projected


def _runtime_source(
    name: str,
    contract: Dict[str, Any],
    facts: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    policy = contract["source_policy"]
    version = sys.version.split()[0]
    if policy.get("name") == "Pint":
        version = importlib.metadata.version("Pint")
    elif policy.get("name") == "tzdata":
        version = _tzdata_version()
    elif policy.get("name") == "jsonschema-pyyaml-defusedxml-csv":
        versions = []
        for package in ("jsonschema", "PyYAML", "defusedxml"):
            try:
                versions.append(f"{package}={importlib.metadata.version(package)}")
            except importlib.metadata.PackageNotFoundError:
                versions.append(f"{package}=unavailable-host-fallback")
        version = ";".join(versions) + f";python-csv={sys.version.split()[0]}"
    return {
        "kind": policy["kind"],
        "name": policy["name"],
        "version": version,
        "as_of": (facts or {}).get("as_of"),
    }


def _normalize_input(args: Dict[str, Any], schema: Dict[str, Any]) -> Dict[str, Any]:
    normalized = dict(args)
    for name, prop in schema.get("properties", {}).items():
        if name not in normalized and isinstance(prop, dict) and "default" in prop:
            normalized[name] = prop["default"]
    return normalized


def _schema_error(value: Any, schema: Dict[str, Any]) -> str:
    try:
        Draft202012Validator.check_schema(schema)
        errors = sorted(
            Draft202012Validator(schema).iter_errors(value),
            key=lambda item: (
                tuple(str(part) for part in item.absolute_path),
                item.message,
            ),
        )
    except SchemaError:
        return "schema_malformed"
    if not errors:
        return ""
    first = errors[0]
    path = ".".join(str(part) for part in first.absolute_path) or "value"
    return f"{path}:{first.validator or 'invalid'}"


def _structured_facts(name: str, args: Dict[str, Any], result: str) -> Dict[str, Any]:
    if name in {
        "calendar_facts", "time_facts", "timezone_convert",
        "decimal_finance", "exact_probability", "structured_validate",
        "rust_compile_check",
    }:
        facts = json.loads(result)
        if not isinstance(facts, dict):
            raise ValueError(f"{name} did not return an object")
        return facts
    if name == "gcd_lcm":
        a = int(args["a"])
        b = int(args["b"])
        gcd_value = math.gcd(abs(a), abs(b))
        return {
            "a": a,
            "b": b,
            "operation": str(args["operation"]),
            "gcd": gcd_value,
            "lcm": abs(a * b) // gcd_value if gcd_value else 0,
        }
    if name == "unit_convert":
        from pint import UnitRegistry
        value = args["value"]
        converted = (value * UnitRegistry()(args["from_unit"])).to(args["to_unit"])
        return {
            "value": value,
            "from_unit": args["from_unit"],
            "converted_value": float(converted.magnitude),
            "to_unit": args["to_unit"],
            "rendered": result,
        }
    raise ValueError(f"No structured result builder for {name}")

for _tn in _TOOL_DESCRIPTIONS:
    if _tn not in _TOOL_ACCESS_KIND:
        logger.warning(
            "mcp_server: tool '%s' has no _TOOL_ACCESS_KIND entry — "
            "defaulting to '%s' for telemetry classification", _tn, _DEFAULT_ACCESS_KIND,
        )

# ─── DISABLED TOOLS PERSISTENCE ───────────────────────────────────────────────
# Note: Disabling only applies to the REST /invoke path (which LangGraph
# uses). The MCP SSE endpoint /mcp (for Claude Desktop) uses FastMCP's
# internal registry and is not affected here.

_DISABLED_TOOLS_PATH = Path("/app/disabled_tools.json")
_disabled_tools_lock = threading.Lock()


def _load_disabled_tools() -> set:
    try:
        if _DISABLED_TOOLS_PATH.exists():
            return set(json.loads(_DISABLED_TOOLS_PATH.read_text(encoding="utf-8")).get("disabled", []))
    except Exception:
        pass
    return set()


def _save_disabled_tools(disabled: set) -> None:
    _DISABLED_TOOLS_PATH.write_text(
        json.dumps({"disabled": sorted(disabled)}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


_disabled_tools: set = _load_disabled_tools()

# ─── FASTAPI REST SHIM (used internally by LangGraph mcp_node) ──────────────

class InvokeRequest(BaseModel):
    tool: str
    args: Dict[str, Any] = {}


rest_app = FastAPI(title="Precision Tools REST Shim", version="1.0")


@rest_app.get("/health")
def health():
    return {"status": "ok", "tools": list(_TOOL_REGISTRY.keys())}


@rest_app.get("/downloads/{filename}")
def download_file(filename: str):
    """Serve generated files for download. Files auto-expire after 24h."""
    from fastapi.responses import FileResponse
    from fastapi import HTTPException as _HTTPException
    safe = re.sub(r'[^\w\-.]', '', filename)
    path = _GENERATED_DIR / safe
    if not path.exists() or not path.is_file():
        raise _HTTPException(status_code=404, detail="File not found or expired")
    media_types = {
        ".html": "text/html",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".md": "text/markdown",
        ".txt": "text/plain",
        ".pdf": "application/pdf",
    }
    mt = media_types.get(path.suffix, "application/octet-stream")
    return FileResponse(path, media_type=mt, filename=safe)


_PY_TO_JSON_TYPE: Dict[Any, str] = {
    str:   "string",
    int:   "integer",
    float: "number",
    bool:  "boolean",
    list:  "array",
    dict:  "object",
}


def _annotation_schema(annotation: Any, default: Any) -> Dict[str, Any]:
    """Translate the small set of runtime annotations used by MCP tools."""
    if annotation in (inspect.Parameter.empty, Any):
        schema: Dict[str, Any] = {}
    else:
        origin = get_origin(annotation)
        args = get_args(annotation)
        if origin in (list, List):
            schema = {"type": "array"}
        elif origin in (dict, Dict):
            schema = {"type": "object"}
        elif origin is Literal:
            values = list(args)
            json_types = sorted({_PY_TO_JSON_TYPE.get(type(value), "string") for value in values})
            schema = {
                "type": json_types[0] if len(json_types) == 1 else json_types,
                "enum": values,
            }
        elif origin is Union:
            json_types = []
            for member in args:
                if member is type(None):
                    json_types.append("null")
                else:
                    member_origin = get_origin(member)
                    if member_origin in (list, List):
                        json_types.append("array")
                    elif member_origin in (dict, Dict):
                        json_types.append("object")
                    else:
                        json_types.append(_PY_TO_JSON_TYPE.get(member, "string"))
            schema = {"type": list(dict.fromkeys(json_types))}
        else:
            schema = {"type": _PY_TO_JSON_TYPE.get(annotation, "string")}
    if default is None:
        raw_type = schema.get("type")
        if isinstance(raw_type, str):
            schema["type"] = [raw_type, "null"]
        elif isinstance(raw_type, list) and "null" not in raw_type:
            schema["type"] = [*raw_type, "null"]
    return schema


def _input_schema(func) -> Dict[str, Any]:
    """Derive a JSON-Schema-compatible inputSchema from a Python function signature.

    Used to make each tool's parameter contract visible to Open WebUI Tool Servers
    and to populate MCP_TOOL_SCHEMAS for pre-call argument validation in the pipeline.
    Ignores *args / **kwargs; defaults become non-required properties.
    """
    props: Dict[str, Any] = {}
    required: List[str] = []
    try:
        sig = inspect.signature(func)
        for pname, param in sig.parameters.items():
            if param.kind in (param.VAR_POSITIONAL, param.VAR_KEYWORD):
                continue
            prop = _annotation_schema(
                param.annotation,
                None if param.default is None else inspect.Parameter.empty,
            )
            if param.default is not inspect.Parameter.empty:
                prop["default"] = param.default
            else:
                required.append(pname)
            props[pname] = prop
    except (ValueError, TypeError):
        pass
    return {
        "type": "object",
        "properties": props,
        "required": required,
        "additionalProperties": False,
    }


def build_tools_catalog() -> Dict[str, Any]:
    """Build the REST discovery document independently of FastAPI wiring."""
    with _disabled_tools_lock:
        disabled = set(_disabled_tools)
    tools = []
    for name, desc in _TOOL_DESCRIPTIONS.items():
        func = _TOOL_REGISTRY.get(name)
        contract = _TOOL_CONTRACTS.get(name)
        schema = (
            contract["inputSchema"]
            if contract
            else _input_schema(func) if func else {
                "type": "object",
                "properties": {},
                "required": [],
                "additionalProperties": False,
            }
        )
        output_schema = (
            contract["outputSchema"] if contract else {"type": "string"}
        )
        contract_hash = _contract_hash(name, contract) if contract else ""
        tools.append({
            "name":        name,
            "description": desc,
            "enabled":     name not in disabled,
            "inputSchema": schema,                    # Open WebUI Tool Server format
            "args":        schema["properties"],      # pipeline MCP_TOOL_SCHEMAS compat
            "required_args": schema["required"],      # pipeline pre-call validation
            "access_kind": _TOOL_ACCESS_KIND.get(name, _DEFAULT_ACCESS_KIND),  # telemetry only
            "outputSchema": output_schema,
            "contract_id": contract["contract_id"] if contract else f"legacy.{name}",
            "contract_version": contract["contract_version"] if contract else "0.0.0",
            "contract_hash": contract_hash,
            "determinism": contract["determinism"] if contract else "unclassified",
            "source_policy": contract["source_policy"] if contract else {"kind": "unclassified"},
            "structured_result": bool(contract),
            "normalization_policy": contract.get("normalization_policy", {}) if contract else {},
            "retry_policy": contract.get("retry_policy", {}) if contract else {},
            "cache_policy": contract.get("cache_policy", {"mode": "legacy"}) if contract else {"mode": "legacy"},
            "limits": contract.get("limits", {"max_result_chars": 65536}) if contract else {"max_result_chars": 65536},
            "evidence_policy": contract.get("evidence_policy", {}) if contract else {},
        })
    return {"tools": tools}


@rest_app.get("/tools")
def list_tools():
    return build_tools_catalog()


@rest_app.post("/tools/{name}/toggle")
def toggle_tool(name: str):
    from fastapi import HTTPException as _HTTPException
    if name not in _TOOL_REGISTRY:
        raise _HTTPException(status_code=404, detail=f"Unknown tool: {name}")
    with _disabled_tools_lock:
        if name in _disabled_tools:
            _disabled_tools.discard(name)
            now_enabled = True
        else:
            _disabled_tools.add(name)
            now_enabled = False
        _save_disabled_tools(_disabled_tools)
    return {"ok": True, "name": name, "enabled": now_enabled}


async def execute_tool(req: InvokeRequest) -> Dict[str, Any]:
    """Execute one REST tool contract independently of FastAPI wiring."""
    if req.tool not in _TOOL_REGISTRY:
        return {
            "error": f"Unknown tool: '{req.tool}'. Available: {list(_TOOL_REGISTRY.keys())}",
            "error_code": "unknown_tool",
            "tool": req.tool,
        }
    with _disabled_tools_lock:
        if req.tool in _disabled_tools:
            return {
                "error": f"Tool '{req.tool}' is disabled.",
                "error_code": "tool_disabled",
                "reason": "disabled",
                "tool": req.tool,
            }
    try:
        import inspect
        func = _TOOL_REGISTRY[req.tool]
        contract = _TOOL_CONTRACTS.get(req.tool)
        input_schema = contract["inputSchema"] if contract else _input_schema(func)
        normalized_args = _normalize_input(req.args, input_schema)
        input_error = _schema_error(normalized_args, input_schema)
        if input_error:
            return {
                "error": f"Input schema validation failed for '{req.tool}': {input_error}",
                "error_code": "input_schema_invalid",
                "tool": req.tool,
            }
        if inspect.iscoroutinefunction(func):
            result = await func(**normalized_args)
        else:
            result = func(**normalized_args)
        max_result_chars = int(
            (contract or {}).get("limits", {}).get("max_result_chars", 65536)
        )
        if len(str(result)) > max_result_chars:
            return {
                "error": f"Result for '{req.tool}' exceeds the contract size limit",
                "error_code": "tool_result_too_large",
                "tool": req.tool,
            }
        response: Dict[str, Any] = {"result": result, "tool": req.tool}
        if contract:
            if not isinstance(result, str) or result.casefold().startswith(("error:", "fehler:")):
                return {
                    "error": f"Structured tool '{req.tool}' returned an error result",
                    "error_code": "tool_result_error",
                    "tool": req.tool,
                }
            facts = _structured_facts(req.tool, normalized_args, result)
            output_error = _schema_error(facts, contract["outputSchema"])
            if output_error:
                return {
                    "error": f"Output schema validation failed for '{req.tool}': {output_error}",
                    "error_code": "output_schema_invalid",
                    "tool": req.tool,
                }
            contract_hash = _contract_hash(req.tool, contract)
            evidence_input = _evidence_input(normalized_args, contract)
            result_hash = _canonical_hash(
                {
                    "contract_hash": contract_hash,
                    "input_normalized": evidence_input,
                    "facts": facts,
                }
            )
            response["structured_result"] = {
                "status": "completed",
                "tool": req.tool,
                "contract_id": contract["contract_id"],
                "contract_version": contract["contract_version"],
                "contract_hash": contract_hash,
                "input_normalized": evidence_input,
                "facts": facts,
                "determinism": contract["determinism"],
                "source": _runtime_source(req.tool, contract, facts),
                "warnings": [],
                "result_hash": result_hash,
            }
        return response
    except TypeError as e:
        return {
            "error": f"Wrong arguments for '{req.tool}': {e}",
            "error_code": "wrong_arguments",
            "tool": req.tool,
        }
    except Exception as e:
        return {
            "error": f"Error in '{req.tool}': {e}",
            "error_code": "tool_execution_error",
            "tool": req.tool,
        }


@rest_app.post("/invoke")
async def invoke_tool(req: InvokeRequest):
    return await execute_tool(req)


# Mount MCP SSE app at /mcp (for Claude Desktop, MCP clients, etc.)
rest_app.mount("/mcp", mcp.sse_app())


if __name__ == "__main__":
    uvicorn.run(rest_app, host="0.0.0.0", port=8003, log_level="info")
